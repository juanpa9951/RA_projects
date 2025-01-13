#####  THIS FOR CALIBRATION METHOD 2
def master_pattern_cross(image_size,Left_centering,Top_centering,Axis_Limit,background_color,layer_color,close_image,Axis_Limit_Y,x_base,y_base,x_step,y_step):

    import matplotlib.pyplot as plt
    from pptx import Presentation
    from pptx.util import Inches
    import numpy as np


    # POWERPOINT SLIDE
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank slide layout

    ############# BEGGIN THE PLOTTING

    Axis_Limit_x = Axis_Limit
    Axis_Limit_y = Axis_Limit_Y
    # x_base=500    #500
    # y_base=500    #100

    for k in range(1,y_step):   ## divisiones eje Y     25
            for i in range(1,x_step):  ## divisiones eje X     32

                        x_lenght=x_base*i                  ####modo x
                        y_lenght=y_base*k

                        x_layer=[x_lenght]        ####modo x
                        y_layer=[y_lenght]        ####modo x

                        Layer_Name_plot='X='+str(x_lenght)+" Y="+str(y_lenght)

                        # Plotting the polygon
                        plt.figure(figsize=(image_size, image_size))
                        plt.plot(x_layer, y_layer,color='blue', linestyle='-', linewidth=0.01, marker='x',markersize=4) # 'bo-' means blue circles connected by lines   o-0.3   x-4
                        title='pattern'
                        plt.title(title)


                        plt.xlim(0, Axis_Limit_x)  # Set x-axis limit
                        plt.ylim(0, Axis_Limit_y)  # Set y-axis limit
                        plt.style.use('dark_background')
                        ax = plt.gca()
                        ax.set_facecolor(background_color)
                        ax.spines['top'].set_color(layer_color)
                        ax.spines['bottom'].set_color(layer_color)
                        ax.spines['right'].set_color(layer_color)
                        ax.spines['left'].set_color(layer_color)
                        ax.tick_params(axis='x', colors=layer_color)
                        ax.tick_params(axis='y', colors=layer_color)
                        ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                        plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
                        plt.savefig('Layers.png')
                        # add slide
                        slide = prs.slides.add_slide(slide_layout)
                        img_path = 'Layers.png'
                        left = Inches(Left_centering)  # Adjust position as needed in icnhes
                        top = Inches(Top_centering)  # Adjust position as needed in inches
                        slide.shapes.add_picture(img_path, left, top)
                        if close_image==1:
                         plt.close('all')

                        # prs.save('Layers.pptx')

    prs.save('Master_Pattern_cross.pptx')
def master_pattern_dot(image_size,Left_centering,Top_centering,Axis_Limit,background_color,layer_color,close_image,Axis_Limit_Y,x_base,y_base,x_step,y_step):

    import matplotlib.pyplot as plt
    from pptx import Presentation
    from pptx.util import Inches
    import numpy as np


    # POWERPOINT SLIDE
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank slide layout

    ############# BEGGIN THE PLOTTING

    Axis_Limit_x = Axis_Limit
    Axis_Limit_y = Axis_Limit_Y
    # x_base=500    #500
    # y_base=500    #100

    for k in range(1,y_step):   ## divisiones eje Y     25
            for i in range(1,x_step):  ## divisiones eje X     32
                        x_lenght=x_base*i                  ####modo x
                        y_lenght=y_base*k                  ####modo x

                        x_layer=[x_lenght]               ####modo x
                        y_layer=[y_lenght]        ####modo x


                        Layer_Name_plot='X='+str(x_lenght)+" Y="+str(y_lenght)

                        # Plotting the polygon
                        plt.figure(figsize=(image_size, image_size))
                        plt.plot(x_layer, y_layer,color='blue', linestyle='-', linewidth=0.01, marker='o',markersize=0.3) # 'bo-' means blue circles connected by lines   o-0.3   x-4
                        title='pattern'
                        plt.title(title)


                        plt.xlim(0, Axis_Limit_x)  # Set x-axis limit
                        plt.ylim(0, Axis_Limit_y)  # Set y-axis limit
                        plt.style.use('dark_background')
                        ax = plt.gca()
                        ax.set_facecolor(background_color)
                        ax.spines['top'].set_color(layer_color)
                        ax.spines['bottom'].set_color(layer_color)
                        ax.spines['right'].set_color(layer_color)
                        ax.spines['left'].set_color(layer_color)
                        ax.tick_params(axis='x', colors=layer_color)
                        ax.tick_params(axis='y', colors=layer_color)
                        ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                        plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
                        plt.savefig('Layers.png')
                        # add slide
                        slide = prs.slides.add_slide(slide_layout)
                        img_path = 'Layers.png'
                        left = Inches(Left_centering)  # Adjust position as needed in icnhes
                        top = Inches(Top_centering)  # Adjust position as needed in inches
                        slide.shapes.add_picture(img_path, left, top)
                        if close_image==1:
                         plt.close('all')

                        # prs.save('Layers.pptx')

    prs.save('Master_Pattern_dot.pptx')
def master_pattern_square(image_size,Left_centering,Top_centering,Axis_Limit,background_color,layer_color,close_image,Axis_Limit_Y,x_base,y_base,x_step,y_step):

    import matplotlib.pyplot as plt
    from pptx import Presentation
    from pptx.util import Inches
    import numpy as np


    # POWERPOINT SLIDE
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank slide layout

    ############# BEGGIN THE PLOTTING

    Axis_Limit_x = Axis_Limit
    Axis_Limit_y = Axis_Limit_Y
    # x_base=400    #500
    # y_base=300    #100

    for k in range(1,y_step):   ## divisiones eje Y     25
            for i in range(1,x_step):  ## divisiones eje X     32
                        x_lenght=x_base*i                  ####modo x
                        y_lenght=y_base*k                  ####modo x

                        x_layer=[0,x_lenght,x_lenght]               ####modo x
                        y_layer=[y_lenght,y_lenght,0]        ####modo x


                        Layer_Name_plot='X='+str(x_lenght)+" Y="+str(y_lenght)

                        # Plotting the polygon
                        plt.figure(figsize=(image_size, image_size))
                        plt.plot(x_layer, y_layer,color='blue', linestyle='-', linewidth=0.3, marker='x',markersize=0.2) # 'bo-' means blue circles connected by lines
                        title='pattern'
                        plt.title(title)


                        plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
                        plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
                        plt.style.use('dark_background')
                        ax = plt.gca()
                        ax.set_facecolor(background_color)
                        ax.spines['top'].set_color(layer_color)
                        ax.spines['bottom'].set_color(layer_color)
                        ax.spines['right'].set_color(layer_color)
                        ax.spines['left'].set_color(layer_color)
                        ax.tick_params(axis='x', colors=layer_color)
                        ax.tick_params(axis='y', colors=layer_color)
                        ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                        plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
                        plt.savefig('Layers.png')
                        # add slide
                        slide = prs.slides.add_slide(slide_layout)
                        img_path = 'Layers.png'
                        left = Inches(Left_centering)  # Adjust position as needed in icnhes
                        top = Inches(Top_centering)  # Adjust position as needed in inches
                        slide.shapes.add_picture(img_path, left, top)
                        if close_image==1:
                         plt.close('all')

                        # prs.save('Layers.pptx')

    prs.save('Master_Pattern_square.pptx')
def master_pattern_crossV1(image_size,Left_centering,Top_centering,Axis_Limit,background_color,layer_color,close_image,Axis_Limit_Y,x_base,y_base,x_step,y_step,x_init,x_final,y_init):

    import matplotlib.pyplot as plt
    from pptx import Presentation
    from pptx.util import Inches
    import numpy as np


    # POWERPOINT SLIDE
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank slide layout

    ############# BEGGIN THE PLOTTING

    Axis_Limit_x = Axis_Limit
    Axis_Limit_y = Axis_Limit_Y
    # x_base=500    #500
    # y_base=500    #100

    for k in range(0,y_step):   ## divisiones eje Y     25
            for i in range(0,x_step):  ## divisiones eje X     32
                        if i==0:
                            x_lenght=x_init                  ####modo x
                            y_lenght=y_base*k                  ####modo x
                        elif i==x_step-1:
                            x_lenght=x_final                  ####modo x
                            y_lenght=y_base*k
                        else:
                            x_lenght=x_base*i                  ####modo x
                            y_lenght=y_base*k

                        if k==0:
                            y_lenght = y_init

                        x_layer=[x_lenght]        ####modo x
                        y_layer=[y_lenght]        ####modo x


                        Layer_Name_plot='X='+str(x_lenght)+" Y="+str(y_lenght)

                        # Plotting the polygon
                        plt.figure(figsize=(image_size, image_size))
                        plt.plot(x_layer, y_layer,color='blue', linestyle='-', linewidth=0.01, marker='x',markersize=4) # 'bo-' means blue circles connected by lines   o-0.3   x-4
                        title='pattern'
                        plt.title(title)


                        plt.xlim(0, Axis_Limit_x)  # Set x-axis limit
                        plt.ylim(0, Axis_Limit_y)  # Set y-axis limit
                        plt.style.use('dark_background')
                        ax = plt.gca()
                        ax.set_facecolor(background_color)
                        ax.spines['top'].set_color(layer_color)
                        ax.spines['bottom'].set_color(layer_color)
                        ax.spines['right'].set_color(layer_color)
                        ax.spines['left'].set_color(layer_color)
                        ax.tick_params(axis='x', colors=layer_color)
                        ax.tick_params(axis='y', colors=layer_color)
                        ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                        plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
                        plt.savefig('Layers.png')
                        # add slide
                        slide = prs.slides.add_slide(slide_layout)
                        img_path = 'Layers.png'
                        left = Inches(Left_centering)  # Adjust position as needed in icnhes
                        top = Inches(Top_centering)  # Adjust position as needed in inches
                        slide.shapes.add_picture(img_path, left, top)
                        if close_image==1:
                         plt.close('all')

                        # prs.save('Layers.pptx')

    prs.save('Master_Pattern_cross.pptx')
def master_pattern_dotV1(image_size,Left_centering,Top_centering,Axis_Limit,background_color,layer_color,close_image,Axis_Limit_Y,x_base,y_base,x_step,y_step,x_init,x_final,y_init):

    import matplotlib.pyplot as plt
    from pptx import Presentation
    from pptx.util import Inches
    import numpy as np


    # POWERPOINT SLIDE
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank slide layout

    ############# BEGGIN THE PLOTTING

    Axis_Limit_x = Axis_Limit
    Axis_Limit_y = Axis_Limit_Y
    # x_base=500    #500
    # y_base=500    #100

    for k in range(0,y_step):   ## divisiones eje Y     25
            for i in range(0,x_step):  ## divisiones eje X     32
                        if i==0:
                            x_lenght=x_init                  ####modo x
                            y_lenght=y_base*k                  ####modo x
                        elif i==x_step-1:
                            x_lenght=x_final                  ####modo x
                            y_lenght=y_base*k
                        else:
                            x_lenght=x_base*i                  ####modo x
                            y_lenght=y_base*k

                        if k==0:
                            y_lenght = y_init

                        x_layer=[x_lenght]               ####modo x
                        y_layer=[y_lenght]        ####modo x


                        Layer_Name_plot='X='+str(x_lenght)+" Y="+str(y_lenght)

                        # Plotting the polygon
                        plt.figure(figsize=(image_size, image_size))
                        plt.plot(x_layer, y_layer,color='blue', linestyle='-', linewidth=0.01, marker='o',markersize=0.3) # 'bo-' means blue circles connected by lines   o-0.3   x-4
                        title='pattern'
                        plt.title(title)


                        plt.xlim(0, Axis_Limit_x)  # Set x-axis limit
                        plt.ylim(0, Axis_Limit_y)  # Set y-axis limit
                        plt.style.use('dark_background')
                        ax = plt.gca()
                        ax.set_facecolor(background_color)
                        ax.spines['top'].set_color(layer_color)
                        ax.spines['bottom'].set_color(layer_color)
                        ax.spines['right'].set_color(layer_color)
                        ax.spines['left'].set_color(layer_color)
                        ax.tick_params(axis='x', colors=layer_color)
                        ax.tick_params(axis='y', colors=layer_color)
                        ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                        plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
                        plt.savefig('Layers.png')
                        # add slide
                        slide = prs.slides.add_slide(slide_layout)
                        img_path = 'Layers.png'
                        left = Inches(Left_centering)  # Adjust position as needed in icnhes
                        top = Inches(Top_centering)  # Adjust position as needed in inches
                        slide.shapes.add_picture(img_path, left, top)
                        if close_image==1:
                         plt.close('all')

                        # prs.save('Layers.pptx')

    prs.save('Master_Pattern_dot.pptx')
def master_pattern_squareV1(image_size,Left_centering,Top_centering,Axis_Limit,background_color,layer_color,close_image,Axis_Limit_Y,x_base,y_base,x_step,y_step,x_init,x_final,y_init):

    import matplotlib.pyplot as plt
    from pptx import Presentation
    from pptx.util import Inches
    import numpy as np


    # POWERPOINT SLIDE
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank slide layout

    ############# BEGGIN THE PLOTTING

    Axis_Limit_x = Axis_Limit
    Axis_Limit_y = Axis_Limit_Y
    # x_base=400    #500
    # y_base=300    #100

    for k in range(0,y_step):   ## divisiones eje Y     25
            for i in range(0,x_step):  ## divisiones eje X     32

                        if i==0:
                            x_lenght=x_init                  ####modo x
                            y_lenght=y_base*k                  ####modo x
                        elif i==x_step-1:
                            x_lenght=x_final                  ####modo x
                            y_lenght=y_base*k
                        else:
                            x_lenght=x_base*i                  ####modo x
                            y_lenght=y_base*k

                        if k==0:
                            y_lenght = y_init

                        x_layer=[0,x_lenght,x_lenght]               ####modo x
                        y_layer=[y_lenght,y_lenght,0]        ####modo x


                        Layer_Name_plot='X='+str(x_lenght)+" Y="+str(y_lenght)

                        # Plotting the polygon
                        plt.figure(figsize=(image_size, image_size))
                        plt.plot(x_layer, y_layer,color='blue', linestyle='-', linewidth=0.3, marker='x',markersize=0.2) # 'bo-' means blue circles connected by lines
                        title='pattern'
                        plt.title(title)


                        plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
                        plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
                        plt.style.use('dark_background')
                        ax = plt.gca()
                        ax.set_facecolor(background_color)
                        ax.spines['top'].set_color(layer_color)
                        ax.spines['bottom'].set_color(layer_color)
                        ax.spines['right'].set_color(layer_color)
                        ax.spines['left'].set_color(layer_color)
                        ax.tick_params(axis='x', colors=layer_color)
                        ax.tick_params(axis='y', colors=layer_color)
                        ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                        plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
                        plt.savefig('Layers.png')
                        # add slide
                        slide = prs.slides.add_slide(slide_layout)
                        img_path = 'Layers.png'
                        left = Inches(Left_centering)  # Adjust position as needed in icnhes
                        top = Inches(Top_centering)  # Adjust position as needed in inches
                        slide.shapes.add_picture(img_path, left, top)
                        if close_image==1:
                         plt.close('all')

                        # prs.save('Layers.pptx')

    prs.save('Master_Pattern_square.pptx')


image_size = 17  # in inches    12.9
Left_centering = -3.7  # in inches   -1.60
Top_centering = -8.85  # in inches   -4.0
Axis_Limit  = 9800  # in MM    # 10000
Axis_Limit_Y=9800  # ONLY FOR V19,  LAST OFFICIAL 10000

background_color = 'white'
layer_color = 'blue'
close_image = 1  # 1- close all images, 2- open all images

x_base=1000
y_base=500
x_step=10
y_step=10
mode=1    #  0- normal     1- limit values
x_init=10
x_final=9790
y_init=10

####  CHOOSE EITHER X OR Y
if mode==0:
    master_pattern_cross(image_size, Left_centering, Top_centering, Axis_Limit, background_color, layer_color,close_image, Axis_Limit_Y, x_base, y_base, x_step, y_step)
    master_pattern_dot(image_size, Left_centering, Top_centering, Axis_Limit, background_color, layer_color,close_image, Axis_Limit_Y, x_base, y_base, x_step, y_step)
    master_pattern_square(image_size, Left_centering, Top_centering, Axis_Limit, background_color, layer_color,close_image, Axis_Limit_Y, x_base, y_base, x_step, y_step)
    print("ready files")
else:
    master_pattern_crossV1(image_size,Left_centering,Top_centering,Axis_Limit,background_color,layer_color,close_image,Axis_Limit_Y,x_base,y_base,x_step,y_step,x_init,x_final,y_init)
    master_pattern_dotV1(image_size,Left_centering,Top_centering,Axis_Limit,background_color,layer_color,close_image,Axis_Limit_Y,x_base,y_base,x_step,y_step,x_init,x_final,y_init)
    master_pattern_squareV1(image_size,Left_centering,Top_centering,Axis_Limit,background_color,layer_color,close_image,Axis_Limit_Y,x_base,y_base,x_step,y_step,x_init,x_final,y_init)
    print("ready files")



