import ezdxf
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches


##............READING THE DXF
doc = ezdxf.readfile("V236_LW_LRO081.dxf")

# POWERPOINT SLIDE
prs = Presentation()
slide_layout = prs.slide_layouts[6]  # Blank slide layout
Axis_Limit=3200  # in dxf unit  3200
Left_centering=1.2   #in inches

# Extract entities (splines)
msp = doc.modelspace()
splines = msp.query('SPLINE')
i=0
x=[]
y=[]
# Iterate over splines
for spline in splines:
    # Check if the spline is a BSpline
    if spline.dxftype() == 'SPLINE':  # here we look for entity type SPLINE
        # Get control points of the spline
        control_points = spline._control_points

        # Print control points
        for point in control_points:
            i=i+1
            print("Vertice:", point,'vertice number', i)
            x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
            y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline

# OPTIONAL TRANSFORMATION FOR CENTERING
x=[i - min(x)*0.9 for i in x]
y=[i - min(y)*0.9 for i in y]
Axis_Limit=max(max(x),max(y))+100
#..............................................

x_init=x[0]
y_init=y[0]
z=0
layer_qty=0
check=0
for i in range(0,len(x)):    # loop through all the vertices of the splines found
    #print(i)
    if x[i]==x_init and y[i]==y_init:   # layer detection condition is when there is a closed loop, initial and final points are equal
        check=check+1 # the initial point of the loop
        if check==2: # the final point of the loop
            x_layer=x[z:i+ 1]
            y_layer = y[z:i + 1]
            z=i+1
            layer_qty=layer_qty+1
            check=0
            if i<(len(x)-1):
             x_init = x[i+1]
             y_init = y[i+1]
            print('Layer ', layer_qty)
            print('X_coordenates ', x_layer)
            print('y_coordenates ', y_layer)

            # Plotting the polygon
            plt.figure(figsize=(7, 7))
            plt.plot(x_layer, y_layer, color='red', linestyle='-', linewidth=1, marker='o', markersize=3) # 'bo-' means blue circles connected by lines
            plt.fill(x_layer, y_layer, alpha=1)  # Fill the polygon
            title='Layer '+str(layer_qty)
            plt.title(title)
            plt.grid(True, linestyle='--', linewidth=0.5)  # Set linewidth to 1.5
            plt.xlim(0, Axis_Limit)  # Set x-axis limit from 0 to 4
            plt.ylim(0, Axis_Limit)  # Set y-axis limit from 0 to 4
            plt.savefig('Layers.png')
            # add slide
            slide = prs.slides.add_slide(slide_layout)
            img_path = 'Layers.png'
            left = Inches(Left_centering)  # Adjust position as needed in icnhes
            top = Inches(0)  # Adjust position as needed in inches
            slide.shapes.add_picture(img_path, left, top)
# Save the PowerPoint presentation
prs.save('Layers.pptx')


print('total layers = ',layer_qty)

















