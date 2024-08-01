def read_and_plot_layerv13(file_path,name,image_size,Left_centering,Top_centering,Axis_Limit,scale_mode,Reduce_factor,background_color,layer_color,close_image,pattern_mode,raspberry,Destination_path):
    ### VERSION USING LAYER NAMES, ITS IMPORTANT TO HAVE LAYER COLORS ASSIGNED IN THE DXF
    import ezdxf
    import pandas as pd
    import matplotlib.pyplot as plt
    from pptx import Presentation
    from pptx.util import Inches


    ##............READING THE DXF
    doc = ezdxf.readfile(file_path)

    # POWERPOINT SLIDE
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank slide layout

    # Extract entities (splines)
    msp = doc.modelspace()
    # for entity in msp:
    #     # Get the type of the entity
    #     print(entity.dxftype())
    splines = msp.query('SPLINE')
    polylines= msp.query('POLYLINE')
    lw_polylines= msp.query('LWPOLYLINE')
    i=0
    x=[]
    y=[]
    x_tag=[]
    y_tag=[]
    x_datum=[]
    y_datum=[]
    ####.................... Iterate over splines/polylines/lw-polyinies
    for spline in splines:
            # Check if the spline is a BSpline
            if spline.dxftype() == 'SPLINE':  # here we look for entity type SPLINE
                if spline.dxf.color==1 or spline.dxf.color==3:   # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline._control_points

                    # Print control points
                    for point in control_points:
                        i=i+1
                        #print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline


    for spline in polylines:
        # Check if the spline is a BSpline
            if spline.dxftype() == 'POLYLINE':  # here we look for entity type SPLINE
                if spline.dxf.color==1 or spline.dxf.color==3:   # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline.points()
                    # Print control points
                    for point in control_points:
                        i=i+1
                        #print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline

    qty_spline=0
    Layer_Names=[]
    Layer_points=[]
    for spline in lw_polylines:
        # Check if the spline is a BSpline
            if spline.dxftype() == 'LWPOLYLINE':  # here we look for entity type SPLINE
                 #print(spline.dxf.color)
                 layer_name = spline.dxf.layer  # Get the layer name
                 layer = doc.layers.get(layer_name)  # Get the layer object
                 color = layer.dxf.color  # Get the color from the layer
                 #print (color)
                 ######## GET THE GREEN A RED LAYERS
                 if color==1 or color==3:    # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline.get_points('xy')

                    # Print control points
                    for point in control_points:
                        i=i+1
                        #print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline
                        Layer_points.append(layer_name)
                    #####................................................
                    if qty_spline == 0:
                        Layer_Names.append(layer_name)
                    elif Layer_Names[-1] != layer_name:
                        Layer_Names.append(layer_name)
                    qty_spline=qty_spline+1

                 ################### GET THE TAGGING LINES
                 if color == 7 or color == 255:
                     control_points = spline.get_points('xy')
                     for point in control_points:
                            i=i+1
                            #print("Vertice:", point,'vertice number', i)
                            x_tag.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                            y_tag.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline

                 ################### GET THE DATUM LINES
                 if color == 141:
                     control_points = spline.get_points('xy')
                     for point in control_points:
                            i=i+1
                            #print("Vertice:", point,'vertice number', i)
                            x_datum.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                            y_datum.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline




            #####...................1st scale reduction  OPTIONAL...........................
    # Reduce_factor=1
    x = [i / Reduce_factor for i in x]
    y = [i / Reduce_factor for i in y]

    excel_table_calib = r'C:\Users\Juan Pablo Lopez\PycharmProjects\ProjectJP\Table_Calib.xlsx'
    # .....LOAD THE TABLE CALIBRATION DATA.......................................................................................................
    if raspberry==1:
        X_table = pd.read_excel(excel_table_calib, sheet_name='X_axis_RP', header=0)
        Y_table = pd.read_excel(excel_table_calib, sheet_name='Y_axis_RP', header=0)
    else:
        X_table = pd.read_excel(excel_table_calib, sheet_name='X_axis', header=0)
        Y_table = pd.read_excel(excel_table_calib, sheet_name='Y_axis', header=0)

    if scale_mode==1:
        for i in range(0,len(x)):
            target_value = x[i]
            df=X_table
            # Find indices of bracketing points
            idx = df['Real_Measure'].searchsorted(target_value)

            # Check if target is within data range
            if idx == 0 or idx == len(df):
                print("Target value outside data range for interpolation X ",target_value)
            else:
                # Extract bracketed values
                a_prev = df.loc[idx - 1, 'Real_Measure']
                a_next = df.loc[idx, 'Real_Measure']
                b_prev = df.loc[idx - 1, 'Autocad_Measure']
                b_next = df.loc[idx, 'Autocad_Measure']

                # Perform linear interpolation
                interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                # print(f"Interpolated value for {target_value}: {interpolated_value}")
                x[i]=interpolated_value

        for i in range(0,len(y)):
            target_value = y[i]
            df=Y_table
            # Find indices of bracketing points
            idx = df['Real_Measure'].searchsorted(target_value)

            # Check if target is within data range
            if idx == 0 or idx == len(df):
                print("Target value outside data range for interpolation Y ",target_value)
            else:
                # Extract bracketed values
                a_prev = df.loc[idx - 1, 'Real_Measure']
                a_next = df.loc[idx, 'Real_Measure']
                b_prev = df.loc[idx - 1, 'Autocad_Measure']
                b_next = df.loc[idx, 'Autocad_Measure']

                # Perform linear interpolation
                interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                # print(f"Interpolated value for {target_value}: {interpolated_value}")
                y[i]=interpolated_value

        if len(x_tag)>0:
            for i in range(0,len(x_tag)):
                target_value = x_tag[i]
                df=X_table
                # Find indices of bracketing points
                idx = df['Real_Measure'].searchsorted(target_value)

                # Check if target is within data range
                if idx == 0 or idx == len(df):
                    print("Target value outside data range for interpolation X_Tag ",target_value)
                else:
                    # Extract bracketed values
                    a_prev = df.loc[idx - 1, 'Real_Measure']
                    a_next = df.loc[idx, 'Real_Measure']
                    b_prev = df.loc[idx - 1, 'Autocad_Measure']
                    b_next = df.loc[idx, 'Autocad_Measure']

                    # Perform linear interpolation
                    interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                    # print(f"Interpolated value for {target_value}: {interpolated_value}")
                    x_tag[i]=interpolated_value

            for i in range(0,len(y_tag)):
                target_value = y_tag[i]
                df=Y_table
                # Find indices of bracketing points
                idx = df['Real_Measure'].searchsorted(target_value)

                # Check if target is within data range
                if idx == 0 or idx == len(df):
                    print("Target value outside data range for interpolation Y_tag ",target_value)
                else:
                    # Extract bracketed values
                    a_prev = df.loc[idx - 1, 'Real_Measure']
                    a_next = df.loc[idx, 'Real_Measure']
                    b_prev = df.loc[idx - 1, 'Autocad_Measure']
                    b_next = df.loc[idx, 'Autocad_Measure']

                    # Perform linear interpolation
                    interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                    # print(f"Interpolated value for {target_value}: {interpolated_value}")
                    y_tag[i]=interpolated_value

        if len(x_datum) > 0:
            for i in range(0, len(x_datum)):
                target_value = x_datum[i]
                df = X_table
                # Find indices of bracketing points
                idx = df['Real_Measure'].searchsorted(target_value)

                # Check if target is within data range
                if idx == 0 or idx == len(df):
                    print("Target value outside data range for interpolation X_Datum ",target_value)
                else:
                    # Extract bracketed values
                    a_prev = df.loc[idx - 1, 'Real_Measure']
                    a_next = df.loc[idx, 'Real_Measure']
                    b_prev = df.loc[idx - 1, 'Autocad_Measure']
                    b_next = df.loc[idx, 'Autocad_Measure']

                    # Perform linear interpolation
                    interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                    # print(f"Interpolated value for {target_value}: {interpolated_value}")
                    x_datum[i] = interpolated_value

            for i in range(0, len(y_datum)):
                target_value = y_datum[i]
                df = Y_table
                # Find indices of bracketing points
                idx = df['Real_Measure'].searchsorted(target_value)

                # Check if target is within data range
                if idx == 0 or idx == len(df):
                    print("Target value outside data range for interpolation Y_Datum ",target_value)
                else:
                    # Extract bracketed values
                    a_prev = df.loc[idx - 1, 'Real_Measure']
                    a_next = df.loc[idx, 'Real_Measure']
                    b_prev = df.loc[idx - 1, 'Autocad_Measure']
                    b_next = df.loc[idx, 'Autocad_Measure']

                    # Perform linear interpolation
                    interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                    # print(f"Interpolated value for {target_value}: {interpolated_value}")
                    y_datum[i] = interpolated_value

    ############# BEGGIN THE PLOTTING

    Axis_Limit_x = Axis_Limit
    Axis_Limit_y = Axis_Limit

    if pattern_mode==1:
        x=x_tag
        y=y_tag
        x_tag=[]
        y_tag=[]

    layer_qty=0
    x_layer=[]
    y_layer=[]
    for Layer_i in Layer_Names:    # loop through all the vertices of the splines found
        for i in range(0, len(x)):
            if Layer_points[i]==Layer_i:
                x_layer.append(x[i])
                y_layer.append(y[i])
        Layer_Name_plot = Layer_i
        layer_qty = layer_qty + 1
        # Plotting the polygon
        plt.figure(figsize=(image_size, image_size))
        plt.scatter(x_layer, y_layer, marker='o', s=0.1, color='blue')  # 'bo-' means blue circles connected by lines
        plt.fill(x_layer, y_layer, facecolor=layer_color, alpha=1)  # Fill the polygon
        title = 'Layer '+str(layer_qty)
        plt.title(title)
        plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
        plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
        plt.style.use('dark_background')
        ax = plt.gca()
        ax.set_facecolor(background_color)
        ax.spines['top'].set_color(layer_color)
        ax.spines['bottom'].set_color(layer_color)
        ax.spines['right'].set_color(layer_color)
        ax.spines['left'].set_color(layer_color)
        ax.tick_params(axis='x', colors=layer_color)
        ax.tick_params(axis='y', colors=layer_color)
        ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
        plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
        plt.savefig('Layers.png')
        # add slide
        slide = prs.slides.add_slide(slide_layout)
        img_path = 'Layers.png'
        left = Inches(Left_centering)  # Adjust position as needed in icnhes
        top = Inches(Top_centering)  # Adjust position as needed in inches
        slide.shapes.add_picture(img_path, left, top)
        if close_image == 1:
            plt.close('all')
        #print(title)
        x_layer = []   ### reset x_layer
        y_layer = []   ### reset y_layer




    if len(x_tag)>0:
                x_layer=x_tag
                y_layer=y_tag
                Layer_Name_plot='TAGGING'

                # Plotting the polygon
                plt.figure(figsize=(image_size, image_size))
                plt.scatter(x_layer, y_layer,marker='x', s=3, color='blue') # 'bo-' means blue circles connected by lines
                title='Tagging'
                plt.title(title)
                plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
                plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
                plt.style.use('dark_background')
                ax = plt.gca()
                ax.set_facecolor(background_color)
                ax.spines['top'].set_color(layer_color)
                ax.spines['bottom'].set_color(layer_color)
                ax.spines['right'].set_color(layer_color)
                ax.spines['left'].set_color(layer_color)
                ax.tick_params(axis='x', colors=layer_color)
                ax.tick_params(axis='y', colors=layer_color)
                ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
                plt.savefig('Layers.png')
                # add slide
                slide = prs.slides.add_slide(slide_layout)
                img_path = 'Layers.png'
                left = Inches(Left_centering)  # Adjust position as needed in icnhes
                top = Inches(Top_centering)  # Adjust position as needed in inches
                slide.shapes.add_picture(img_path, left, top)
                if close_image==1:
                 plt.close('all')

                # prs.save('Layers.pptx')



    if len(x_datum)>0:
                x_layer=x_datum
                y_layer=y_datum
                Layer_Name_plot='DATUM LINE'

                # Plotting the polygon
                plt.figure(figsize=(image_size, image_size))
                plt.scatter(x_layer, y_layer,marker='.', s=5, color='blue') # 'bo-' means blue circles connected by lines
                title='Tagging'
                plt.title(title)
                plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
                plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
                plt.style.use('dark_background')
                ax = plt.gca()
                ax.set_facecolor(background_color)
                ax.spines['top'].set_color(layer_color)
                ax.spines['bottom'].set_color(layer_color)
                ax.spines['right'].set_color(layer_color)
                ax.spines['left'].set_color(layer_color)
                ax.tick_params(axis='x', colors=layer_color)
                ax.tick_params(axis='y', colors=layer_color)
                ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
                plt.savefig('Layers.png')
                # add slide
                slide = prs.slides.add_slide(slide_layout)
                img_path = 'Layers.png'
                left = Inches(Left_centering)  # Adjust position as needed in icnhes
                top = Inches(Top_centering)  # Adjust position as needed in inches
                slide.shapes.add_picture(img_path, left, top)
                if close_image==1:
                 plt.close('all')

                # prs.save('Layers.pptx')


    # Save the PowerPoint presentation
    Name0=name[:-4]
    Name2 = Name0 + '.pptx'
    out_Name = f"{Destination_path}\{Name2}"
    prs.save(out_Name)


    print('total layers = ',layer_qty)
def read_and_plot_layerv17(file_path,name,image_size,Left_centering,Top_centering,Axis_Limit,scale_mode,Reduce_factor,background_color,layer_color,close_image,pattern_mode,raspberry,Destination_path):
    ### VERSION USING LAYER NAMES (MUST HAVE COLOR ASSIGNED) AND DESFASE APPLIED, not correctly applied because the desfase should occur before applying the scale (interpolation), not after
    import ezdxf
    import pandas as pd
    import matplotlib.pyplot as plt
    from pptx import Presentation
    from pptx.util import Inches


    ##............READING THE DXF
    doc = ezdxf.readfile(file_path)

    # POWERPOINT SLIDE
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank slide layout

    # Extract entities (splines)
    msp = doc.modelspace()
    # for entity in msp:
    #     # Get the type of the entity
    #     print(entity.dxftype())
    splines = msp.query('SPLINE')
    polylines= msp.query('POLYLINE')
    lw_polylines= msp.query('LWPOLYLINE')
    i=0
    x=[]
    y=[]
    x_tag=[]
    y_tag=[]
    x_datum=[]
    y_datum=[]
    ####.................... Iterate over splines/polylines/lw-polyinies
    for spline in splines:
            # Check if the spline is a BSpline
            if spline.dxftype() == 'SPLINE':  # here we look for entity type SPLINE
                if spline.dxf.color==1 or spline.dxf.color==3:   # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline._control_points

                    # Print control points
                    for point in control_points:
                        i=i+1
                        #print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline


    for spline in polylines:
        # Check if the spline is a BSpline
            if spline.dxftype() == 'POLYLINE':  # here we look for entity type SPLINE
                if spline.dxf.color==1 or spline.dxf.color==3:   # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline.points()
                    # Print control points
                    for point in control_points:
                        i=i+1
                        #print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline

    qty_spline=0
    Layer_Names=[]
    Layer_points=[]
    for spline in lw_polylines:
        # Check if the spline is a BSpline
            if spline.dxftype() == 'LWPOLYLINE':  # here we look for entity type SPLINE
                 #print(spline.dxf.color)
                 layer_name = spline.dxf.layer  # Get the layer name
                 layer = doc.layers.get(layer_name)  # Get the layer object
                 color = layer.dxf.color  # Get the color from the layer
                 #print (color)
                 ######## GET THE GREEN A RED LAYERS
                 if color==1 or color==3:    # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline.get_points('xy')

                    # Print control points
                    for point in control_points:
                        i=i+1
                        #print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline
                        Layer_points.append(layer_name)
                    #####................................................
                    if qty_spline == 0:
                        Layer_Names.append(layer_name)
                    elif Layer_Names[-1] != layer_name:
                        Layer_Names.append(layer_name)
                    qty_spline=qty_spline+1

                 ################### GET THE TAGGING LINES
                 if color == 7 or color == 255:
                     control_points = spline.get_points('xy')
                     for point in control_points:
                            i=i+1
                            #print("Vertice:", point,'vertice number', i)
                            x_tag.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                            y_tag.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline

                 ################### GET THE DATUM LINES
                 if color == 141:
                     control_points = spline.get_points('xy')
                     for point in control_points:
                            i=i+1
                            #print("Vertice:", point,'vertice number', i)
                            x_datum.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                            y_datum.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline




            #####...................1st scale reduction  OPTIONAL...........................
    # Reduce_factor=1
    x = [i / Reduce_factor for i in x]
    y = [i / Reduce_factor for i in y]

    #### calculo desfases
    cx=min(x)
    cy=min(y)

    ####  tabla de calibracion
    excel_table_calib = r'C:\Users\Juan Pablo Lopez\PycharmProjects\ProjectJP\Table_Calib.xlsx'
    # .....LOAD THE TABLE CALIBRATION DATA.......................................................................................................
    if raspberry==1:
        X_table = pd.read_excel(excel_table_calib, sheet_name='X_axis_RP', header=0)
        Y_table = pd.read_excel(excel_table_calib, sheet_name='Y_axis_RP', header=0)
    else:
        X_table = pd.read_excel(excel_table_calib, sheet_name='X_axis', header=0)
        Y_table = pd.read_excel(excel_table_calib, sheet_name='Y_axis', header=0)

    if scale_mode==1:
        for i in range(0,len(x)):
            target_value = x[i]
            df=X_table
            # Find indices of bracketing points
            idx = df['Real_Measure'].searchsorted(target_value)

            # Check if target is within data range
            if idx == 0 or idx == len(x):
                print("Target value outside data range for interpolation")
            else:
                # Extract bracketed values
                a_prev = df.loc[idx - 1, 'Real_Measure']
                a_next = df.loc[idx, 'Real_Measure']
                b_prev = df.loc[idx - 1, 'Autocad_Measure']
                b_next = df.loc[idx, 'Autocad_Measure']

                # Perform linear interpolation
                interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                # print(f"Interpolated value for {target_value}: {interpolated_value}")
                x[i]=interpolated_value

        for i in range(0,len(y)):
            target_value = y[i]
            df=Y_table
            # Find indices of bracketing points
            idx = df['Real_Measure'].searchsorted(target_value)

            # Check if target is within data range
            if idx == 0 or idx == len(y):
                print("Target value outside data range for interpolation")
            else:
                # Extract bracketed values
                a_prev = df.loc[idx - 1, 'Real_Measure']
                a_next = df.loc[idx, 'Real_Measure']
                b_prev = df.loc[idx - 1, 'Autocad_Measure']
                b_next = df.loc[idx, 'Autocad_Measure']

                # Perform linear interpolation
                interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                # print(f"Interpolated value for {target_value}: {interpolated_value}")
                y[i]=interpolated_value

        if len(x_tag)>0:
            for i in range(0,len(x_tag)):
                target_value = x_tag[i]
                df=X_table
                # Find indices of bracketing points
                idx = df['Real_Measure'].searchsorted(target_value)

                # Check if target is within data range
                if idx == 0 or idx == len(x_tag):
                    print("Target value outside data range for interpolation")
                else:
                    # Extract bracketed values
                    a_prev = df.loc[idx - 1, 'Real_Measure']
                    a_next = df.loc[idx, 'Real_Measure']
                    b_prev = df.loc[idx - 1, 'Autocad_Measure']
                    b_next = df.loc[idx, 'Autocad_Measure']

                    # Perform linear interpolation
                    interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                    # print(f"Interpolated value for {target_value}: {interpolated_value}")
                    x_tag[i]=interpolated_value

            for i in range(0,len(y_tag)):
                target_value = y_tag[i]
                df=Y_table
                # Find indices of bracketing points
                idx = df['Real_Measure'].searchsorted(target_value)

                # Check if target is within data range
                if idx == 0 or idx == len(y_tag):
                    print("Target value outside data range for interpolation")
                else:
                    # Extract bracketed values
                    a_prev = df.loc[idx - 1, 'Real_Measure']
                    a_next = df.loc[idx, 'Real_Measure']
                    b_prev = df.loc[idx - 1, 'Autocad_Measure']
                    b_next = df.loc[idx, 'Autocad_Measure']

                    # Perform linear interpolation
                    interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                    # print(f"Interpolated value for {target_value}: {interpolated_value}")
                    y_tag[i]=interpolated_value

        if len(x_datum) > 0:
            for i in range(0, len(x_datum)):
                target_value = x_datum[i]
                df = X_table
                # Find indices of bracketing points
                idx = df['Real_Measure'].searchsorted(target_value)

                # Check if target is within data range
                if idx == 0 or idx == len(x_tag):
                    print("Target value outside data range for interpolation")
                else:
                    # Extract bracketed values
                    a_prev = df.loc[idx - 1, 'Real_Measure']
                    a_next = df.loc[idx, 'Real_Measure']
                    b_prev = df.loc[idx - 1, 'Autocad_Measure']
                    b_next = df.loc[idx, 'Autocad_Measure']

                    # Perform linear interpolation
                    interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                    # print(f"Interpolated value for {target_value}: {interpolated_value}")
                    x_datum[i] = interpolated_value

            for i in range(0, len(y_datum)):
                target_value = y_datum[i]
                df = Y_table
                # Find indices of bracketing points
                idx = df['Real_Measure'].searchsorted(target_value)

                # Check if target is within data range
                if idx == 0 or idx == len(y_tag):
                    print("Target value outside data range for interpolation")
                else:
                    # Extract bracketed values
                    a_prev = df.loc[idx - 1, 'Real_Measure']
                    a_next = df.loc[idx, 'Real_Measure']
                    b_prev = df.loc[idx - 1, 'Autocad_Measure']
                    b_next = df.loc[idx, 'Autocad_Measure']

                    # Perform linear interpolation
                    interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                    # print(f"Interpolated value for {target_value}: {interpolated_value}")
                    y_datum[i] = interpolated_value

            ###### apply scale to the desfase
            target_value = cx
            df = X_table
            # Find indices of bracketing points
            idx = df['Real_Measure'].searchsorted(target_value)

            # Check if target is within data range
            if idx == 0 or idx == len(x_tag):
                print("Target value outside data range for interpolation")
            else:
                # Extract bracketed values
                a_prev = df.loc[idx - 1, 'Real_Measure']
                a_next = df.loc[idx, 'Real_Measure']
                b_prev = df.loc[idx - 1, 'Autocad_Measure']
                b_next = df.loc[idx, 'Autocad_Measure']

                # Perform linear interpolation
                interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                # print(f"Interpolated value for {target_value}: {interpolated_value}")
                cx = interpolated_value


            target_value = cy
            df = Y_table
            # Find indices of bracketing points
            idx = df['Real_Measure'].searchsorted(target_value)

            # Check if target is within data range
            if idx == 0 or idx == len(x_tag):
                print("Target value outside data range for interpolation")
            else:
                # Extract bracketed values
                a_prev = df.loc[idx - 1, 'Real_Measure']
                a_next = df.loc[idx, 'Real_Measure']
                b_prev = df.loc[idx - 1, 'Autocad_Measure']
                b_next = df.loc[idx, 'Autocad_Measure']

                # Perform linear interpolation
                interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                # print(f"Interpolated value for {target_value}: {interpolated_value}")
                cy = interpolated_value



    #### APPLY THE DESFASE
    cx0=min(x)
    cy0=min(y)
    x = [i -cx0 + cx for i in x]  # BRING X VALUES TO THE ORIGIN and apply desfase
    y = [i -cy0 + cy for i in y]  # BRING Y VALUES TO THE ORIGIN and apply desfase
    x_tag = [i -cx0 + cx for i in x_tag]  # BRING X VALUES TO THE ORIGIN and apply desfase
    y_tag = [i -cy0 + cy for i in y_tag]  # BRING Y VALUES TO THE ORIGIN and apply desfase
    x_datum = [i -cx0 + cx for i in x_datum]  # BRING X VALUES TO THE ORIGIN and apply desfase
    y_datum = [i -cy0 + cy for i in y_datum]  # BRING Y VALUES TO THE ORIGIN and apply desfase

    ############# BEGGIN THE PLOTTING
    Axis_Limit_x = Axis_Limit
    Axis_Limit_y = Axis_Limit

    if pattern_mode==1:
        x=x_tag
        y=y_tag
        x_tag=[]
        y_tag=[]

    layer_qty=0
    x_layer=[]
    y_layer=[]
    for Layer_i in Layer_Names:    # loop through all the vertices of the splines found
        for i in range(0, len(x)):
            if Layer_points[i]==Layer_i:
                x_layer.append(x[i])
                y_layer.append(y[i])
        Layer_Name_plot = Layer_i
        layer_qty = layer_qty + 1
        # Plotting the polygon
        plt.figure(figsize=(image_size, image_size))
        plt.scatter(x_layer, y_layer, marker='o', s=0.1, color='blue')  # 'bo-' means blue circles connected by lines
        plt.fill(x_layer, y_layer, facecolor=layer_color, alpha=1)  # Fill the polygon
        title = 'Layer '+str(layer_qty)
        plt.title(title)
        plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
        plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
        plt.style.use('dark_background')
        ax = plt.gca()
        ax.set_facecolor(background_color)
        ax.spines['top'].set_color(layer_color)
        ax.spines['bottom'].set_color(layer_color)
        ax.spines['right'].set_color(layer_color)
        ax.spines['left'].set_color(layer_color)
        ax.tick_params(axis='x', colors=layer_color)
        ax.tick_params(axis='y', colors=layer_color)
        ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
        plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
        plt.savefig('Layers.png')
        # add slide
        slide = prs.slides.add_slide(slide_layout)
        img_path = 'Layers.png'
        left = Inches(Left_centering)  # Adjust position as needed in icnhes
        top = Inches(Top_centering)  # Adjust position as needed in inches
        slide.shapes.add_picture(img_path, left, top)
        if close_image == 1:
            plt.close('all')
        #print(title)
        x_layer = []   ### reset x_layer
        y_layer = []   ### reset y_layer




    if len(x_tag)>0:
                x_layer=x_tag
                y_layer=y_tag
                Layer_Name_plot='TAGGING'

                # Plotting the polygon
                plt.figure(figsize=(image_size, image_size))
                plt.scatter(x_layer, y_layer,marker='x', s=3, color='blue') # 'bo-' means blue circles connected by lines
                title='Tagging'
                plt.title(title)
                plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
                plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
                plt.style.use('dark_background')
                ax = plt.gca()
                ax.set_facecolor(background_color)
                ax.spines['top'].set_color(layer_color)
                ax.spines['bottom'].set_color(layer_color)
                ax.spines['right'].set_color(layer_color)
                ax.spines['left'].set_color(layer_color)
                ax.tick_params(axis='x', colors=layer_color)
                ax.tick_params(axis='y', colors=layer_color)
                ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
                plt.savefig('Layers.png')
                # add slide
                slide = prs.slides.add_slide(slide_layout)
                img_path = 'Layers.png'
                left = Inches(Left_centering)  # Adjust position as needed in icnhes
                top = Inches(Top_centering)  # Adjust position as needed in inches
                slide.shapes.add_picture(img_path, left, top)
                if close_image==1:
                 plt.close('all')

                # prs.save('Layers.pptx')



    if len(x_datum)>0:
                x_layer=x_datum
                y_layer=y_datum
                Layer_Name_plot='DATUM LINE'

                # Plotting the polygon
                plt.figure(figsize=(image_size, image_size))
                plt.scatter(x_layer, y_layer,marker='.', s=5, color='blue') # 'bo-' means blue circles connected by lines
                title='Tagging'
                plt.title(title)
                plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
                plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
                plt.style.use('dark_background')
                ax = plt.gca()
                ax.set_facecolor(background_color)
                ax.spines['top'].set_color(layer_color)
                ax.spines['bottom'].set_color(layer_color)
                ax.spines['right'].set_color(layer_color)
                ax.spines['left'].set_color(layer_color)
                ax.tick_params(axis='x', colors=layer_color)
                ax.tick_params(axis='y', colors=layer_color)
                ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
                plt.savefig('Layers.png')
                # add slide
                slide = prs.slides.add_slide(slide_layout)
                img_path = 'Layers.png'
                left = Inches(Left_centering)  # Adjust position as needed in icnhes
                top = Inches(Top_centering)  # Adjust position as needed in inches
                slide.shapes.add_picture(img_path, left, top)
                if close_image==1:
                 plt.close('all')

                # prs.save('Layers.pptx')


    # Save the PowerPoint presentation
    Name0=name[:-4]
    Name2 = Name0 + '.pptx'
    out_Name = f"{Destination_path}\{Name2}"
    prs.save(out_Name)


    print('total layers = ',layer_qty)
def read_and_plot_layerv18(file_path,name,image_size,Left_centering,Top_centering,Axis_Limit,scale_mode,Reduce_factor,background_color,layer_color,close_image,pattern_mode,raspberry,Destination_path):
    ### VERSION USING LAYER NAMES (MUST HAVE COLOR ASSIGNED), DESFASE APPLIED, FULL DATUM LINE DRAWN, still not working properly
    import ezdxf
    import pandas as pd
    import matplotlib.pyplot as plt
    from pptx import Presentation
    from pptx.util import Inches


    ##............READING THE DXF
    doc = ezdxf.readfile(file_path)

    # POWERPOINT SLIDE
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank slide layout

    # Extract entities (splines)
    msp = doc.modelspace()
    # for entity in msp:
    #     # Get the type of the entity
    #     print(entity.dxftype())
    splines = msp.query('SPLINE')
    polylines= msp.query('POLYLINE')
    lw_polylines= msp.query('LWPOLYLINE')
    i=0
    x=[]
    y=[]
    x_tag=[]
    y_tag=[]
    x_datum=[]
    y_datum=[]
    ####.................... Iterate over splines/polylines/lw-polyinies
    for spline in splines:
            # Check if the spline is a BSpline
            if spline.dxftype() == 'SPLINE':  # here we look for entity type SPLINE
                if spline.dxf.color==1 or spline.dxf.color==3:   # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline._control_points

                    # Print control points
                    for point in control_points:
                        i=i+1
                        #print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline


    for spline in polylines:
        # Check if the spline is a BSpline
            if spline.dxftype() == 'POLYLINE':  # here we look for entity type SPLINE
                if spline.dxf.color==1 or spline.dxf.color==3:   # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline.points()
                    # Print control points
                    for point in control_points:
                        i=i+1
                        #print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline

    qty_spline=0
    Layer_Names=[]
    Layer_points=[]
    for spline in lw_polylines:
        # Check if the spline is a BSpline
            if spline.dxftype() == 'LWPOLYLINE':  # here we look for entity type SPLINE
                 #print(spline.dxf.color)
                 layer_name = spline.dxf.layer  # Get the layer name
                 layer = doc.layers.get(layer_name)  # Get the layer object
                 color = layer.dxf.color  # Get the color from the layer
                 #print (color)
                 ######## GET THE GREEN A RED LAYERS
                 if color==1 or color==3:    # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline.get_points('xy')

                    # Print control points
                    for point in control_points:
                        i=i+1
                        #print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline
                        Layer_points.append(layer_name)
                    #####................................................
                    if qty_spline == 0:
                        Layer_Names.append(layer_name)
                    elif Layer_Names[-1] != layer_name:
                        Layer_Names.append(layer_name)
                    qty_spline=qty_spline+1

                 ################### GET THE TAGGING LINES
                 if color == 7 or color == 255:
                     control_points = spline.get_points('xy')
                     for point in control_points:
                            i=i+1
                            #print("Vertice:", point,'vertice number', i)
                            x_tag.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                            y_tag.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline

                 ################### GET THE DATUM LINES
                 if color == 141:
                     control_points = spline.get_points('xy')
                     for point in control_points:
                            i=i+1
                            #print("Vertice:", point,'vertice number', i)
                            x_datum.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                            y_datum.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline




            #####...................1st scale reduction  OPTIONAL...........................
    # Reduce_factor=1
    x = [i / Reduce_factor for i in x]
    y = [i / Reduce_factor for i in y]

    #### calculo desfases
    cx=min(x)
    cy=min(y)

    ####  tabla de calibracion
    excel_table_calib = r'C:\Users\Juan Pablo Lopez\PycharmProjects\ProjectJP\Table_Calib.xlsx'
    # .....LOAD THE TABLE CALIBRATION DATA.......................................................................................................
    if raspberry==1:
        X_table = pd.read_excel(excel_table_calib, sheet_name='X_axis_RP', header=0)
        Y_table = pd.read_excel(excel_table_calib, sheet_name='Y_axis_RP', header=0)
    else:
        X_table = pd.read_excel(excel_table_calib, sheet_name='X_axis', header=0)
        Y_table = pd.read_excel(excel_table_calib, sheet_name='Y_axis', header=0)

    if scale_mode==1:
        for i in range(0,len(x)):
            target_value = x[i]
            df=X_table
            # Find indices of bracketing points
            idx = df['Real_Measure'].searchsorted(target_value)

            # Check if target is within data range
            if idx == 0 or idx == len(x):
                print("Target value outside data range for interpolation")
            else:
                # Extract bracketed values
                a_prev = df.loc[idx - 1, 'Real_Measure']
                a_next = df.loc[idx, 'Real_Measure']
                b_prev = df.loc[idx - 1, 'Autocad_Measure']
                b_next = df.loc[idx, 'Autocad_Measure']

                # Perform linear interpolation
                interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                # print(f"Interpolated value for {target_value}: {interpolated_value}")
                x[i]=interpolated_value

        for i in range(0,len(y)):
            target_value = y[i]
            df=Y_table
            # Find indices of bracketing points
            idx = df['Real_Measure'].searchsorted(target_value)

            # Check if target is within data range
            if idx == 0 or idx == len(y):
                print("Target value outside data range for interpolation")
            else:
                # Extract bracketed values
                a_prev = df.loc[idx - 1, 'Real_Measure']
                a_next = df.loc[idx, 'Real_Measure']
                b_prev = df.loc[idx - 1, 'Autocad_Measure']
                b_next = df.loc[idx, 'Autocad_Measure']

                # Perform linear interpolation
                interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                # print(f"Interpolated value for {target_value}: {interpolated_value}")
                y[i]=interpolated_value

        if len(x_tag)>0:
            for i in range(0,len(x_tag)):
                target_value = x_tag[i]
                df=X_table
                # Find indices of bracketing points
                idx = df['Real_Measure'].searchsorted(target_value)

                # Check if target is within data range
                if idx == 0 or idx == len(x_tag):
                    print("Target value outside data range for interpolation")
                else:
                    # Extract bracketed values
                    a_prev = df.loc[idx - 1, 'Real_Measure']
                    a_next = df.loc[idx, 'Real_Measure']
                    b_prev = df.loc[idx - 1, 'Autocad_Measure']
                    b_next = df.loc[idx, 'Autocad_Measure']

                    # Perform linear interpolation
                    interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                    # print(f"Interpolated value for {target_value}: {interpolated_value}")
                    x_tag[i]=interpolated_value

            for i in range(0,len(y_tag)):
                target_value = y_tag[i]
                df=Y_table
                # Find indices of bracketing points
                idx = df['Real_Measure'].searchsorted(target_value)

                # Check if target is within data range
                if idx == 0 or idx == len(y_tag):
                    print("Target value outside data range for interpolation")
                else:
                    # Extract bracketed values
                    a_prev = df.loc[idx - 1, 'Real_Measure']
                    a_next = df.loc[idx, 'Real_Measure']
                    b_prev = df.loc[idx - 1, 'Autocad_Measure']
                    b_next = df.loc[idx, 'Autocad_Measure']

                    # Perform linear interpolation
                    interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                    # print(f"Interpolated value for {target_value}: {interpolated_value}")
                    y_tag[i]=interpolated_value

        if len(x_datum) > 0:
            for i in range(0, len(x_datum)):
                target_value = x_datum[i]
                df = X_table
                # Find indices of bracketing points
                idx = df['Real_Measure'].searchsorted(target_value)

                # Check if target is within data range
                if idx == 0 or idx == len(x_tag):
                    print("Target value outside data range for interpolation")
                else:
                    # Extract bracketed values
                    a_prev = df.loc[idx - 1, 'Real_Measure']
                    a_next = df.loc[idx, 'Real_Measure']
                    b_prev = df.loc[idx - 1, 'Autocad_Measure']
                    b_next = df.loc[idx, 'Autocad_Measure']

                    # Perform linear interpolation
                    interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                    # print(f"Interpolated value for {target_value}: {interpolated_value}")
                    x_datum[i] = interpolated_value

            for i in range(0, len(y_datum)):
                target_value = y_datum[i]
                df = Y_table
                # Find indices of bracketing points
                idx = df['Real_Measure'].searchsorted(target_value)

                # Check if target is within data range
                if idx == 0 or idx == len(y_tag):
                    print("Target value outside data range for interpolation")
                else:
                    # Extract bracketed values
                    a_prev = df.loc[idx - 1, 'Real_Measure']
                    a_next = df.loc[idx, 'Real_Measure']
                    b_prev = df.loc[idx - 1, 'Autocad_Measure']
                    b_next = df.loc[idx, 'Autocad_Measure']

                    # Perform linear interpolation
                    interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                    # print(f"Interpolated value for {target_value}: {interpolated_value}")
                    y_datum[i] = interpolated_value

            ###### apply scale to the desfase
            target_value = cx
            df = X_table
            # Find indices of bracketing points
            idx = df['Real_Measure'].searchsorted(target_value)

            # Check if target is within data range
            if idx == 0 or idx == len(x_tag):
                print("Target value outside data range for interpolation")
            else:
                # Extract bracketed values
                a_prev = df.loc[idx - 1, 'Real_Measure']
                a_next = df.loc[idx, 'Real_Measure']
                b_prev = df.loc[idx - 1, 'Autocad_Measure']
                b_next = df.loc[idx, 'Autocad_Measure']

                # Perform linear interpolation
                interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                # print(f"Interpolated value for {target_value}: {interpolated_value}")
                cx = interpolated_value


            target_value = cy
            df = Y_table
            # Find indices of bracketing points
            idx = df['Real_Measure'].searchsorted(target_value)

            # Check if target is within data range
            if idx == 0 or idx == len(x_tag):
                print("Target value outside data range for interpolation")
            else:
                # Extract bracketed values
                a_prev = df.loc[idx - 1, 'Real_Measure']
                a_next = df.loc[idx, 'Real_Measure']
                b_prev = df.loc[idx - 1, 'Autocad_Measure']
                b_next = df.loc[idx, 'Autocad_Measure']

                # Perform linear interpolation
                interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                # print(f"Interpolated value for {target_value}: {interpolated_value}")
                cy = interpolated_value



    #### APPLY THE DESFASE

    x = [i -min(x) + cx for i in x]  # BRING X VALUES TO THE ORIGIN and apply desfase
    y = [i -min(y) + cy for i in y]  # BRING Y VALUES TO THE ORIGIN and apply desfase
    x_tag = [i -min(x) + cx for i in x_tag]  # BRING X VALUES TO THE ORIGIN and apply desfase
    y_tag = [i -min(y) + cy for i in y_tag]  # BRING Y VALUES TO THE ORIGIN and apply desfase
    x_datum = [i -min(x) + cx for i in x_datum]  # BRING X VALUES TO THE ORIGIN and apply desfase
    y_datum = [i -min(y) + cy for i in y_datum]  # BRING Y VALUES TO THE ORIGIN and apply desfase

    ############# BEGGIN THE PLOTTING
    Axis_Limit_x = Axis_Limit
    Axis_Limit_y = Axis_Limit

    if pattern_mode==1:
        x=x_tag
        y=y_tag
        x_tag=[]
        y_tag=[]

    layer_qty=0
    x_layer=[]
    y_layer=[]
    for Layer_i in Layer_Names:    # loop through all the vertices of the splines found
        for i in range(0, len(x)):
            if Layer_points[i]==Layer_i:
                x_layer.append(x[i])
                y_layer.append(y[i])
        Layer_Name_plot = Layer_i
        layer_qty = layer_qty + 1
        # Plotting the polygon
        plt.figure(figsize=(image_size, image_size))
        plt.scatter(x_layer, y_layer, marker='o', s=0.1, color='blue')  # 'bo-' means blue circles connected by lines
        plt.fill(x_layer, y_layer, facecolor=layer_color, alpha=1)  # Fill the polygon
        title = 'Layer '+str(layer_qty)
        plt.title(title)
        plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
        plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
        plt.style.use('dark_background')
        ax = plt.gca()
        ax.set_facecolor(background_color)
        ax.spines['top'].set_color(layer_color)
        ax.spines['bottom'].set_color(layer_color)
        ax.spines['right'].set_color(layer_color)
        ax.spines['left'].set_color(layer_color)
        ax.tick_params(axis='x', colors=layer_color)
        ax.tick_params(axis='y', colors=layer_color)
        ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
        plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
        plt.savefig('Layers.png')
        # add slide
        slide = prs.slides.add_slide(slide_layout)
        img_path = 'Layers.png'
        left = Inches(Left_centering)  # Adjust position as needed in icnhes
        top = Inches(Top_centering)  # Adjust position as needed in inches
        slide.shapes.add_picture(img_path, left, top)
        if close_image == 1:
            plt.close('all')
        #print(title)
        x_layer = []   ### reset x_layer
        y_layer = []   ### reset y_layer




    if len(x_tag)>0:
                x_layer=x_tag
                y_layer=y_tag
                Layer_Name_plot='TAGGING'

                # Plotting the polygon
                plt.figure(figsize=(image_size, image_size))
                plt.scatter(x_layer, y_layer,marker='x', s=3, color='blue') # 'bo-' means blue circles connected by lines
                title='Tagging'
                plt.title(title)
                plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
                plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
                plt.style.use('dark_background')
                ax = plt.gca()
                ax.set_facecolor(background_color)
                ax.spines['top'].set_color(layer_color)
                ax.spines['bottom'].set_color(layer_color)
                ax.spines['right'].set_color(layer_color)
                ax.spines['left'].set_color(layer_color)
                ax.tick_params(axis='x', colors=layer_color)
                ax.tick_params(axis='y', colors=layer_color)
                ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
                plt.savefig('Layers.png')
                # add slide
                slide = prs.slides.add_slide(slide_layout)
                img_path = 'Layers.png'
                left = Inches(Left_centering)  # Adjust position as needed in icnhes
                top = Inches(Top_centering)  # Adjust position as needed in inches
                slide.shapes.add_picture(img_path, left, top)
                if close_image==1:
                 plt.close('all')

                # prs.save('Layers.pptx')



    if len(x_datum)>0:
                x_datum1=[]
                y_datum1=[]
                for i in range(0,len(x_datum)):
                    if i < len(x_datum)-1:
                        diff=abs(x_datum[i+1]-x_datum[i])
                        if diff<=1:
                            x_datum1.append(x_datum[i])
                            y_datum1.append(y_datum[i])
                            k=i+2
                x_datum2=x_datum[k:]
                y_datum2=y_datum[k:]

                ##### PLOT THE FIRST DATUM LINE
                Layer_Name_plot='DATUM LINE 1'
                x_layer=x_datum1
                y_layer=y_datum1
                # Plotting the polygon
                plt.figure(figsize=(image_size, image_size))
                plt.plot(x_layer, y_layer,color=layer_color, linestyle='-', linewidth=1, marker='o', markersize=0.1) # 'bo-' means blue circles connected by lines
                title='Tagging'
                plt.title(title)
                plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
                plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
                plt.style.use('dark_background')
                ax = plt.gca()
                ax.set_facecolor(background_color)
                ax.spines['top'].set_color(layer_color)
                ax.spines['bottom'].set_color(layer_color)
                ax.spines['right'].set_color(layer_color)
                ax.spines['left'].set_color(layer_color)
                ax.tick_params(axis='x', colors=layer_color)
                ax.tick_params(axis='y', colors=layer_color)
                ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
                plt.savefig('Layers.png')
                # add slide
                slide = prs.slides.add_slide(slide_layout)
                img_path = 'Layers.png'
                left = Inches(Left_centering)  # Adjust position as needed in icnhes
                top = Inches(Top_centering)  # Adjust position as needed in inches
                slide.shapes.add_picture(img_path, left, top)
                if close_image==1:
                 plt.close('all')
                # prs.save('Layers.pptx')

                ##### PLOT THE SECOND DATUM LINE
                Layer_Name_plot='DATUM LINE 2'
                x_layer=x_datum2
                y_layer=y_datum2
                # Plotting the polygon
                plt.figure(figsize=(image_size, image_size))
                plt.plot(x_layer, y_layer,color=layer_color, linestyle='-', linewidth=1, marker='o', markersize=0.1) # 'bo-' means blue circles connected by lines
                title='Tagging'
                plt.title(title)
                plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
                plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
                plt.style.use('dark_background')
                ax = plt.gca()
                ax.set_facecolor(background_color)
                ax.spines['top'].set_color(layer_color)
                ax.spines['bottom'].set_color(layer_color)
                ax.spines['right'].set_color(layer_color)
                ax.spines['left'].set_color(layer_color)
                ax.tick_params(axis='x', colors=layer_color)
                ax.tick_params(axis='y', colors=layer_color)
                ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
                plt.savefig('Layers.png')
                # add slide
                slide = prs.slides.add_slide(slide_layout)
                img_path = 'Layers.png'
                left = Inches(Left_centering)  # Adjust position as needed in icnhes
                top = Inches(Top_centering)  # Adjust position as needed in inches
                slide.shapes.add_picture(img_path, left, top)
                if close_image==1:
                 plt.close('all')
                # prs.save('Layers.pptx')



    # Save the PowerPoint presentation
    Name0=name[:-4]
    Name2 = Name0 + '.pptx'
    out_Name = f"{Destination_path}\{Name2}"
    prs.save(out_Name)


    print('total layers = ',layer_qty)
def read_and_plot_layerv19(file_path,name,image_size,Left_centering,Top_centering,Axis_Limit,scale_mode,Reduce_factor,background_color,layer_color,close_image,pattern_mode,raspberry,Destination_path,Axis_Limit_Y):
    ### V13 with separate x-axis and y-axis limit in the plot
    import ezdxf
    import pandas as pd
    import matplotlib.pyplot as plt
    from pptx import Presentation
    from pptx.util import Inches


    ##............READING THE DXF
    doc = ezdxf.readfile(file_path)

    # POWERPOINT SLIDE
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank slide layout

    # Extract entities (splines)
    msp = doc.modelspace()
    # for entity in msp:
    #     # Get the type of the entity
    #     print(entity.dxftype())
    splines = msp.query('SPLINE')
    polylines= msp.query('POLYLINE')
    lw_polylines= msp.query('LWPOLYLINE')
    i=0
    x=[]
    y=[]
    x_tag=[]
    y_tag=[]
    x_datum=[]
    y_datum=[]
    ####.................... Iterate over splines/polylines/lw-polyinies
    for spline in splines:
            # Check if the spline is a BSpline
            if spline.dxftype() == 'SPLINE':  # here we look for entity type SPLINE
                if spline.dxf.color==1 or spline.dxf.color==3:   # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline._control_points

                    # Print control points
                    for point in control_points:
                        i=i+1
                        #print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline


    for spline in polylines:
        # Check if the spline is a BSpline
            if spline.dxftype() == 'POLYLINE':  # here we look for entity type SPLINE
                if spline.dxf.color==1 or spline.dxf.color==3:   # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline.points()
                    # Print control points
                    for point in control_points:
                        i=i+1
                        #print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline

    qty_spline=0
    Layer_Names=[]
    Layer_points=[]
    for spline in lw_polylines:
        # Check if the spline is a BSpline
            if spline.dxftype() == 'LWPOLYLINE':  # here we look for entity type SPLINE
                 #print(spline.dxf.color)
                 layer_name = spline.dxf.layer  # Get the layer name
                 layer = doc.layers.get(layer_name)  # Get the layer object
                 color = layer.dxf.color  # Get the color from the layer
                 #print (color)
                 ######## GET THE GREEN A RED LAYERS
                 if color==1 or color==3:    # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline.get_points('xy')

                    # Print control points
                    for point in control_points:
                        i=i+1
                        #print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline
                        Layer_points.append(layer_name)
                    #####................................................
                    if qty_spline == 0:
                        Layer_Names.append(layer_name)
                    elif Layer_Names[-1] != layer_name:
                        Layer_Names.append(layer_name)
                    qty_spline=qty_spline+1

                 ################### GET THE TAGGING LINES
                 if color == 7 or color == 255:
                     control_points = spline.get_points('xy')
                     for point in control_points:
                            i=i+1
                            #print("Vertice:", point,'vertice number', i)
                            x_tag.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                            y_tag.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline

                 ################### GET THE DATUM LINES
                 if color == 141:
                     control_points = spline.get_points('xy')
                     for point in control_points:
                            i=i+1
                            #print("Vertice:", point,'vertice number', i)
                            x_datum.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                            y_datum.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline




            #####...................1st scale reduction  OPTIONAL...........................
    # Reduce_factor=1
    x = [i / Reduce_factor for i in x]
    y = [i / Reduce_factor for i in y]

    excel_table_calib = r'C:\Users\Juan Pablo Lopez\PycharmProjects\ProjectJP\Table_Calib.xlsx'
    # .....LOAD THE TABLE CALIBRATION DATA.......................................................................................................
    if raspberry==1:
        X_table = pd.read_excel(excel_table_calib, sheet_name='X_axis_RP', header=0)
        Y_table = pd.read_excel(excel_table_calib, sheet_name='Y_axis_RP', header=0)
    else:
        X_table = pd.read_excel(excel_table_calib, sheet_name='X_axis', header=0)
        Y_table = pd.read_excel(excel_table_calib, sheet_name='Y_axis', header=0)

    if scale_mode==1:
        for i in range(0,len(x)):
            target_value = x[i]
            df=X_table
            # Find indices of bracketing points
            idx = df['Real_Measure'].searchsorted(target_value)

            # Check if target is within data range
            if idx == 0 or idx == len(df):
                print("Target value outside data range for interpolation X ",target_value)
            else:
                # Extract bracketed values
                a_prev = df.loc[idx - 1, 'Real_Measure']
                a_next = df.loc[idx, 'Real_Measure']
                b_prev = df.loc[idx - 1, 'Autocad_Measure']
                b_next = df.loc[idx, 'Autocad_Measure']

                # Perform linear interpolation
                interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                # print(f"Interpolated value for {target_value}: {interpolated_value}")
                x[i]=interpolated_value

        for i in range(0,len(y)):
            target_value = y[i]
            df=Y_table
            # Find indices of bracketing points
            idx = df['Real_Measure'].searchsorted(target_value)

            # Check if target is within data range
            if idx == 0 or idx == len(df):
                print("Target value outside data range for interpolation Y ",target_value)
            else:
                # Extract bracketed values
                a_prev = df.loc[idx - 1, 'Real_Measure']
                a_next = df.loc[idx, 'Real_Measure']
                b_prev = df.loc[idx - 1, 'Autocad_Measure']
                b_next = df.loc[idx, 'Autocad_Measure']

                # Perform linear interpolation
                interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                # print(f"Interpolated value for {target_value}: {interpolated_value}")
                y[i]=interpolated_value

        if len(x_tag)>0:
            for i in range(0,len(x_tag)):
                target_value = x_tag[i]
                df=X_table
                # Find indices of bracketing points
                idx = df['Real_Measure'].searchsorted(target_value)

                # Check if target is within data range
                if idx == 0 or idx == len(df):
                    print("Target value outside data range for interpolation X_Tag ",target_value)
                else:
                    # Extract bracketed values
                    a_prev = df.loc[idx - 1, 'Real_Measure']
                    a_next = df.loc[idx, 'Real_Measure']
                    b_prev = df.loc[idx - 1, 'Autocad_Measure']
                    b_next = df.loc[idx, 'Autocad_Measure']

                    # Perform linear interpolation
                    interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                    # print(f"Interpolated value for {target_value}: {interpolated_value}")
                    x_tag[i]=interpolated_value

            for i in range(0,len(y_tag)):
                target_value = y_tag[i]
                df=Y_table
                # Find indices of bracketing points
                idx = df['Real_Measure'].searchsorted(target_value)

                # Check if target is within data range
                if idx == 0 or idx == len(df):
                    print("Target value outside data range for interpolation Y_tag ",target_value)
                else:
                    # Extract bracketed values
                    a_prev = df.loc[idx - 1, 'Real_Measure']
                    a_next = df.loc[idx, 'Real_Measure']
                    b_prev = df.loc[idx - 1, 'Autocad_Measure']
                    b_next = df.loc[idx, 'Autocad_Measure']

                    # Perform linear interpolation
                    interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                    # print(f"Interpolated value for {target_value}: {interpolated_value}")
                    y_tag[i]=interpolated_value

        if len(x_datum) > 0:
            for i in range(0, len(x_datum)):
                target_value = x_datum[i]
                df = X_table
                # Find indices of bracketing points
                idx = df['Real_Measure'].searchsorted(target_value)

                # Check if target is within data range
                if idx == 0 or idx == len(df):
                    print("Target value outside data range for interpolation X_Datum ",target_value)
                else:
                    # Extract bracketed values
                    a_prev = df.loc[idx - 1, 'Real_Measure']
                    a_next = df.loc[idx, 'Real_Measure']
                    b_prev = df.loc[idx - 1, 'Autocad_Measure']
                    b_next = df.loc[idx, 'Autocad_Measure']

                    # Perform linear interpolation
                    interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                    # print(f"Interpolated value for {target_value}: {interpolated_value}")
                    x_datum[i] = interpolated_value

            for i in range(0, len(y_datum)):
                target_value = y_datum[i]
                df = Y_table
                # Find indices of bracketing points
                idx = df['Real_Measure'].searchsorted(target_value)

                # Check if target is within data range
                if idx == 0 or idx == len(df):
                    print("Target value outside data range for interpolation Y_Datum ",target_value)
                else:
                    # Extract bracketed values
                    a_prev = df.loc[idx - 1, 'Real_Measure']
                    a_next = df.loc[idx, 'Real_Measure']
                    b_prev = df.loc[idx - 1, 'Autocad_Measure']
                    b_next = df.loc[idx, 'Autocad_Measure']

                    # Perform linear interpolation
                    interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                    # print(f"Interpolated value for {target_value}: {interpolated_value}")
                    y_datum[i] = interpolated_value

    ############# BEGGIN THE PLOTTING

    Axis_Limit_x = Axis_Limit
    Axis_Limit_y = Axis_Limit_Y

    if pattern_mode==1:
        x=x_tag
        y=y_tag
        x_tag=[]
        y_tag=[]

    layer_qty=0
    x_layer=[]
    y_layer=[]
    for Layer_i in Layer_Names:    # loop through all the vertices of the splines found
        for i in range(0, len(x)):
            if Layer_points[i]==Layer_i:
                x_layer.append(x[i])
                y_layer.append(y[i])
        Layer_Name_plot = Layer_i
        layer_qty = layer_qty + 1
        # Plotting the polygon
        plt.figure(figsize=(image_size, image_size))
        plt.scatter(x_layer, y_layer, marker='o', s=0.1, color='blue')  # 'bo-' means blue circles connected by lines
        plt.fill(x_layer, y_layer, facecolor=layer_color, alpha=1)  # Fill the polygon
        title = 'Layer '+str(layer_qty)
        plt.title(title)
        plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
        plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
        plt.style.use('dark_background')
        ax = plt.gca()
        ax.set_facecolor(background_color)
        ax.spines['top'].set_color(layer_color)
        ax.spines['bottom'].set_color(layer_color)
        ax.spines['right'].set_color(layer_color)
        ax.spines['left'].set_color(layer_color)
        ax.tick_params(axis='x', colors=layer_color)
        ax.tick_params(axis='y', colors=layer_color)
        ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
        plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
        plt.savefig('Layers.png')
        # add slide
        slide = prs.slides.add_slide(slide_layout)
        img_path = 'Layers.png'
        left = Inches(Left_centering)  # Adjust position as needed in icnhes
        top = Inches(Top_centering)  # Adjust position as needed in inches
        slide.shapes.add_picture(img_path, left, top)
        if close_image == 1:
            plt.close('all')
        #print(title)
        x_layer = []   ### reset x_layer
        y_layer = []   ### reset y_layer




    if len(x_tag)>0:
                x_layer=x_tag
                y_layer=y_tag
                Layer_Name_plot='TAGGING'

                # Plotting the polygon
                plt.figure(figsize=(image_size, image_size))
                plt.scatter(x_layer, y_layer,marker='x', s=3, color='blue') # 'bo-' means blue circles connected by lines
                title='Tagging'
                plt.title(title)
                plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
                plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
                plt.style.use('dark_background')
                ax = plt.gca()
                ax.set_facecolor(background_color)
                ax.spines['top'].set_color(layer_color)
                ax.spines['bottom'].set_color(layer_color)
                ax.spines['right'].set_color(layer_color)
                ax.spines['left'].set_color(layer_color)
                ax.tick_params(axis='x', colors=layer_color)
                ax.tick_params(axis='y', colors=layer_color)
                ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
                plt.savefig('Layers.png')
                # add slide
                slide = prs.slides.add_slide(slide_layout)
                img_path = 'Layers.png'
                left = Inches(Left_centering)  # Adjust position as needed in icnhes
                top = Inches(Top_centering)  # Adjust position as needed in inches
                slide.shapes.add_picture(img_path, left, top)
                if close_image==1:
                 plt.close('all')

                # prs.save('Layers.pptx')



    if len(x_datum)>0:
                x_layer=x_datum
                y_layer=y_datum
                Layer_Name_plot='DATUM LINE'

                # Plotting the polygon
                plt.figure(figsize=(image_size, image_size))
                plt.scatter(x_layer, y_layer,marker='.', s=5, color='blue') # 'bo-' means blue circles connected by lines
                title='Tagging'
                plt.title(title)
                plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
                plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
                plt.style.use('dark_background')
                ax = plt.gca()
                ax.set_facecolor(background_color)
                ax.spines['top'].set_color(layer_color)
                ax.spines['bottom'].set_color(layer_color)
                ax.spines['right'].set_color(layer_color)
                ax.spines['left'].set_color(layer_color)
                ax.tick_params(axis='x', colors=layer_color)
                ax.tick_params(axis='y', colors=layer_color)
                ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
                plt.savefig('Layers.png')
                # add slide
                slide = prs.slides.add_slide(slide_layout)
                img_path = 'Layers.png'
                left = Inches(Left_centering)  # Adjust position as needed in icnhes
                top = Inches(Top_centering)  # Adjust position as needed in inches
                slide.shapes.add_picture(img_path, left, top)
                if close_image==1:
                 plt.close('all')

                # prs.save('Layers.pptx')


    # Save the PowerPoint presentation
    Name0=name[:-4]
    Name2 = Name0 + '.pptx'
    out_Name = f"{Destination_path}\{Name2}"
    prs.save(out_Name)


    print('total layers = ',layer_qty)
def read_and_plot_layerv20(file_path,name,image_size,Left_centering,Top_centering,Axis_Limit,scale_mode,Reduce_factor,background_color,layer_color,close_image,pattern_mode,raspberry,Destination_path,Axis_Limit_Y):
    ### this is V19 with Surface Map calibration
    import ezdxf
    import pandas as pd
    import matplotlib.pyplot as plt
    from pptx import Presentation
    from pptx.util import Inches
    from dxf_surface_map import euclidean_distance,find_closest_tupleV2


    ##............READING THE DXF
    doc = ezdxf.readfile(file_path)

    # POWERPOINT SLIDE
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank slide layout

    # Extract entities (splines)
    msp = doc.modelspace()
    # for entity in msp:
    #     # Get the type of the entity
    #     print(entity.dxftype())
    splines = msp.query('SPLINE')
    polylines= msp.query('POLYLINE')
    lw_polylines= msp.query('LWPOLYLINE')
    i=0
    x=[]
    y=[]
    x_tag=[]
    y_tag=[]
    x_datum=[]
    y_datum=[]
    ####.................... Iterate over splines/polylines/lw-polyinies
    for spline in splines:
            # Check if the spline is a BSpline
            if spline.dxftype() == 'SPLINE':  # here we look for entity type SPLINE
                if spline.dxf.color==1 or spline.dxf.color==3:   # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline._control_points

                    # Print control points
                    for point in control_points:
                        i=i+1
                        #print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline


    for spline in polylines:
        # Check if the spline is a BSpline
            if spline.dxftype() == 'POLYLINE':  # here we look for entity type SPLINE
                if spline.dxf.color==1 or spline.dxf.color==3:   # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline.points()
                    # Print control points
                    for point in control_points:
                        i=i+1
                        #print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline

    qty_spline=0
    Layer_Names=[]
    Layer_points=[]
    for spline in lw_polylines:
        # Check if the spline is a BSpline
            if spline.dxftype() == 'LWPOLYLINE':  # here we look for entity type SPLINE
                 #print(spline.dxf.color)
                 layer_name = spline.dxf.layer  # Get the layer name
                 layer = doc.layers.get(layer_name)  # Get the layer object
                 color = layer.dxf.color  # Get the color from the layer
                 #print (color)
                 ######## GET THE GREEN A RED LAYERS
                 if color==1 or color==3:    # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline.get_points('xy')

                    # Print control points
                    for point in control_points:
                        i=i+1
                        #print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline
                        Layer_points.append(layer_name)
                    #####................................................
                    if qty_spline == 0:
                        Layer_Names.append(layer_name)
                    elif Layer_Names[-1] != layer_name:
                        Layer_Names.append(layer_name)
                    qty_spline=qty_spline+1

                 ################### GET THE TAGGING LINES
                 if color == 7 or color == 255:
                     control_points = spline.get_points('xy')
                     for point in control_points:
                            i=i+1
                            #print("Vertice:", point,'vertice number', i)
                            x_tag.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                            y_tag.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline

                 ################### GET THE DATUM LINES
                 if color == 141:
                     control_points = spline.get_points('xy')
                     for point in control_points:
                            i=i+1
                            #print("Vertice:", point,'vertice number', i)
                            x_datum.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                            y_datum.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline

    # .....LOAD THE TABLE CALIBRATION DATA.......................................................................................................
    excel_table_calib = r'C:\Users\Juan Pablo Lopez\PycharmProjects\ProjectJP\Table_Calib.xlsx'
    Surface_map = pd.read_excel(excel_table_calib, sheet_name='Surface', header=0)
    tuples_list_real=[]
    tuples_list_autocad=[]
    for i in range(0,len(Surface_map)):
        x_tup_autocad = Surface_map['Autocad_X'][i]
        y_tup_autocad = Surface_map['Autocad_Y'][i]
        tuples_list_autocad.append((x_tup_autocad,y_tup_autocad))
        x_tup_real = Surface_map['Real_X'][i]
        y_tup_real = Surface_map['Real_Y'][i]
        tuples_list_real.append((x_tup_real,y_tup_real))

    #####......................... APPLY THE CALIBRATION......................................
    if scale_mode==1:
        for i in range(0,len(x)):
            input_tuple = (x[i],y[i])
            xy_tup,idx_low,idx_high,tup_low,tup_high = find_closest_tupleV2(tuples_list_real,tuples_list_autocad,input_tuple)
            x[i] = xy_tup[0]
            y[i] = xy_tup[1]
            #print("Original ",input_tuple," Converted ",xy_tup," tup_low=",tup_low," idx low ",idx_low," tup_high=",tup_high, " idx high ",idx_high)

        if len(x_tag)>0:
            for i in range(0,len(x_tag)):
                input_tuple = (x_tag[i], y_tag[i])
                xy_tup, idx_low, idx_high, tup_low, tup_high = find_closest_tupleV2(tuples_list_real, tuples_list_autocad, input_tuple)
                x_tag[i] = xy_tup[0]
                y_tag[i] = xy_tup[1]
                # print("Original ", input_tuple, " Converted ", xy_tup, " tup_low=", tup_low, " idx low ", idx_low," tup_high=", tup_high, " idx high ", idx_high)

        if len(x_datum) > 0:
            for i in range(0, len(x_datum)):
                input_tuple = (x_datum[i], y_datum[i])
                xy_tup, idx_low, idx_high, tup_low, tup_high = find_closest_tupleV2(tuples_list_real, tuples_list_autocad, input_tuple)
                x_datum[i] = xy_tup[0]
                y_datum[i] = xy_tup[1]
                # print("Original ", input_tuple, " Converted ", xy_tup, " tup_low=", tup_low, " idx low ", idx_low," tup_high=", tup_high, " idx high ", idx_high)


    ############# BEGGIN THE PLOTTING...............................

    Axis_Limit_x = Axis_Limit
    Axis_Limit_y = Axis_Limit_Y

    if pattern_mode==1:
        x=x_tag
        y=y_tag
        x_tag=[]
        y_tag=[]

    layer_qty=0
    x_layer=[]
    y_layer=[]
    for Layer_i in Layer_Names:    # loop through all the vertices of the splines found
        for i in range(0, len(x)):
            if Layer_points[i]==Layer_i:
                x_layer.append(x[i])
                y_layer.append(y[i])
        Layer_Name_plot = Layer_i
        layer_qty = layer_qty + 1
        # Plotting the polygon
        plt.figure(figsize=(image_size, image_size))
        plt.scatter(x_layer, y_layer, marker='o', s=0.05, color='blue')  #  s=0.1  'bo-' means blue circles connected by lines
        plt.fill(x_layer, y_layer, facecolor=layer_color, alpha=1)  # Fill the polygon
        title = 'Layer '+str(layer_qty)
        plt.title(title)
        plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
        plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
        plt.style.use('dark_background')
        ax = plt.gca()
        ax.set_facecolor(background_color)
        ax.spines['top'].set_color(layer_color)
        ax.spines['bottom'].set_color(layer_color)
        ax.spines['right'].set_color(layer_color)
        ax.spines['left'].set_color(layer_color)
        ax.tick_params(axis='x', colors=layer_color)
        ax.tick_params(axis='y', colors=layer_color)
        ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
        plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
        plt.savefig('Layers.png')
        # add slide
        slide = prs.slides.add_slide(slide_layout)
        img_path = 'Layers.png'
        left = Inches(Left_centering)  # Adjust position as needed in icnhes
        top = Inches(Top_centering)  # Adjust position as needed in inches
        slide.shapes.add_picture(img_path, left, top)
        if close_image == 1:
            plt.close('all')
        #print(title)
        x_layer = []   ### reset x_layer
        y_layer = []   ### reset y_layer




    if len(x_tag)>0:
                x_layer=x_tag
                y_layer=y_tag
                Layer_Name_plot='TAGGING'

                # Plotting the polygon
                plt.figure(figsize=(image_size, image_size))
                plt.scatter(x_layer, y_layer,marker='x', s=3, color='blue') # 'bo-' means blue circles connected by lines
                title='Tagging'
                plt.title(title)
                plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
                plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
                plt.style.use('dark_background')
                ax = plt.gca()
                ax.set_facecolor(background_color)
                ax.spines['top'].set_color(layer_color)
                ax.spines['bottom'].set_color(layer_color)
                ax.spines['right'].set_color(layer_color)
                ax.spines['left'].set_color(layer_color)
                ax.tick_params(axis='x', colors=layer_color)
                ax.tick_params(axis='y', colors=layer_color)
                ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
                plt.savefig('Layers.png')
                # add slide
                slide = prs.slides.add_slide(slide_layout)
                img_path = 'Layers.png'
                left = Inches(Left_centering)  # Adjust position as needed in icnhes
                top = Inches(Top_centering)  # Adjust position as needed in inches
                slide.shapes.add_picture(img_path, left, top)
                if close_image==1:
                 plt.close('all')

                # prs.save('Layers.pptx')



    if len(x_datum)>0:
                x_layer=x_datum
                y_layer=y_datum
                Layer_Name_plot='DATUM LINE'

                # Plotting the polygon
                plt.figure(figsize=(image_size, image_size))
                plt.scatter(x_layer, y_layer,marker='.', s=5, color='blue') # 'bo-' means blue circles connected by lines
                title='Tagging'
                plt.title(title)
                plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
                plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
                plt.style.use('dark_background')
                ax = plt.gca()
                ax.set_facecolor(background_color)
                ax.spines['top'].set_color(layer_color)
                ax.spines['bottom'].set_color(layer_color)
                ax.spines['right'].set_color(layer_color)
                ax.spines['left'].set_color(layer_color)
                ax.tick_params(axis='x', colors=layer_color)
                ax.tick_params(axis='y', colors=layer_color)
                ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
                plt.savefig('Layers.png')
                # add slide
                slide = prs.slides.add_slide(slide_layout)
                img_path = 'Layers.png'
                left = Inches(Left_centering)  # Adjust position as needed in icnhes
                top = Inches(Top_centering)  # Adjust position as needed in inches
                slide.shapes.add_picture(img_path, left, top)
                if close_image==1:
                 plt.close('all')

                # prs.save('Layers.pptx')


    # Save the PowerPoint presentation
    Name0=name[:-4]
    Name2 = Name0 + '.pptx'
    out_Name = f"{Destination_path}\{Name2}"
    prs.save(out_Name)


    print('total layers = ',layer_qty)
def read_and_plot_layerv21(file_path,name,image_size,Left_centering,Top_centering,Axis_Limit,scale_mode,Reduce_factor,background_color,layer_color,close_image,pattern_mode,raspberry,Destination_path,Axis_Limit_Y):
    ### this is V19 with Surface Map calibration and stacks re-positioning (offset)
    import ezdxf
    import pandas as pd
    import matplotlib.pyplot as plt
    from pptx import Presentation
    from pptx.util import Inches
    from dxf_surface_map import euclidean_distance,find_closest_tupleV2


    ##............READING THE DXF
    doc = ezdxf.readfile(file_path)

    # POWERPOINT SLIDE
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank slide layout

    # Extract entities (splines)
    msp = doc.modelspace()
    for entity in msp:
        # Get the type of the entity
        print(entity.dxftype())
    splines = msp.query('SPLINE')
    polylines= msp.query('POLYLINE')
    lw_polylines= msp.query('LWPOLYLINE')
    i=0
    x=[]
    y=[]
    x_tag=[]
    y_tag=[]
    x_datum=[]
    y_datum=[]
    ####.................... Iterate over splines/polylines/lw-polyinies
    for spline in splines:
            # Check if the spline is a BSpline
            if spline.dxftype() == 'SPLINE':  # here we look for entity type SPLINE
                if spline.dxf.color==1 or spline.dxf.color==3:   # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline._control_points

                    # Print control points
                    for point in control_points:
                        i=i+1
                        #print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline


    for spline in polylines:
        # Check if the spline is a BSpline
            if spline.dxftype() == 'POLYLINE':  # here we look for entity type SPLINE
                if spline.dxf.color==1 or spline.dxf.color==3:   # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline.points()
                    # Print control points
                    for point in control_points:
                        i=i+1
                        #print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline

    qty_spline=0
    Layer_Names=[]
    Layer_points=[]
    for spline in lw_polylines:
        # Check if the spline is a BSpline
            if spline.dxftype() == 'LWPOLYLINE':  # here we look for entity type SPLINE
                 #print(spline.dxf.color)
                 layer_name = spline.dxf.layer  # Get the layer name
                 layer = doc.layers.get(layer_name)  # Get the layer object
                 color = layer.dxf.color  # Get the color from the layer
                 print (layer_name," color ",color)
                 ######## GET THE GREEN A RED LAYERS
                 if color==1 or color==3:    # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline.get_points('xy')

                    # Print control points
                    for point in control_points:
                        i=i+1
                        #print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline
                        Layer_points.append(layer_name)
                    #####................................................
                    if qty_spline == 0:
                        Layer_Names.append(layer_name)
                    elif Layer_Names[-1] != layer_name:
                        Layer_Names.append(layer_name)
                    qty_spline=qty_spline+1

                 ################### GET THE TAGGING LINES
                 if color == 7 or color == 255:
                     control_points = spline.get_points('xy')
                     for point in control_points:
                            i=i+1
                            #print("Vertice:", point,'vertice number', i)
                            x_tag.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                            y_tag.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline

                 ################### GET THE DATUM LINES
                 if color == 141:
                     control_points = spline.get_points('xy')
                     for point in control_points:
                            i=i+1
                            #print("Vertice:", point,'vertice number', i)
                            x_datum.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                            y_datum.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline



    #### BRING TO THE ORIGIN
    cx1=min(x)
    cy1=min(y)
    x = [i -cx1 for i in x]  # BRING X VALUES TO THE ORIGIN
    y = [i -cy1 for i in y]  # BRING Y VALUES TO THE ORIGIN
    x_tag = [i -cx1 for i in x_tag]  # BRING X VALUES TO THE ORIGIN
    y_tag = [i -cy1 for i in y_tag]  # BRING Y VALUES TO THE ORIGIN
    x_datum = [i -cx1 for i in x_datum]  # BRING X VALUES TO THE ORIGIN
    y_datum = [i -cy1 for i in y_datum]  # BRING Y VALUES TO THE ORIGIN


    #### APPLY THE OFFSET (DESFASE)
    x_lim=7040
    y_lim=2330
    cx=x_lim-max(x)
    cy=y_lim-max(y)
    x = [i + cx for i in x]  # apply offset desfase
    y = [i + cy for i in y]  # apply offset desfase
    x_tag = [i + cx for i in x_tag]  # apply offset desfase
    y_tag = [i + cy for i in y_tag]  # apply offset desfase
    x_datum = [i + cx for i in x_datum]  # apply offset desfase
    y_datum = [i + cy for i in y_datum]  # apply offset desfase

    # .....LOAD THE TABLE CALIBRATION DATA.......................................................................................................
    excel_table_calib = r'C:\Users\Juan Pablo Lopez\PycharmProjects\ProjectJP\Table_Calib.xlsx'
    Surface_map = pd.read_excel(excel_table_calib, sheet_name='Surface', header=0)
    tuples_list_real=[]
    tuples_list_autocad=[]
    for i in range(0,len(Surface_map)):
        x_tup_autocad = Surface_map['Autocad_X'][i]
        y_tup_autocad = Surface_map['Autocad_Y'][i]
        tuples_list_autocad.append((x_tup_autocad,y_tup_autocad))
        x_tup_real = Surface_map['Real_X'][i]
        y_tup_real = Surface_map['Real_Y'][i]
        tuples_list_real.append((x_tup_real,y_tup_real))

    #####......................... APPLY THE CALIBRATION......................................
    if scale_mode==1:
        for i in range(0,len(x)):
            input_tuple = (x[i],y[i])
            xy_tup,idx_low,idx_high,tup_low,tup_high = find_closest_tupleV2(tuples_list_real,tuples_list_autocad,input_tuple)
            x[i] = xy_tup[0]
            y[i] = xy_tup[1]
            #print("Original ",input_tuple," Converted ",xy_tup," tup_low=",tup_low," idx low ",idx_low," tup_high=",tup_high, " idx high ",idx_high)

        if len(x_tag)>0:
            for i in range(0,len(x_tag)):
                input_tuple = (x_tag[i], y_tag[i])
                # print(i,input_tuple)
                xy_tup, idx_low, idx_high, tup_low, tup_high = find_closest_tupleV2(tuples_list_real, tuples_list_autocad, input_tuple)
                x_tag[i] = xy_tup[0]
                y_tag[i] = xy_tup[1]
                # print("Original ", input_tuple, " Converted ", xy_tup, " tup_low=", tup_low, " idx low ", idx_low," tup_high=", tup_high, " idx high ", idx_high)

        if len(x_datum) > 0:
            for i in range(0, len(x_datum)):
                input_tuple = (x_datum[i], y_datum[i])
                xy_tup, idx_low, idx_high, tup_low, tup_high = find_closest_tupleV2(tuples_list_real, tuples_list_autocad, input_tuple)
                x_datum[i] = xy_tup[0]
                y_datum[i] = xy_tup[1]
                # print("Original ", input_tuple, " Converted ", xy_tup, " tup_low=", tup_low, " idx low ", idx_low," tup_high=", tup_high, " idx high ", idx_high)


    ############# BEGGIN THE PLOTTING...............................

    Axis_Limit_x = Axis_Limit
    Axis_Limit_y = Axis_Limit_Y

    if pattern_mode==1:
        x=x_tag
        y=y_tag
        x_tag=[]
        y_tag=[]

    layer_qty=0
    x_layer=[]
    y_layer=[]
    for Layer_i in Layer_Names:    # loop through all the vertices of the splines found
        for i in range(0, len(x)):
            if Layer_points[i]==Layer_i:
                x_layer.append(x[i])
                y_layer.append(y[i])
        Layer_Name_plot = Layer_i
        layer_qty = layer_qty + 1
        # Plotting the polygon
        plt.figure(figsize=(image_size, image_size))
        plt.scatter(x_layer, y_layer, marker='o', s=0.05, color='blue')  #  s=0.1  'bo-' means blue circles connected by lines
        plt.fill(x_layer, y_layer, facecolor=layer_color, alpha=1)  # Fill the polygon
        title = 'Layer '+str(layer_qty)
        plt.title(title)
        plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
        plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
        plt.style.use('dark_background')
        ax = plt.gca()
        ax.set_facecolor(background_color)
        ax.spines['top'].set_color(layer_color)
        ax.spines['bottom'].set_color(layer_color)
        ax.spines['right'].set_color(layer_color)
        ax.spines['left'].set_color(layer_color)
        ax.tick_params(axis='x', colors=layer_color)
        ax.tick_params(axis='y', colors=layer_color)
        ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
        plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
        plt.savefig('Layers.png')
        # add slide
        slide = prs.slides.add_slide(slide_layout)
        img_path = 'Layers.png'
        left = Inches(Left_centering)  # Adjust position as needed in icnhes
        top = Inches(Top_centering)  # Adjust position as needed in inches
        slide.shapes.add_picture(img_path, left, top)
        if close_image == 1:
            plt.close('all')
        #print(title)
        x_layer = []   ### reset x_layer
        y_layer = []   ### reset y_layer




    if len(x_tag)>0:
                x_layer=x_tag
                y_layer=y_tag
                Layer_Name_plot='TAGGING'

                # Plotting the polygon
                plt.figure(figsize=(image_size, image_size))
                plt.scatter(x_layer, y_layer,marker='x', s=3, color='blue') # 'bo-' means blue circles connected by lines
                title='Tagging'
                plt.title(title)
                plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
                plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
                plt.style.use('dark_background')
                ax = plt.gca()
                ax.set_facecolor(background_color)
                ax.spines['top'].set_color(layer_color)
                ax.spines['bottom'].set_color(layer_color)
                ax.spines['right'].set_color(layer_color)
                ax.spines['left'].set_color(layer_color)
                ax.tick_params(axis='x', colors=layer_color)
                ax.tick_params(axis='y', colors=layer_color)
                ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
                plt.savefig('Layers.png')
                # add slide
                slide = prs.slides.add_slide(slide_layout)
                img_path = 'Layers.png'
                left = Inches(Left_centering)  # Adjust position as needed in icnhes
                top = Inches(Top_centering)  # Adjust position as needed in inches
                slide.shapes.add_picture(img_path, left, top)
                if close_image==1:
                 plt.close('all')

                # prs.save('Layers.pptx')



    if len(x_datum)>0:
                x_layer=x_datum
                y_layer=y_datum
                Layer_Name_plot='DATUM LINE'

                # Plotting the polygon
                plt.figure(figsize=(image_size, image_size))
                plt.scatter(x_layer, y_layer,marker='.', s=5, color='blue') # 'bo-' means blue circles connected by lines
                title='Tagging'
                plt.title(title)
                plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
                plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
                plt.style.use('dark_background')
                ax = plt.gca()
                ax.set_facecolor(background_color)
                ax.spines['top'].set_color(layer_color)
                ax.spines['bottom'].set_color(layer_color)
                ax.spines['right'].set_color(layer_color)
                ax.spines['left'].set_color(layer_color)
                ax.tick_params(axis='x', colors=layer_color)
                ax.tick_params(axis='y', colors=layer_color)
                ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
                plt.savefig('Layers.png')
                # add slide
                slide = prs.slides.add_slide(slide_layout)
                img_path = 'Layers.png'
                left = Inches(Left_centering)  # Adjust position as needed in icnhes
                top = Inches(Top_centering)  # Adjust position as needed in inches
                slide.shapes.add_picture(img_path, left, top)
                if close_image==1:
                 plt.close('all')

                # prs.save('Layers.pptx')


    # Save the PowerPoint presentation
    Name0=name[:-4]
    Name2 = Name0 + '.pptx'
    out_Name = f"{Destination_path}\{Name2}"
    prs.save(out_Name)


    print('total layers = ',layer_qty)
def read_and_plot_layerv22(file_path,name,image_size,Left_centering,Top_centering,Axis_Limit,scale_mode,Reduce_factor,background_color,layer_color,close_image,pattern_mode,raspberry,Destination_path,Axis_Limit_Y):
    ### this is V20 with Surface Map calibration, stacks re-positioning (offset) and black screen overlay, not working yet
    import ezdxf
    import pandas as pd
    import matplotlib.pyplot as plt
    from pptx import Presentation
    from pptx.util import Inches
    from dxf_surface_map import euclidean_distance,find_closest_tupleV2
    from PIL import Image


    ##............READING THE DXF
    doc = ezdxf.readfile(file_path)

    # POWERPOINT SLIDE
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank slide layout

    # Extract entities (splines)
    msp = doc.modelspace()
    # for entity in msp:
    #     # Get the type of the entity
    #     print(entity.dxftype())
    splines = msp.query('SPLINE')
    polylines= msp.query('POLYLINE')
    lw_polylines= msp.query('LWPOLYLINE')
    i=0
    x=[]
    y=[]
    x_tag=[]
    y_tag=[]
    x_datum=[]
    y_datum=[]
    ####.................... Iterate over splines/polylines/lw-polyinies
    for spline in splines:
            # Check if the spline is a BSpline
            if spline.dxftype() == 'SPLINE':  # here we look for entity type SPLINE
                if spline.dxf.color==1 or spline.dxf.color==3:   # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline._control_points

                    # Print control points
                    for point in control_points:
                        i=i+1
                        #print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline


    for spline in polylines:
        # Check if the spline is a BSpline
            if spline.dxftype() == 'POLYLINE':  # here we look for entity type SPLINE
                if spline.dxf.color==1 or spline.dxf.color==3:   # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline.points()
                    # Print control points
                    for point in control_points:
                        i=i+1
                        #print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline

    qty_spline=0
    Layer_Names=[]
    Layer_points=[]
    for spline in lw_polylines:
        # Check if the spline is a BSpline
            if spline.dxftype() == 'LWPOLYLINE':  # here we look for entity type SPLINE
                 #print(spline.dxf.color)
                 layer_name = spline.dxf.layer  # Get the layer name
                 layer = doc.layers.get(layer_name)  # Get the layer object
                 color = layer.dxf.color  # Get the color from the layer
                 #print (color)
                 ######## GET THE GREEN A RED LAYERS
                 if color==1 or color==3:    # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline.get_points('xy')

                    # Print control points
                    for point in control_points:
                        i=i+1
                        #print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline
                        Layer_points.append(layer_name)
                    #####................................................
                    if qty_spline == 0:
                        Layer_Names.append(layer_name)
                    elif Layer_Names[-1] != layer_name:
                        Layer_Names.append(layer_name)
                    qty_spline=qty_spline+1

                 ################### GET THE TAGGING LINES
                 if color == 7 or color == 255:
                     control_points = spline.get_points('xy')
                     for point in control_points:
                            i=i+1
                            #print("Vertice:", point,'vertice number', i)
                            x_tag.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                            y_tag.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline

                 ################### GET THE DATUM LINES
                 if color == 141:
                     control_points = spline.get_points('xy')
                     for point in control_points:
                            i=i+1
                            #print("Vertice:", point,'vertice number', i)
                            x_datum.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                            y_datum.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline



    #### BRING TO THE ORIGIN
    cx1=min(x)
    cy1=min(y)
    x = [i -cx1 for i in x]  # BRING X VALUES TO THE ORIGIN
    y = [i -cy1 for i in y]  # BRING Y VALUES TO THE ORIGIN
    x_tag = [i -cx1 for i in x_tag]  # BRING X VALUES TO THE ORIGIN
    y_tag = [i -cy1 for i in y_tag]  # BRING Y VALUES TO THE ORIGIN
    x_datum = [i -cx1 for i in x_datum]  # BRING X VALUES TO THE ORIGIN
    y_datum = [i -cy1 for i in y_datum]  # BRING Y VALUES TO THE ORIGIN


    #### APPLY THE OFFSET (DESFASE)
    x_lim=7040
    y_lim=2330
    cx=x_lim-max(x)
    cy=y_lim-max(y)
    x = [i + cx for i in x]  # apply offset desfase
    y = [i + cy for i in y]  # apply offset desfase
    x_tag = [i + cx for i in x_tag]  # apply offset desfase
    y_tag = [i + cy for i in y_tag]  # apply offset desfase
    x_datum = [i + cx for i in x_datum]  # apply offset desfase
    y_datum = [i + cy for i in y_datum]  # apply offset desfase

    # .....LOAD THE TABLE CALIBRATION DATA.......................................................................................................
    excel_table_calib = r'C:\Users\Juan Pablo Lopez\PycharmProjects\ProjectJP\Table_Calib.xlsx'
    Surface_map = pd.read_excel(excel_table_calib, sheet_name='Surface', header=0)
    tuples_list_real=[]
    tuples_list_autocad=[]
    for i in range(0,len(Surface_map)):
        x_tup_autocad = Surface_map['Autocad_X'][i]
        y_tup_autocad = Surface_map['Autocad_Y'][i]
        tuples_list_autocad.append((x_tup_autocad,y_tup_autocad))
        x_tup_real = Surface_map['Real_X'][i]
        y_tup_real = Surface_map['Real_Y'][i]
        tuples_list_real.append((x_tup_real,y_tup_real))

    #####......................... APPLY THE CALIBRATION......................................
    if scale_mode==1:
        for i in range(0,len(x)):
            input_tuple = (x[i],y[i])
            xy_tup,idx_low,idx_high,tup_low,tup_high = find_closest_tupleV2(tuples_list_real,tuples_list_autocad,input_tuple)
            x[i] = xy_tup[0]
            y[i] = xy_tup[1]
            #print("Original ",input_tuple," Converted ",xy_tup," tup_low=",tup_low," idx low ",idx_low," tup_high=",tup_high, " idx high ",idx_high)

        if len(x_tag)>0:
            for i in range(0,len(x_tag)):
                input_tuple = (x_tag[i], y_tag[i])
                #print(i,input_tuple)
                xy_tup, idx_low, idx_high, tup_low, tup_high = find_closest_tupleV2(tuples_list_real, tuples_list_autocad, input_tuple)
                x_tag[i] = xy_tup[0]
                y_tag[i] = xy_tup[1]
                # print("Original ", input_tuple, " Converted ", xy_tup, " tup_low=", tup_low, " idx low ", idx_low," tup_high=", tup_high, " idx high ", idx_high)

        if len(x_datum) > 0:
            for i in range(0, len(x_datum)):
                input_tuple = (x_datum[i], y_datum[i])
                xy_tup, idx_low, idx_high, tup_low, tup_high = find_closest_tupleV2(tuples_list_real, tuples_list_autocad, input_tuple)
                x_datum[i] = xy_tup[0]
                y_datum[i] = xy_tup[1]
                # print("Original ", input_tuple, " Converted ", xy_tup, " tup_low=", tup_low, " idx low ", idx_low," tup_high=", tup_high, " idx high ", idx_high)


    ############# BEGGIN THE PLOTTING...............................

    Axis_Limit_x = Axis_Limit
    Axis_Limit_y = Axis_Limit_Y

    if pattern_mode==1:
        x=x_tag
        y=y_tag
        x_tag=[]
        y_tag=[]

    layer_qty=0
    x_layer=[]
    y_layer=[]
    for Layer_i in Layer_Names:    # loop through all the vertices of the splines found
        for i in range(0, len(x)):
            if Layer_points[i]==Layer_i:
                x_layer.append(x[i])
                y_layer.append(y[i])
        Layer_Name_plot = Layer_i
        layer_qty = layer_qty + 1
        # Plotting the polygon
        plt.figure(figsize=(image_size, image_size))
        plt.scatter(x_layer, y_layer, marker='o', s=0.05, color='blue')  #  s=0.1  'bo-' means blue circles connected by lines
        plt.fill(x_layer, y_layer, facecolor=layer_color, alpha=1)  # Fill the polygon
        title = 'Layer '+str(layer_qty)
        plt.title(title)
        plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
        plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
        plt.style.use('dark_background')
        ax = plt.gca()
        ax.set_facecolor(background_color)
        ax.spines['top'].set_color(layer_color)
        ax.spines['bottom'].set_color(layer_color)
        ax.spines['right'].set_color(layer_color)
        ax.spines['left'].set_color(layer_color)
        ax.tick_params(axis='x', colors=layer_color)
        ax.tick_params(axis='y', colors=layer_color)
        ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
        plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
        plt.savefig('Layers.png')

        #####..... overlay the black screen in the background
        base_image = Image.open("Layers.png")
        ###................X AXIS OVERLAY....................
        # Open the image to be overlayed
        overlay_image = Image.open("overlay_black_screen.PNG")
        # Define the position to overlay the image
        x_position = 1000  # 142
        y_position = 1000  # 1050
        # Paste the overlay image onto the base image
        base_image.paste(overlay_image, (x_position, y_position), overlay_image)
        # Save or display the resulting image
        base_image.save("Layers.png")
        #####...........

        # add slide
        slide = prs.slides.add_slide(slide_layout)
        img_path = 'Layers.png'
        left = Inches(Left_centering)  # Adjust position as needed in icnhes
        top = Inches(Top_centering)  # Adjust position as needed in inches
        slide.shapes.add_picture(img_path, left, top)
        if close_image == 1:
            plt.close('all')
        #print(title)
        x_layer = []   ### reset x_layer
        y_layer = []   ### reset y_layer




    if len(x_tag)>0:
                x_layer=x_tag
                y_layer=y_tag
                Layer_Name_plot='TAGGING'

                # Plotting the polygon
                plt.figure(figsize=(image_size, image_size))
                plt.scatter(x_layer, y_layer,marker='x', s=3, color='blue') # 'bo-' means blue circles connected by lines
                title='Tagging'
                plt.title(title)
                plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
                plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
                plt.style.use('dark_background')
                ax = plt.gca()
                ax.set_facecolor(background_color)
                ax.spines['top'].set_color(layer_color)
                ax.spines['bottom'].set_color(layer_color)
                ax.spines['right'].set_color(layer_color)
                ax.spines['left'].set_color(layer_color)
                ax.tick_params(axis='x', colors=layer_color)
                ax.tick_params(axis='y', colors=layer_color)
                ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
                plt.savefig('Layers.png')
                # add slide
                slide = prs.slides.add_slide(slide_layout)
                img_path = 'Layers.png'
                left = Inches(Left_centering)  # Adjust position as needed in icnhes
                top = Inches(Top_centering)  # Adjust position as needed in inches
                slide.shapes.add_picture(img_path, left, top)
                if close_image==1:
                 plt.close('all')

                # prs.save('Layers.pptx')



    if len(x_datum)>0:
                x_layer=x_datum
                y_layer=y_datum
                Layer_Name_plot='DATUM LINE'

                # Plotting the polygon
                plt.figure(figsize=(image_size, image_size))
                plt.scatter(x_layer, y_layer,marker='.', s=5, color='blue') # 'bo-' means blue circles connected by lines
                title='Tagging'
                plt.title(title)
                plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
                plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
                plt.style.use('dark_background')
                ax = plt.gca()
                ax.set_facecolor(background_color)
                ax.spines['top'].set_color(layer_color)
                ax.spines['bottom'].set_color(layer_color)
                ax.spines['right'].set_color(layer_color)
                ax.spines['left'].set_color(layer_color)
                ax.tick_params(axis='x', colors=layer_color)
                ax.tick_params(axis='y', colors=layer_color)
                ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
                plt.savefig('Layers.png')
                # add slide
                slide = prs.slides.add_slide(slide_layout)
                img_path = 'Layers.png'
                left = Inches(Left_centering)  # Adjust position as needed in icnhes
                top = Inches(Top_centering)  # Adjust position as needed in inches
                slide.shapes.add_picture(img_path, left, top)
                if close_image==1:
                 plt.close('all')

                # prs.save('Layers.pptx')


    # Save the PowerPoint presentation
    Name0=name[:-4]
    Name2 = Name0 + '.pptx'
    out_Name = f"{Destination_path}\{Name2}"
    prs.save(out_Name)


    print('total layers = ',layer_qty)
def read_and_plot_layerv23(file_path,name,image_size,Left_centering,Top_centering,Axis_Limit,scale_mode,Reduce_factor,background_color,layer_color,close_image,pattern_mode,raspberry,Destination_path,Axis_Limit_Y):
    ### this is V21 with error detection
    import ezdxf
    import pandas as pd
    import matplotlib.pyplot as plt
    from pptx import Presentation
    from pptx.util import Inches
    from dxf_surface_map import euclidean_distance,find_closest_tupleV2


    ##............READING THE DXF
    doc = ezdxf.readfile(file_path)

    # POWERPOINT SLIDE
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank slide layout

    # Extract entities (splines)
    msp = doc.modelspace()
    # for entity in msp:
    #     # Get the type of the entity
    #     print(entity.dxftype())
    splines = msp.query('SPLINE')
    polylines= msp.query('POLYLINE')
    lw_polylines= msp.query('LWPOLYLINE')
    i=0
    x=[]
    y=[]
    x_tag=[]
    y_tag=[]
    x_datum=[]
    y_datum=[]
    ####.................... Iterate over splines/polylines/lw-polyinies
    for spline in splines:
            # Check if the spline is a BSpline
            if spline.dxftype() == 'SPLINE':  # here we look for entity type SPLINE
                if spline.dxf.color==1 or spline.dxf.color==3:   # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline._control_points

                    # Print control points
                    for point in control_points:
                        i=i+1
                        #print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline


    for spline in polylines:
        # Check if the spline is a BSpline
            if spline.dxftype() == 'POLYLINE':  # here we look for entity type SPLINE
                if spline.dxf.color==1 or spline.dxf.color==3:   # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline.points()
                    # Print control points
                    for point in control_points:
                        i=i+1
                        #print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline

    qty_spline=0
    Layer_Names=[]
    Layer_points=[]
    for spline in lw_polylines:
        # Check if the spline is a BSpline
            if spline.dxftype() == 'LWPOLYLINE':  # here we look for entity type SPLINE
                 #print(spline.dxf.color)
                 layer_name = spline.dxf.layer  # Get the layer name
                 layer = doc.layers.get(layer_name)  # Get the layer object
                 color = layer.dxf.color  # Get the color from the layer
                 print (layer_name," color ",color)
                 ######## GET THE GREEN A RED LAYERS
                 if color==1 or color==3:    # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline.get_points('xy')

                    # Print control points
                    for point in control_points:
                        i=i+1
                        #print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline
                        Layer_points.append(layer_name)
                    #####................................................
                    if qty_spline == 0:
                        Layer_Names.append(layer_name)
                    elif Layer_Names[-1] != layer_name:
                        Layer_Names.append(layer_name)
                    qty_spline=qty_spline+1

                 ################### GET THE TAGGING LINES
                 if color == 7 or color == 255:
                     control_points = spline.get_points('xy')
                     for point in control_points:
                            i=i+1
                            #print("Vertice:", point,'vertice number', i)
                            x_tag.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                            y_tag.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline

                 ################### GET THE DATUM LINES
                 if color == 141:
                     control_points = spline.get_points('xy')
                     for point in control_points:
                            i=i+1
                            #print("Vertice:", point,'vertice number', i)
                            x_datum.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                            y_datum.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline



    #### BRING TO THE ORIGIN
    error_detection=0
    if len(x)==0:   #### if true then layers are not detected, error detection condition
        error_detection=1
        x=[0,0]
        y=[0,0]
    cx1=min(x)
    cy1=min(y)
    x = [i -cx1 for i in x]  # BRING X VALUES TO THE ORIGIN
    y = [i -cy1 for i in y]  # BRING Y VALUES TO THE ORIGIN
    x_tag = [i -cx1 for i in x_tag]  # BRING X VALUES TO THE ORIGIN
    y_tag = [i -cy1 for i in y_tag]  # BRING Y VALUES TO THE ORIGIN
    x_datum = [i -cx1 for i in x_datum]  # BRING X VALUES TO THE ORIGIN
    y_datum = [i -cy1 for i in y_datum]  # BRING Y VALUES TO THE ORIGIN


    #### APPLY THE OFFSET (DESFASE)
    x_lim=7040
    y_lim=2330
    cx=x_lim-max(x)
    if cx<=0:       ####  check if layer is too large
        cx=100
    cy=y_lim-max(y)
    if cy<=0:        ####  check if layer is too large
        cy=10
    x = [i + cx for i in x]  # apply offset desfase
    y = [i + cy for i in y]  # apply offset desfase
    x_tag = [i + cx for i in x_tag]  # apply offset desfase
    y_tag = [i + cy for i in y_tag]  # apply offset desfase
    x_datum = [i + cx for i in x_datum]  # apply offset desfase
    y_datum = [i + cy for i in y_datum]  # apply offset desfase

    # .....LOAD THE TABLE CALIBRATION DATA.......................................................................................................
    excel_table_calib = r'C:\Users\Juan Pablo Lopez\PycharmProjects\ProjectJP\Table_Calib.xlsx'
    Surface_map = pd.read_excel(excel_table_calib, sheet_name='Surface', header=0)
    tuples_list_real=[]
    tuples_list_autocad=[]
    for i in range(0,len(Surface_map)):
        x_tup_autocad = Surface_map['Autocad_X'][i]
        y_tup_autocad = Surface_map['Autocad_Y'][i]
        tuples_list_autocad.append((x_tup_autocad,y_tup_autocad))
        x_tup_real = Surface_map['Real_X'][i]
        y_tup_real = Surface_map['Real_Y'][i]
        tuples_list_real.append((x_tup_real,y_tup_real))

    #####......................... APPLY THE CALIBRATION......................................
    if scale_mode==1 and error_detection==0:
        for i in range(0,len(x)):
            input_tuple = (x[i],y[i])
            # print(input_tuple)
            xy_tup,idx_low,idx_high,tup_low,tup_high = find_closest_tupleV2(tuples_list_real,tuples_list_autocad,input_tuple)
            x[i] = xy_tup[0]
            y[i] = xy_tup[1]
            # print("Original ",input_tuple," Converted ",xy_tup," tup_low=",tup_low," idx low ",idx_low," tup_high=",tup_high, " idx high ",idx_high)

        if len(x_tag)>0:
            for i in range(0,len(x_tag)):
                input_tuple = (x_tag[i], y_tag[i])
                # print(i,input_tuple)
                xy_tup, idx_low, idx_high, tup_low, tup_high = find_closest_tupleV2(tuples_list_real, tuples_list_autocad, input_tuple)
                x_tag[i] = xy_tup[0]
                y_tag[i] = xy_tup[1]
                # print("Original ", input_tuple, " Converted ", xy_tup, " tup_low=", tup_low, " idx low ", idx_low," tup_high=", tup_high, " idx high ", idx_high)

        if len(x_datum) > 0:
            for i in range(0, len(x_datum)):
                input_tuple = (x_datum[i], y_datum[i])
                xy_tup, idx_low, idx_high, tup_low, tup_high = find_closest_tupleV2(tuples_list_real, tuples_list_autocad, input_tuple)
                x_datum[i] = xy_tup[0]
                y_datum[i] = xy_tup[1]
                # print("Original ", input_tuple, " Converted ", xy_tup, " tup_low=", tup_low, " idx low ", idx_low," tup_high=", tup_high, " idx high ", idx_high)


    ############# BEGGIN THE PLOTTING...............................

    Axis_Limit_x = Axis_Limit
    Axis_Limit_y = Axis_Limit_Y

    if pattern_mode==1:
        x=x_tag
        y=y_tag
        x_tag=[]
        y_tag=[]

    layer_qty=0
    x_layer=[]
    y_layer=[]
    for Layer_i in Layer_Names:    # loop through all the vertices of the splines found
        for i in range(0, len(x)):
            if Layer_points[i]==Layer_i:
                x_layer.append(x[i])
                y_layer.append(y[i])
        Layer_Name_plot = Layer_i
        layer_qty = layer_qty + 1
        # Plotting the polygon
        plt.figure(figsize=(image_size, image_size))
        plt.scatter(x_layer, y_layer, marker='o', s=0.05, color='blue')  #  s=0.1  'bo-' means blue circles connected by lines
        plt.fill(x_layer, y_layer, facecolor=layer_color, alpha=1)  # Fill the polygon
        title = 'Layer '+str(layer_qty)
        plt.title(title)
        plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
        plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
        plt.style.use('dark_background')
        ax = plt.gca()
        ax.set_facecolor(background_color)
        ax.spines['top'].set_color(layer_color)
        ax.spines['bottom'].set_color(layer_color)
        ax.spines['right'].set_color(layer_color)
        ax.spines['left'].set_color(layer_color)
        ax.tick_params(axis='x', colors=layer_color)
        ax.tick_params(axis='y', colors=layer_color)
        ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
        plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
        plt.savefig('Layers.png')
        # add slide
        slide = prs.slides.add_slide(slide_layout)
        img_path = 'Layers.png'
        left = Inches(Left_centering)  # Adjust position as needed in icnhes
        top = Inches(Top_centering)  # Adjust position as needed in inches
        slide.shapes.add_picture(img_path, left, top)
        if close_image == 1:
            plt.close('all')
        #print(title)
        x_layer = []   ### reset x_layer
        y_layer = []   ### reset y_layer




    if len(x_tag)>0:
                x_layer=x_tag
                y_layer=y_tag
                Layer_Name_plot='TAGGING'

                # Plotting the polygon
                plt.figure(figsize=(image_size, image_size))
                plt.scatter(x_layer, y_layer,marker='x', s=2, color='blue') # 'bo-' means blue circles connected by lines
                title='Tagging'
                plt.title(title)
                plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
                plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
                plt.style.use('dark_background')
                ax = plt.gca()
                ax.set_facecolor(background_color)
                ax.spines['top'].set_color(layer_color)
                ax.spines['bottom'].set_color(layer_color)
                ax.spines['right'].set_color(layer_color)
                ax.spines['left'].set_color(layer_color)
                ax.tick_params(axis='x', colors=layer_color)
                ax.tick_params(axis='y', colors=layer_color)
                ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
                plt.savefig('Layers.png')
                # add slide
                slide = prs.slides.add_slide(slide_layout)
                img_path = 'Layers.png'
                left = Inches(Left_centering)  # Adjust position as needed in icnhes
                top = Inches(Top_centering)  # Adjust position as needed in inches
                slide.shapes.add_picture(img_path, left, top)
                if close_image==1:
                 plt.close('all')

                # prs.save('Layers.pptx')



    if len(x_datum)>0:
                x_layer=x_datum
                y_layer=y_datum
                Layer_Name_plot='DATUM LINE'

                # Plotting the polygon
                plt.figure(figsize=(image_size, image_size))
                plt.scatter(x_layer, y_layer,marker="^", s=3, color='blue') # 'bo-' means blue circles connected by lines
                # plt.fill(x_layer, y_layer, facecolor=layer_color, alpha=1)  # Fill the polygon
                title='Tagging'
                plt.title(title)
                plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
                plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
                plt.style.use('dark_background')
                ax = plt.gca()
                ax.set_facecolor(background_color)
                ax.spines['top'].set_color(layer_color)
                ax.spines['bottom'].set_color(layer_color)
                ax.spines['right'].set_color(layer_color)
                ax.spines['left'].set_color(layer_color)
                ax.tick_params(axis='x', colors=layer_color)
                ax.tick_params(axis='y', colors=layer_color)
                ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
                plt.savefig('Layers.png')
                # add slide
                slide = prs.slides.add_slide(slide_layout)
                img_path = 'Layers.png'
                left = Inches(Left_centering)  # Adjust position as needed in icnhes
                top = Inches(Top_centering)  # Adjust position as needed in inches
                slide.shapes.add_picture(img_path, left, top)
                if close_image==1:
                 plt.close('all')

                # prs.save('Layers.pptx')


    # Save the PowerPoint presentation
    Name0=name[:-4]
    Name2 = Name0 + '.pptx'
    if error_detection==1:
        Name2='Error_'+Name2
    out_Name = f"{Destination_path}\{Name2}"
    prs.save(out_Name)


    print('total layers = ',layer_qty)
def read_and_plot_layerv24(file_path,name,image_size,Left_centering,Top_centering,Axis_Limit,scale_mode,Reduce_factor,background_color,layer_color,close_image,pattern_mode,raspberry,Destination_path,Axis_Limit_Y):
    ### this is v23 with full datum line
    import ezdxf
    import pandas as pd
    import matplotlib.pyplot as plt
    from pptx import Presentation
    from pptx.util import Inches
    from dxf_surface_map import euclidean_distance,find_closest_tupleV2


    ##............READING THE DXF
    doc = ezdxf.readfile(file_path)

    # POWERPOINT SLIDE
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank slide layout

    # Extract entities (splines)
    msp = doc.modelspace()
    # for entity in msp:
    #     # Get the type of the entity
    #     print(entity.dxftype())
    splines = msp.query('SPLINE')
    polylines= msp.query('POLYLINE')
    lw_polylines= msp.query('LWPOLYLINE')
    i=0
    x=[]
    y=[]
    x_tag=[]
    y_tag=[]
    x_datum=[]
    y_datum=[]
    ####.................... Iterate over splines/polylines/lw-polyinies
    for spline in splines:
            # Check if the spline is a BSpline
            if spline.dxftype() == 'SPLINE':  # here we look for entity type SPLINE
                if spline.dxf.color==1 or spline.dxf.color==3:   # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline._control_points

                    # Print control points
                    for point in control_points:
                        i=i+1
                        #print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline


    for spline in polylines:
        # Check if the spline is a BSpline
            if spline.dxftype() == 'POLYLINE':  # here we look for entity type SPLINE
                if spline.dxf.color==1 or spline.dxf.color==3:   # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline.points()
                    # Print control points
                    for point in control_points:
                        i=i+1
                        #print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline

    qty_spline=0
    Layer_Names=[]
    Layer_points=[]
    for spline in lw_polylines:
        # Check if the spline is a BSpline
            if spline.dxftype() == 'LWPOLYLINE':  # here we look for entity type SPLINE
                 #print(spline.dxf.color)
                 layer_name = spline.dxf.layer  # Get the layer name
                 layer = doc.layers.get(layer_name)  # Get the layer object
                 color = layer.dxf.color  # Get the color from the layer
                 print (layer_name," color ",color)
                 ######## GET THE GREEN A RED LAYERS
                 if color==1 or color==3:    # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline.get_points('xy')

                    # Print control points
                    for point in control_points:
                        i=i+1
                        #print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline
                        Layer_points.append(layer_name)
                    #####................................................
                    if qty_spline == 0:
                        Layer_Names.append(layer_name)
                    elif Layer_Names[-1] != layer_name:
                        Layer_Names.append(layer_name)
                    qty_spline=qty_spline+1

                 ################### GET THE TAGGING LINES
                 if color == 7 or color == 255:
                     control_points = spline.get_points('xy')
                     for point in control_points:
                            i=i+1
                            #print("Vertice:", point,'vertice number', i)
                            x_tag.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                            y_tag.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline

                 ################### GET THE DATUM LINES
                 if color == 141:
                     control_points = spline.get_points('xy')
                     for point in control_points:
                            i=i+1
                            #print("Vertice:", point,'vertice number', i)
                            x_datum.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                            y_datum.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline



    #### BRING TO THE ORIGIN
    error_detection=0
    if len(x)==0:   #### if true then layers are not detected, error detection condition
        error_detection=1
        x=[0,0]
        y=[0,0]
    cx1=min(x)
    cy1=min(y)
    x = [i -cx1 for i in x]  # BRING X VALUES TO THE ORIGIN
    y = [i -cy1 for i in y]  # BRING Y VALUES TO THE ORIGIN
    x_tag = [i -cx1 for i in x_tag]  # BRING X VALUES TO THE ORIGIN
    y_tag = [i -cy1 for i in y_tag]  # BRING Y VALUES TO THE ORIGIN
    x_datum = [i -cx1 for i in x_datum]  # BRING X VALUES TO THE ORIGIN
    y_datum = [i -cy1 for i in y_datum]  # BRING Y VALUES TO THE ORIGIN


    #### APPLY THE OFFSET (DESFASE)
    x_lim=7040
    y_lim=2330
    cx=x_lim-max(x)
    if cx<=0:       ####  check if layer is too large
        cx=100
    cy=y_lim-max(y)
    if cy<=0:        ####  check if layer is too large
        cy=10
    x = [i + cx for i in x]  # apply offset desfase
    y = [i + cy for i in y]  # apply offset desfase
    x_tag = [i + cx for i in x_tag]  # apply offset desfase
    y_tag = [i + cy for i in y_tag]  # apply offset desfase
    x_datum = [i + cx for i in x_datum]  # apply offset desfase
    y_datum = [i + cy for i in y_datum]  # apply offset desfase

    # .....LOAD THE TABLE CALIBRATION DATA.......................................................................................................
    excel_table_calib = r'C:\Users\Juan Pablo Lopez\PycharmProjects\ProjectJP\Table_Calib.xlsx'
    Surface_map = pd.read_excel(excel_table_calib, sheet_name='Surface', header=0)
    tuples_list_real=[]
    tuples_list_autocad=[]
    for i in range(0,len(Surface_map)):
        x_tup_autocad = Surface_map['Autocad_X'][i]
        y_tup_autocad = Surface_map['Autocad_Y'][i]
        tuples_list_autocad.append((x_tup_autocad,y_tup_autocad))
        x_tup_real = Surface_map['Real_X'][i]
        y_tup_real = Surface_map['Real_Y'][i]
        tuples_list_real.append((x_tup_real,y_tup_real))

    #####......................... APPLY THE CALIBRATION......................................
    if scale_mode==1 and error_detection==0:
        for i in range(0,len(x)):
            input_tuple = (x[i],y[i])
            # print(input_tuple)
            xy_tup,idx_low,idx_high,tup_low,tup_high = find_closest_tupleV2(tuples_list_real,tuples_list_autocad,input_tuple)
            x[i] = xy_tup[0]
            y[i] = xy_tup[1]
            # print("Original ",input_tuple," Converted ",xy_tup," tup_low=",tup_low," idx low ",idx_low," tup_high=",tup_high, " idx high ",idx_high)

        if len(x_tag)>0:
            for i in range(0,len(x_tag)):
                input_tuple = (x_tag[i], y_tag[i])
                # print(i,input_tuple)
                xy_tup, idx_low, idx_high, tup_low, tup_high = find_closest_tupleV2(tuples_list_real, tuples_list_autocad, input_tuple)
                x_tag[i] = xy_tup[0]
                y_tag[i] = xy_tup[1]
                # print("Original ", input_tuple, " Converted ", xy_tup, " tup_low=", tup_low, " idx low ", idx_low," tup_high=", tup_high, " idx high ", idx_high)

        if len(x_datum) > 0:
            for i in range(0, len(x_datum)):
                input_tuple = (x_datum[i], y_datum[i])
                xy_tup, idx_low, idx_high, tup_low, tup_high = find_closest_tupleV2(tuples_list_real, tuples_list_autocad, input_tuple)
                x_datum[i] = xy_tup[0]
                y_datum[i] = xy_tup[1]
                # print("Original ", input_tuple, " Converted ", xy_tup, " tup_low=", tup_low, " idx low ", idx_low," tup_high=", tup_high, " idx high ", idx_high)


    ############# BEGGIN THE PLOTTING...............................

    Axis_Limit_x = Axis_Limit
    Axis_Limit_y = Axis_Limit_Y

    if pattern_mode==1:
        x=x_tag
        y=y_tag
        x_tag=[]
        y_tag=[]

    layer_qty=0
    x_layer=[]
    y_layer=[]
    for Layer_i in Layer_Names:    # loop through all the vertices of the splines found
        for i in range(0, len(x)):
            if Layer_points[i]==Layer_i:
                x_layer.append(x[i])
                y_layer.append(y[i])
        Layer_Name_plot = Layer_i
        layer_qty = layer_qty + 1
        # Plotting the polygon
        plt.figure(figsize=(image_size, image_size))
        plt.scatter(x_layer, y_layer, marker='o', s=0.05, color='blue')  #  s=0.1  'bo-' means blue circles connected by lines
        plt.fill(x_layer, y_layer, facecolor=layer_color, alpha=1)  # Fill the polygon
        title = 'Layer '+str(layer_qty)
        plt.title(title)
        plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
        plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
        plt.style.use('dark_background')
        ax = plt.gca()
        ax.set_facecolor(background_color)
        ax.spines['top'].set_color(layer_color)
        ax.spines['bottom'].set_color(layer_color)
        ax.spines['right'].set_color(layer_color)
        ax.spines['left'].set_color(layer_color)
        ax.tick_params(axis='x', colors=layer_color)
        ax.tick_params(axis='y', colors=layer_color)
        ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
        plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
        plt.savefig('Layers.png')
        # add slide
        slide = prs.slides.add_slide(slide_layout)
        img_path = 'Layers.png'
        left = Inches(Left_centering)  # Adjust position as needed in icnhes
        top = Inches(Top_centering)  # Adjust position as needed in inches
        slide.shapes.add_picture(img_path, left, top)
        if close_image == 1:
            plt.close('all')
        #print(title)
        x_layer = []   ### reset x_layer
        y_layer = []   ### reset y_layer




    if len(x_tag)>0:
                x_layer=x_tag
                y_layer=y_tag
                Layer_Name_plot='TAGGING'

                # Plotting the polygon
                plt.figure(figsize=(image_size, image_size))
                plt.scatter(x_layer, y_layer,marker='x', s=3, color='blue') # 'bo-' means blue circles connected by lines
                title='Tagging'
                plt.title(title)
                plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
                plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
                plt.style.use('dark_background')
                ax = plt.gca()
                ax.set_facecolor(background_color)
                ax.spines['top'].set_color(layer_color)
                ax.spines['bottom'].set_color(layer_color)
                ax.spines['right'].set_color(layer_color)
                ax.spines['left'].set_color(layer_color)
                ax.tick_params(axis='x', colors=layer_color)
                ax.tick_params(axis='y', colors=layer_color)
                ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
                plt.savefig('Layers.png')
                # add slide
                slide = prs.slides.add_slide(slide_layout)
                img_path = 'Layers.png'
                left = Inches(Left_centering)  # Adjust position as needed in icnhes
                top = Inches(Top_centering)  # Adjust position as needed in inches
                slide.shapes.add_picture(img_path, left, top)
                if close_image==1:
                 plt.close('all')

                # prs.save('Layers.pptx')



    if len(x_datum)>0:
                x_datum1=[]
                y_datum1=[]
                i=0
                while True:
                    if i < len(x_datum)-1:
                        diff=abs(x_datum[i+1]-x_datum[i])
                        print(diff)
                        if diff<=100:  # 100
                            x_datum1.append(x_datum[i])
                            y_datum1.append(y_datum[i])
                            i=i+1
                        else:
                            break
                    else:
                        break

                x_datum2=x_datum[i+1:]
                y_datum2=y_datum[i+1:]

                if len(x_datum2)==1:
                    print('WARNING DATUM LINE 2', name)

                ##### PLOT THE FIRST DATUM LINE
                Layer_Name_plot='DATUM LINES'
                x_layer=x_datum1
                y_layer=y_datum1
                # Plotting the polygon
                plt.figure(figsize=(image_size, image_size))
                plt.plot(x_layer, y_layer,color=layer_color, linestyle='-', linewidth=1, marker='o', markersize=0.1) # 'bo-' means blue circles connected by lines
                plt.plot(x_datum2, y_datum2, color=layer_color, linestyle='-', linewidth=1, marker='o',markersize=0.1)  # 'bo-' means blue circles connected by lines
                title='Tagging'
                plt.title(title)
                plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
                plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
                plt.style.use('dark_background')
                ax = plt.gca()
                ax.set_facecolor(background_color)
                ax.spines['top'].set_color(layer_color)
                ax.spines['bottom'].set_color(layer_color)
                ax.spines['right'].set_color(layer_color)
                ax.spines['left'].set_color(layer_color)
                ax.tick_params(axis='x', colors=layer_color)
                ax.tick_params(axis='y', colors=layer_color)
                ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
                plt.savefig('Layers.png')
                # add slide
                slide = prs.slides.add_slide(slide_layout)
                img_path = 'Layers.png'
                left = Inches(Left_centering)  # Adjust position as needed in icnhes
                top = Inches(Top_centering)  # Adjust position as needed in inches
                slide.shapes.add_picture(img_path, left, top)
                if close_image==1:
                 plt.close('all')
                # prs.save('Layers.pptx')







    # Save the PowerPoint presentation
    Name0=name[:-4]
    Name2 = Name0 + '.pptx'
    if error_detection==1:
        Name2='Error_'+Name2
    out_Name = f"{Destination_path}\{Name2}"
    prs.save(out_Name)


    print('total layers = ',layer_qty)



Origin_path=r'C:\Users\Juan Pablo Lopez\OneDrive - Rewair A S\Desktop\PFILES\Python_versions\LG projecting files\origin'
Destination_path=r'C:\Users\Juan Pablo Lopez\OneDrive - Rewair A S\Desktop\PFILES\Python_versions\LG projecting files\destiny'
image_size = 12.7  # in inches    12.7
Left_centering = -1.55  # in inches   -1.55
Top_centering = -5.0  # in inches   -4.9
Axis_Limit = 8360  # in MM    last official-8355,   (20mm(-)-8335)   (5mm(+)-8360)
Axis_Limit_Y=8355  # ONLY FOR V19,  last official-8355
Reduce_factor = 1  # default = 1, if not it is used for scaling down the original image by a factor, eg 10,100,1000
background_color = 'white'
layer_color = 'blue'
close_image = 1  # 1- close all images, 0- open all images
scale_mode = 1  # 1-apply scale, 2- no scale
pattern_mode = 0  # 1- loading a pattern dxf file, 0-  normal dxf
raspberry = 1  # 1- using the raspberry pi, 0- using the laptop


import os
DXF_Names = os.listdir(Origin_path)  # this extracts the names of the files inside source folder
for k1 in DXF_Names:  # this is for eliminating the files that are not photos
    if k1.endswith("dxf") == False:  # dxf
        Path2Remove = f"{Origin_path}\{k1}"
        os.remove(Path2Remove)
DXF_Names = os.listdir(Origin_path)  # this is to update the names after deleting non dxf files
print(" TOTAL OF ", len(DXF_Names), " DXF UPLOADED")

for name in DXF_Names:
    dxf_file = f"{Origin_path}\{name}"
    #print(name)
    #print(dxf_file)
    #read_and_plot_layerv13(dxf_file,name, image_size, Left_centering, Top_centering, Axis_Limit, scale_mode, Reduce_factor,background_color, layer_color, close_image, pattern_mode, raspberry,Destination_path)
    read_and_plot_layerv21(dxf_file, name, image_size, Left_centering, Top_centering, Axis_Limit, scale_mode,
                           Reduce_factor, background_color, layer_color, close_image, pattern_mode, raspberry,
                           Destination_path,Axis_Limit_Y)
    print(name+' fully converted')
print('All files converted')


#### V17- OFFSET DESFASE version not correct yet
#### V18- full datum line, not working yet
#### V13- LAST OFFICIAL, works fine
#### V19- This is V13 with separate x-axis and y-axis limit in the plot
#### V20 LAST OFFCIAL, this is V19 with Surface Map calibration
#### V21 this is V20 with re-positioning
#### V22 this is
