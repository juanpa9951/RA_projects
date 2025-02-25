#####  THIS FOR ORIGINAL CALIBRATION METHOD

def master_pattern_x(image_size,Left_centering,Top_centering,Axis_Limit,background_color,layer_color,close_image,Axis_Limit_Y):

    import matplotlib.pyplot as plt
    from pptx import Presentation
    from pptx.util import Inches
    import numpy as np


    # POWERPOINT SLIDE
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank slide layout

    ############# BEGGIN THE PLOTTING

    Axis_Limit_x = Axis_Limit
    Axis_Limit_y = Axis_Limit_Y
    x_base=400    #500
    y_base=300    #100

    for k in range(1,22):   ## divisiones eje Y     25
            for i in range(1,26):  ## divisiones eje X     32
                        x_lenght=x_base*i                  ####modo x
                        y_lenght=y_base*k                  ####modo x

                        x_layer=[0,x_lenght]               ####modo x
                        y_layer=[y_lenght,y_lenght]        ####modo x


                        Layer_Name_plot='X='+str(x_lenght)+" Y="+str(y_lenght)

                        # Plotting the polygon
                        plt.figure(figsize=(image_size, image_size))
                        plt.plot(x_layer, y_layer,color='blue', linestyle='-', linewidth=0.3, marker='x',markersize=0.2) # 'bo-' means blue circles connected by lines
                        title='pattern'
                        plt.title(title)


                        plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
                        plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
                        plt.style.use('dark_background')
                        ax = plt.gca()
                        ax.set_facecolor(background_color)
                        ax.spines['top'].set_color(layer_color)
                        ax.spines['bottom'].set_color(layer_color)
                        ax.spines['right'].set_color(layer_color)
                        ax.spines['left'].set_color(layer_color)
                        ax.tick_params(axis='x', colors=layer_color)
                        ax.tick_params(axis='y', colors=layer_color)
                        ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                        plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
                        plt.savefig('Layers.png')
                        # add slide
                        slide = prs.slides.add_slide(slide_layout)
                        img_path = 'Layers.png'
                        left = Inches(Left_centering)  # Adjust position as needed in icnhes
                        top = Inches(Top_centering)  # Adjust position as needed in inches
                        slide.shapes.add_picture(img_path, left, top)
                        if close_image==1:
                         plt.close('all')

                        # prs.save('Layers.pptx')

    prs.save('Master_Pattern_PS_x.pptx')
def master_pattern_y(image_size,Left_centering,Top_centering,Axis_Limit,background_color,layer_color,close_image,Axis_Limit_Y):

    import matplotlib.pyplot as plt
    from pptx import Presentation
    from pptx.util import Inches
    import numpy as np

    # POWERPOINT SLIDE
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank slide layout

    ############# BEGGIN THE PLOTTING

    Axis_Limit_x = Axis_Limit
    Axis_Limit_y = Axis_Limit_Y
    x_base=400    #500
    y_base=300    #100

    for k in range(1,26):   ## divisiones eje x     25
            for i in range(1,22):  ## divisiones eje y    32

                        x_lenght=x_base*k                    ####modo y
                        y_lenght=y_base*i                    ####modo y

                        x_layer=[x_lenght,x_lenght]          ####modo y
                        y_layer=[0,y_lenght]                 ####modo y

                        Layer_Name_plot='X='+str(x_lenght)+" Y="+str(y_lenght)

                        # Plotting the polygon
                        plt.figure(figsize=(image_size, image_size))
                        plt.plot(x_layer, y_layer,color='blue', linestyle='-', linewidth=0.3, marker='x',markersize=0.2) # 'bo-' means blue circles connected by lines
                        title='pattern'
                        plt.title(title)


                        plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
                        plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
                        plt.style.use('dark_background')
                        ax = plt.gca()
                        ax.set_facecolor(background_color)
                        ax.spines['top'].set_color(layer_color)
                        ax.spines['bottom'].set_color(layer_color)
                        ax.spines['right'].set_color(layer_color)
                        ax.spines['left'].set_color(layer_color)
                        ax.tick_params(axis='x', colors=layer_color)
                        ax.tick_params(axis='y', colors=layer_color)
                        ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                        plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
                        plt.savefig('Layers.png')
                        # add slide
                        slide = prs.slides.add_slide(slide_layout)
                        img_path = 'Layers.png'
                        left = Inches(Left_centering)  # Adjust position as needed in icnhes
                        top = Inches(Top_centering)  # Adjust position as needed in inches
                        slide.shapes.add_picture(img_path, left, top)
                        if close_image==1:
                         plt.close('all')

                        # prs.save('Layers.pptx')

    prs.save('Master_Pattern_PS_y.pptx')


image_size = 12.9  # in inches    12.9
Left_centering = -1.60  # in inches   -1.60
Top_centering = -4.0  # in inches   -4.0
Axis_Limit = 10000  # in MM    # 10000
Axis_Limit_Y=10000  # ONLY FOR V19,  LAST OFFICIAL 10000

background_color = 'white'
layer_color = 'blue'
close_image = 1  # 1- close all images, 2- open all images

####  CHOOSE EITHER X OR Y
master_pattern_x(image_size,Left_centering,Top_centering,Axis_Limit,background_color,layer_color,close_image,Axis_Limit_Y)
print("ready x files")
master_pattern_y(image_size,Left_centering,Top_centering,Axis_Limit,background_color,layer_color,close_image,Axis_Limit_Y)
print("ready y files")

