def read_and_plot_layerv8(file_path,image_size,Left_centering,Top_centering,Axis_Limit,scale_mode,Reduce_factor,background_color,layer_color,close_image):
    import ezdxf
    import pandas as pd
    import matplotlib.pyplot as plt
    from pptx import Presentation
    from pptx.util import Inches


    ##............READING THE DXF
    doc = ezdxf.readfile(file_path)

    # POWERPOINT SLIDE
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank slide layout
    # image_size=12 # in inches
    # Left_centering=-1   #in inches
    # Top_centering=-4     #


    # Extract entities (splines)
    msp = doc.modelspace()
    for entity in msp:
        # Get the type of the entity
        print(entity.dxftype())
    splines = msp.query('SPLINE')
    polylines= msp.query('POLYLINE')
    lw_polylines= msp.query('LWPOLYLINE')
    i=0
    x=[]
    y=[]
    x_tag=[]
    y_tag=[]
    ####.................... Iterate over splines/polylines/lw-polyinies
    for spline in splines:
            # Check if the spline is a BSpline
            if spline.dxftype() == 'SPLINE':  # here we look for entity type SPLINE
                if spline.dxf.color==1 or spline.dxf.color==3:   # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline._control_points

                    # Print control points
                    for point in control_points:
                        i=i+1
                        print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline


    for spline in polylines:
        # Check if the spline is a BSpline
            if spline.dxftype() == 'POLYLINE':  # here we look for entity type SPLINE
                if spline.dxf.color==1 or spline.dxf.color==3:   # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline.points()
                    # Print control points
                    for point in control_points:
                        i=i+1
                        print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline

    qty_spline=0
    Layer_Names=[]
    for spline in lw_polylines:
        # Check if the spline is a BSpline
            if spline.dxftype() == 'LWPOLYLINE':  # here we look for entity type SPLINE
                 print(spline.dxf.color)
                 layer_name = spline.dxf.layer  # Get the layer name
                 layer = doc.layers.get(layer_name)  # Get the layer object
                 color = layer.dxf.color  # Get the color from the layer
                 print (color)
                 ######## GET THE GREEN A RED LAYERS
                 if color==1 or color==3:    # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline.get_points('xy')

                    # Print control points
                    for point in control_points:
                        i=i+1
                        print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline
                 #####................................................
                    if qty_spline == 0:
                        Layer_Names.append(layer_name)
                    elif Layer_Names[-1] != layer_name:
                        Layer_Names.append(layer_name)
                    qty_spline=qty_spline+1

                 ################### GET THE TAGGING LINES
                 if color == 7 or color == 255:
                     control_points = spline.get_points('xy')
                     for point in control_points:
                            i=i+1
                            print("Vertice:", point,'vertice number', i)
                            x_tag.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                            y_tag.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline


            #####...................1st scale reduction  OPTIONAL...........................
    # Reduce_factor=1
    x = [i / Reduce_factor for i in x]
    y = [i / Reduce_factor for i in y]

    excel_table_calib = r'C:\Users\Juan Pablo Lopez\PycharmProjects\ProjectJP\Table_Calib.xlsx'
    # .....LOAD THE TABLE CALIBRATION DATA.......................................................................................................
    X_table = pd.read_excel(excel_table_calib, sheet_name='X_axis', header=0)
    Y_table = pd.read_excel(excel_table_calib, sheet_name='Y_axis', header=0)

    if scale_mode==1:
        for i in range(0,len(x)):
            # X_value_index = (X_table['Real_Measure'] - x[i]).abs().idxmin()
            # X_value_autocad= X_table.at[X_value_index, 'Autocad_Measure']
            # x[i]=X_value_autocad
            target_value = x[i]
            df=X_table
            # Find indices of bracketing points
            idx = df['Real_Measure'].searchsorted(target_value)

            # Check if target is within data range
            if idx == 0 or idx == len(x):
                print("Target value outside data range for interpolation")
            else:
                # Extract bracketed values
                a_prev = df.loc[idx - 1, 'Real_Measure']
                a_next = df.loc[idx, 'Real_Measure']
                b_prev = df.loc[idx - 1, 'Autocad_Measure']
                b_next = df.loc[idx, 'Autocad_Measure']

                # Perform linear interpolation
                interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                # print(f"Interpolated value for {target_value}: {interpolated_value}")
                x[i]=interpolated_value



        for i in range(0,len(y)):
            # X_value_index = (X_table['Real_Measure'] - x[i]).abs().idxmin()
            # X_value_autocad= X_table.at[X_value_index, 'Autocad_Measure']
            # x[i]=X_value_autocad
            target_value = y[i]
            df=Y_table
            # Find indices of bracketing points
            idx = df['Real_Measure'].searchsorted(target_value)

            # Check if target is within data range
            if idx == 0 or idx == len(y):
                print("Target value outside data range for interpolation")
            else:
                # Extract bracketed values
                a_prev = df.loc[idx - 1, 'Real_Measure']
                a_next = df.loc[idx, 'Real_Measure']
                b_prev = df.loc[idx - 1, 'Autocad_Measure']
                b_next = df.loc[idx, 'Autocad_Measure']

                # Perform linear interpolation
                interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                # print(f"Interpolated value for {target_value}: {interpolated_value}")
                y[i]=interpolated_value

    ############# BEGGIN THE PLOTTING

    Axis_Limit_x = Axis_Limit
    Axis_Limit_y = Axis_Limit
    x_init=x[0]
    y_init=y[0]
    z=0
    layer_qty=0
    check=0
    for i in range(0,len(x)):    # loop through all the vertices of the splines found
        #print(i)
        if x[i]==x_init and y[i]==y_init:   # layer detection condition is when there is a closed loop, initial and final points are equal
            check=check+1 # the initial point of the loop
            if check==2: # the final point of the loop
                Layer_Name_plot=Layer_Names[layer_qty]
                x_layer = x[z:i + 1]
                y_layer = y[z:i + 1]
                z=i+1
                layer_qty=layer_qty+1
                check=0
                if i<(len(x)-1):
                 x_init = x[i+1]
                 y_init = y[i+1]
                print('Layer ', layer_qty)
                print('X_coordenates ', x_layer)
                print('y_coordenates ', y_layer)

                # Plotting the polygon
                plt.figure(figsize=(image_size, image_size))
                plt.plot(x_layer, y_layer, color=layer_color, linestyle='-', linewidth=0.1, marker='o', markersize=0.1) # 'bo-' means blue circles connected by lines
                plt.fill(x_layer, y_layer,facecolor = layer_color ,alpha=1)  # Fill the polygon
                title='Layer '+str(layer_qty)
                plt.title(title)
                #plt.grid(True, linestyle='--', linewidth=0.5)  # Set linewidth to 1.5
                plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
                plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
                plt.style.use('dark_background')
                ax = plt.gca()
                ax.set_facecolor(background_color)
                ax.spines['top'].set_color(layer_color)
                ax.spines['bottom'].set_color(layer_color)
                ax.spines['right'].set_color(layer_color)
                ax.spines['left'].set_color(layer_color)
                ax.tick_params(axis='x', colors=layer_color)
                ax.tick_params(axis='y', colors=layer_color)
                ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                plt.text(2000, 1700, Layer_Name_plot, fontsize=15, color='red')
                plt.savefig('Layers.png')
                # add slide
                slide = prs.slides.add_slide(slide_layout)
                img_path = 'Layers.png'
                left = Inches(Left_centering)  # Adjust position as needed in icnhes
                top = Inches(Top_centering)  # Adjust position as needed in inches
                slide.shapes.add_picture(img_path, left, top)
                if close_image==1:
                 plt.close('all')

    if len(x_tag)>0:
        x=x_tag
        print(len(x))
        y=y_tag
        seen = set()
        x = [x for x in x if not (x in seen or seen.add(x))]    ### remove the duplicates
        print(len(x))
        seen = set()
        y = [y for y in y if not (y in seen or seen.add(y))]    ### remove the duplicates


        for i in range(0,len(x)):    # loop through all the vertices of the splines found

                    Layer_Name_plot='TAGGING'
                    if i < (len(x) - 1):
                        x_layer = x[0:2]
                        y_layer = y[0:2]
                        del x[:2]

                    # Plotting the polygon
                    plt.figure(figsize=(image_size, image_size))
                    plt.plot(x_layer, y_layer, color=layer_color, linestyle='-', linewidth=2, marker='o', markersize=0.1) # 'bo-' means blue circles connected by lines
                    # plt.fill(x_layer, y_layer,facecolor = layer_color ,alpha=1)  # Fill the polygon
                    title='Tagging'
                    plt.title(title)
                    plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
                    plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
                    plt.style.use('dark_background')
                    ax = plt.gca()
                    ax.set_facecolor(background_color)
                    ax.spines['top'].set_color(layer_color)
                    ax.spines['bottom'].set_color(layer_color)
                    ax.spines['right'].set_color(layer_color)
                    ax.spines['left'].set_color(layer_color)
                    ax.tick_params(axis='x', colors=layer_color)
                    ax.tick_params(axis='y', colors=layer_color)
                    ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                    plt.text(2000, 1700, Layer_Name_plot, fontsize=15, color='red')
                    plt.savefig('Layers.png')
                    # add slide
                    slide = prs.slides.add_slide(slide_layout)
                    img_path = 'Layers.png'
                    left = Inches(Left_centering)  # Adjust position as needed in icnhes
                    top = Inches(Top_centering)  # Adjust position as needed in inches
                    slide.shapes.add_picture(img_path, left, top)
                    if close_image==1:
                     plt.close('all')

                    prs.save('Layers.pptx')


    # Save the PowerPoint presentation
    prs.save('Layers.pptx')


    print('total layers = ',layer_qty)
def read_and_plot_layerv9(file_path,image_size,Left_centering,Top_centering,Axis_Limit,scale_mode,Reduce_factor,background_color,layer_color,close_image):
    import ezdxf
    import pandas as pd
    import matplotlib.pyplot as plt
    from pptx import Presentation
    from pptx.util import Inches


    ##............READING THE DXF
    doc = ezdxf.readfile(file_path)

    # POWERPOINT SLIDE
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank slide layout

    # Extract entities (splines)
    msp = doc.modelspace()
    for entity in msp:
        # Get the type of the entity
        print(entity.dxftype())
    splines = msp.query('SPLINE')
    polylines= msp.query('POLYLINE')
    lw_polylines= msp.query('LWPOLYLINE')
    i=0
    x=[]
    y=[]
    x_tag=[]
    y_tag=[]
    ####.................... Iterate over splines/polylines/lw-polyinies
    for spline in splines:
            # Check if the spline is a BSpline
            if spline.dxftype() == 'SPLINE':  # here we look for entity type SPLINE
                if spline.dxf.color==1 or spline.dxf.color==3:   # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline._control_points

                    # Print control points
                    for point in control_points:
                        i=i+1
                        print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline


    for spline in polylines:
        # Check if the spline is a BSpline
            if spline.dxftype() == 'POLYLINE':  # here we look for entity type SPLINE
                if spline.dxf.color==1 or spline.dxf.color==3:   # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline.points()
                    # Print control points
                    for point in control_points:
                        i=i+1
                        print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline

    qty_spline=0
    Layer_Names=[]
    for spline in lw_polylines:
        # Check if the spline is a BSpline
            if spline.dxftype() == 'LWPOLYLINE':  # here we look for entity type SPLINE
                 print(spline.dxf.color)
                 layer_name = spline.dxf.layer  # Get the layer name
                 layer = doc.layers.get(layer_name)  # Get the layer object
                 color = layer.dxf.color  # Get the color from the layer
                 print (color)
                 ######## GET THE GREEN A RED LAYERS
                 if color==1 or color==3:    # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline.get_points('xy')

                    # Print control points
                    for point in control_points:
                        i=i+1
                        print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline
                 #####................................................
                    if qty_spline == 0:
                        Layer_Names.append(layer_name)
                    elif Layer_Names[-1] != layer_name:
                        Layer_Names.append(layer_name)
                    qty_spline=qty_spline+1

                 ################### GET THE TAGGING LINES
                 if color == 7 or color == 255:
                     control_points = spline.get_points('xy')
                     for point in control_points:
                            i=i+1
                            print("Vertice:", point,'vertice number', i)
                            x_tag.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                            y_tag.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline


            #####...................1st scale reduction  OPTIONAL...........................
    # Reduce_factor=1
    x = [i / Reduce_factor for i in x]
    y = [i / Reduce_factor for i in y]

    excel_table_calib = r'C:\Users\Juan Pablo Lopez\PycharmProjects\ProjectJP\Table_Calib.xlsx'
    # .....LOAD THE TABLE CALIBRATION DATA.......................................................................................................
    X_table = pd.read_excel(excel_table_calib, sheet_name='X_axis', header=0)
    Y_table = pd.read_excel(excel_table_calib, sheet_name='Y_axis', header=0)

    if scale_mode==1:
        for i in range(0,len(x)):
            target_value = x[i]
            df=X_table
            # Find indices of bracketing points
            idx = df['Real_Measure'].searchsorted(target_value)

            # Check if target is within data range
            if idx == 0 or idx == len(x):
                print("Target value outside data range for interpolation")
            else:
                # Extract bracketed values
                a_prev = df.loc[idx - 1, 'Real_Measure']
                a_next = df.loc[idx, 'Real_Measure']
                b_prev = df.loc[idx - 1, 'Autocad_Measure']
                b_next = df.loc[idx, 'Autocad_Measure']

                # Perform linear interpolation
                interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                # print(f"Interpolated value for {target_value}: {interpolated_value}")
                x[i]=interpolated_value

        for i in range(0,len(y)):
            target_value = y[i]
            df=Y_table
            # Find indices of bracketing points
            idx = df['Real_Measure'].searchsorted(target_value)

            # Check if target is within data range
            if idx == 0 or idx == len(y):
                print("Target value outside data range for interpolation")
            else:
                # Extract bracketed values
                a_prev = df.loc[idx - 1, 'Real_Measure']
                a_next = df.loc[idx, 'Real_Measure']
                b_prev = df.loc[idx - 1, 'Autocad_Measure']
                b_next = df.loc[idx, 'Autocad_Measure']

                # Perform linear interpolation
                interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                # print(f"Interpolated value for {target_value}: {interpolated_value}")
                y[i]=interpolated_value

        for i in range(0,len(x_tag)):
            target_value = x_tag[i]
            df=X_table
            # Find indices of bracketing points
            idx = df['Real_Measure'].searchsorted(target_value)

            # Check if target is within data range
            if idx == 0 or idx == len(x_tag):
                print("Target value outside data range for interpolation")
            else:
                # Extract bracketed values
                a_prev = df.loc[idx - 1, 'Real_Measure']
                a_next = df.loc[idx, 'Real_Measure']
                b_prev = df.loc[idx - 1, 'Autocad_Measure']
                b_next = df.loc[idx, 'Autocad_Measure']

                # Perform linear interpolation
                interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                # print(f"Interpolated value for {target_value}: {interpolated_value}")
                x_tag[i]=interpolated_value



        for i in range(0,len(y_tag)):
            target_value = y_tag[i]
            df=Y_table
            # Find indices of bracketing points
            idx = df['Real_Measure'].searchsorted(target_value)

            # Check if target is within data range
            if idx == 0 or idx == len(y_tag):
                print("Target value outside data range for interpolation")
            else:
                # Extract bracketed values
                a_prev = df.loc[idx - 1, 'Real_Measure']
                a_next = df.loc[idx, 'Real_Measure']
                b_prev = df.loc[idx - 1, 'Autocad_Measure']
                b_next = df.loc[idx, 'Autocad_Measure']

                # Perform linear interpolation
                interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                # print(f"Interpolated value for {target_value}: {interpolated_value}")
                y_tag[i]=interpolated_value

    ############# BEGGIN THE PLOTTING

    Axis_Limit_x = Axis_Limit
    Axis_Limit_y = Axis_Limit
    x_init=x[0]
    y_init=y[0]
    z=0
    layer_qty=0
    check=0
    for i in range(0,len(x)):    # loop through all the vertices of the splines found
        #print(i)
        if x[i]==x_init and y[i]==y_init:   # layer detection condition is when there is a closed loop, initial and final points are equal
            check=check+1 # the initial point of the loop
            if check==2: # the final point of the loop
                Layer_Name_plot=Layer_Names[layer_qty]
                x_layer = x[z:i + 1]
                y_layer = y[z:i + 1]
                z=i+1
                layer_qty=layer_qty+1
                check=0
                if i<(len(x)-1):
                 x_init = x[i+1]
                 y_init = y[i+1]
                print('Layer ', layer_qty)
                print('X_coordenates ', x_layer)
                print('y_coordenates ', y_layer)

                # Plotting the polygon
                plt.figure(figsize=(image_size, image_size))
                plt.plot(x_layer, y_layer, color=layer_color, linestyle='-', linewidth=0.1, marker='o', markersize=0.1) # 'bo-' means blue circles connected by lines
                plt.fill(x_layer, y_layer,facecolor = layer_color ,alpha=1)  # Fill the polygon
                title='Layer '+str(layer_qty)
                plt.title(title)
                #plt.grid(True, linestyle='--', linewidth=0.5)  # Set linewidth to 1.5
                plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
                plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
                plt.style.use('dark_background')
                ax = plt.gca()
                ax.set_facecolor(background_color)
                ax.spines['top'].set_color(layer_color)
                ax.spines['bottom'].set_color(layer_color)
                ax.spines['right'].set_color(layer_color)
                ax.spines['left'].set_color(layer_color)
                ax.tick_params(axis='x', colors=layer_color)
                ax.tick_params(axis='y', colors=layer_color)
                ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
                plt.savefig('Layers.png')
                # add slide
                slide = prs.slides.add_slide(slide_layout)
                img_path = 'Layers.png'
                left = Inches(Left_centering)  # Adjust position as needed in icnhes
                top = Inches(Top_centering)  # Adjust position as needed in inches
                slide.shapes.add_picture(img_path, left, top)
                if close_image==1:
                 plt.close('all')

    if len(x_tag)>0:
                x_layer=x_tag
                y_layer=y_tag
                Layer_Name_plot='TAGGING'

                # Plotting the polygon
                plt.figure(figsize=(image_size, image_size))
                plt.scatter(x_layer, y_layer,marker='x', s=3, color='blue') # 'bo-' means blue circles connected by lines
                title='Tagging'
                plt.title(title)
                plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
                plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
                plt.style.use('dark_background')
                ax = plt.gca()
                ax.set_facecolor(background_color)
                ax.spines['top'].set_color(layer_color)
                ax.spines['bottom'].set_color(layer_color)
                ax.spines['right'].set_color(layer_color)
                ax.spines['left'].set_color(layer_color)
                ax.tick_params(axis='x', colors=layer_color)
                ax.tick_params(axis='y', colors=layer_color)
                ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
                plt.savefig('Layers.png')
                # add slide
                slide = prs.slides.add_slide(slide_layout)
                img_path = 'Layers.png'
                left = Inches(Left_centering)  # Adjust position as needed in icnhes
                top = Inches(Top_centering)  # Adjust position as needed in inches
                slide.shapes.add_picture(img_path, left, top)
                if close_image==1:
                 plt.close('all')

                # prs.save('Layers.pptx')


    # Save the PowerPoint presentation
    prs.save('Layers.pptx')


    print('total layers = ',layer_qty)
def read_and_plot_layerv10(file_path,image_size,Left_centering,Top_centering,Axis_Limit,scale_mode,Reduce_factor,background_color,layer_color,close_image,pattern_mode,raspberry):
    ### Added datum lines in the form of scatter plot
    import ezdxf
    import pandas as pd
    import matplotlib.pyplot as plt
    from pptx import Presentation
    from pptx.util import Inches


    ##............READING THE DXF
    doc = ezdxf.readfile(file_path)

    # POWERPOINT SLIDE
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank slide layout

    # Extract entities (splines)
    msp = doc.modelspace()
    for entity in msp:
        # Get the type of the entity
        print(entity.dxftype())
    splines = msp.query('SPLINE')
    polylines= msp.query('POLYLINE')
    lw_polylines= msp.query('LWPOLYLINE')
    i=0
    x=[]
    y=[]
    x_tag=[]
    y_tag=[]
    x_datum=[]
    y_datum=[]
    ####.................... Iterate over splines/polylines/lw-polyinies
    for spline in splines:
            # Check if the spline is a BSpline
            if spline.dxftype() == 'SPLINE':  # here we look for entity type SPLINE
                if spline.dxf.color==1 or spline.dxf.color==3:   # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline._control_points

                    # Print control points
                    for point in control_points:
                        i=i+1
                        print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline


    for spline in polylines:
        # Check if the spline is a BSpline
            if spline.dxftype() == 'POLYLINE':  # here we look for entity type SPLINE
                if spline.dxf.color==1 or spline.dxf.color==3:   # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline.points()
                    # Print control points
                    for point in control_points:
                        i=i+1
                        print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline

    qty_spline=0
    Layer_Names=[]
    for spline in lw_polylines:
        # Check if the spline is a BSpline
            if spline.dxftype() == 'LWPOLYLINE':  # here we look for entity type SPLINE
                 print(spline.dxf.color)
                 layer_name = spline.dxf.layer  # Get the layer name
                 layer = doc.layers.get(layer_name)  # Get the layer object
                 color = layer.dxf.color  # Get the color from the layer
                 print (color)
                 ######## GET THE GREEN A RED LAYERS
                 if color==1 or color==3:    # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline.get_points('xy')

                    # Print control points
                    for point in control_points:
                        i=i+1
                        print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline
                    #####................................................
                    if qty_spline == 0:
                        Layer_Names.append(layer_name)
                    elif Layer_Names[-1] != layer_name:
                        Layer_Names.append(layer_name)
                    qty_spline=qty_spline+1

                 ################### GET THE TAGGING LINES
                 if color == 7 or color == 255:
                     control_points = spline.get_points('xy')
                     for point in control_points:
                            i=i+1
                            print("Vertice:", point,'vertice number', i)
                            x_tag.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                            y_tag.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline

                 ################### GET THE DATUM LINES
                 if color == 141:
                     control_points = spline.get_points('xy')
                     for point in control_points:
                            i=i+1
                            print("Vertice:", point,'vertice number', i)
                            x_datum.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                            y_datum.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline




            #####...................1st scale reduction  OPTIONAL...........................
    # Reduce_factor=1
    x = [i / Reduce_factor for i in x]
    y = [i / Reduce_factor for i in y]

    excel_table_calib = r'C:\Users\Juan Pablo Lopez\PycharmProjects\ProjectJP\Table_Calib.xlsx'
    # .....LOAD THE TABLE CALIBRATION DATA.......................................................................................................
    if raspberry==1:
        X_table = pd.read_excel(excel_table_calib, sheet_name='X_axis_RP', header=0)
        Y_table = pd.read_excel(excel_table_calib, sheet_name='Y_axis_RP', header=0)
    else:
        X_table = pd.read_excel(excel_table_calib, sheet_name='X_axis', header=0)
        Y_table = pd.read_excel(excel_table_calib, sheet_name='Y_axis', header=0)

    if scale_mode==1:
        for i in range(0,len(x)):
            target_value = x[i]
            df=X_table
            # Find indices of bracketing points
            idx = df['Real_Measure'].searchsorted(target_value)

            # Check if target is within data range
            if idx == 0 or idx == len(x):
                print("Target value outside data range for interpolation")
            else:
                # Extract bracketed values
                a_prev = df.loc[idx - 1, 'Real_Measure']
                a_next = df.loc[idx, 'Real_Measure']
                b_prev = df.loc[idx - 1, 'Autocad_Measure']
                b_next = df.loc[idx, 'Autocad_Measure']

                # Perform linear interpolation
                interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                # print(f"Interpolated value for {target_value}: {interpolated_value}")
                x[i]=interpolated_value

        for i in range(0,len(y)):
            target_value = y[i]
            df=Y_table
            # Find indices of bracketing points
            idx = df['Real_Measure'].searchsorted(target_value)

            # Check if target is within data range
            if idx == 0 or idx == len(y):
                print("Target value outside data range for interpolation")
            else:
                # Extract bracketed values
                a_prev = df.loc[idx - 1, 'Real_Measure']
                a_next = df.loc[idx, 'Real_Measure']
                b_prev = df.loc[idx - 1, 'Autocad_Measure']
                b_next = df.loc[idx, 'Autocad_Measure']

                # Perform linear interpolation
                interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                # print(f"Interpolated value for {target_value}: {interpolated_value}")
                y[i]=interpolated_value

        if len(x_tag)>0:
            for i in range(0,len(x_tag)):
                target_value = x_tag[i]
                df=X_table
                # Find indices of bracketing points
                idx = df['Real_Measure'].searchsorted(target_value)

                # Check if target is within data range
                if idx == 0 or idx == len(x_tag):
                    print("Target value outside data range for interpolation")
                else:
                    # Extract bracketed values
                    a_prev = df.loc[idx - 1, 'Real_Measure']
                    a_next = df.loc[idx, 'Real_Measure']
                    b_prev = df.loc[idx - 1, 'Autocad_Measure']
                    b_next = df.loc[idx, 'Autocad_Measure']

                    # Perform linear interpolation
                    interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                    # print(f"Interpolated value for {target_value}: {interpolated_value}")
                    x_tag[i]=interpolated_value

            for i in range(0,len(y_tag)):
                target_value = y_tag[i]
                df=Y_table
                # Find indices of bracketing points
                idx = df['Real_Measure'].searchsorted(target_value)

                # Check if target is within data range
                if idx == 0 or idx == len(y_tag):
                    print("Target value outside data range for interpolation")
                else:
                    # Extract bracketed values
                    a_prev = df.loc[idx - 1, 'Real_Measure']
                    a_next = df.loc[idx, 'Real_Measure']
                    b_prev = df.loc[idx - 1, 'Autocad_Measure']
                    b_next = df.loc[idx, 'Autocad_Measure']

                    # Perform linear interpolation
                    interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                    # print(f"Interpolated value for {target_value}: {interpolated_value}")
                    y_tag[i]=interpolated_value

        if len(x_datum) > 0:
            for i in range(0, len(x_datum)):
                target_value = x_datum[i]
                df = X_table
                # Find indices of bracketing points
                idx = df['Real_Measure'].searchsorted(target_value)

                # Check if target is within data range
                if idx == 0 or idx == len(x_tag):
                    print("Target value outside data range for interpolation")
                else:
                    # Extract bracketed values
                    a_prev = df.loc[idx - 1, 'Real_Measure']
                    a_next = df.loc[idx, 'Real_Measure']
                    b_prev = df.loc[idx - 1, 'Autocad_Measure']
                    b_next = df.loc[idx, 'Autocad_Measure']

                    # Perform linear interpolation
                    interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                    # print(f"Interpolated value for {target_value}: {interpolated_value}")
                    x_datum[i] = interpolated_value

            for i in range(0, len(y_datum)):
                target_value = y_datum[i]
                df = Y_table
                # Find indices of bracketing points
                idx = df['Real_Measure'].searchsorted(target_value)

                # Check if target is within data range
                if idx == 0 or idx == len(y_tag):
                    print("Target value outside data range for interpolation")
                else:
                    # Extract bracketed values
                    a_prev = df.loc[idx - 1, 'Real_Measure']
                    a_next = df.loc[idx, 'Real_Measure']
                    b_prev = df.loc[idx - 1, 'Autocad_Measure']
                    b_next = df.loc[idx, 'Autocad_Measure']

                    # Perform linear interpolation
                    interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                    # print(f"Interpolated value for {target_value}: {interpolated_value}")
                    y_datum[i] = interpolated_value

    ############# BEGGIN THE PLOTTING

    Axis_Limit_x = Axis_Limit
    Axis_Limit_y = Axis_Limit

    if pattern_mode==1:
        x=x_tag
        y=y_tag
        x_tag=[]
        y_tag=[]

    x_init=x[0]
    y_init=y[0]
    z=0
    layer_qty=0
    check=0

    for i in range(0,len(x)):    # loop through all the vertices of the splines found
        #print(i)
        if x[i]==x_init and y[i]==y_init:   # layer detection condition is when there is a closed loop, initial and final points are equal
            check=check+1 # the initial point of the loop
            if check==2: # the final point of the loop
                if pattern_mode==1:
                    Layer_Name_plot = 'PATTERN'
                else:
                    Layer_Name_plot=Layer_Names[layer_qty]
                x_layer = x[z:i + 1]
                y_layer = y[z:i + 1]
                z=i+1
                layer_qty=layer_qty+1
                check=0
                if i<(len(x)-1):
                 x_init = x[i+1]
                 y_init = y[i+1]
                print('Layer ', layer_qty)
                print('X_coordenates ', x_layer)
                print('y_coordenates ', y_layer)

                # Plotting the polygon
                plt.figure(figsize=(image_size, image_size))
                plt.plot(x_layer, y_layer, color=layer_color, linestyle='-', linewidth=0.1, marker='o', markersize=0.1) # 'bo-' means blue circles connected by lines
                plt.fill(x_layer, y_layer,facecolor = layer_color ,alpha=1)  # Fill the polygon
                title='Layer '+str(layer_qty)
                plt.title(title)
                #plt.grid(True, linestyle='--', linewidth=0.5)  # Set linewidth to 1.5
                plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
                plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
                plt.style.use('dark_background')
                ax = plt.gca()
                ax.set_facecolor(background_color)
                ax.spines['top'].set_color(layer_color)
                ax.spines['bottom'].set_color(layer_color)
                ax.spines['right'].set_color(layer_color)
                ax.spines['left'].set_color(layer_color)
                ax.tick_params(axis='x', colors=layer_color)
                ax.tick_params(axis='y', colors=layer_color)
                ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
                plt.savefig('Layers.png')
                # add slide
                slide = prs.slides.add_slide(slide_layout)
                img_path = 'Layers.png'
                left = Inches(Left_centering)  # Adjust position as needed in icnhes
                top = Inches(Top_centering)  # Adjust position as needed in inches
                slide.shapes.add_picture(img_path, left, top)
                if close_image==1:
                 plt.close('all')

    if len(x_tag)>0:
                x_layer=x_tag
                y_layer=y_tag
                Layer_Name_plot='TAGGING'

                # Plotting the polygon
                plt.figure(figsize=(image_size, image_size))
                plt.scatter(x_layer, y_layer,marker='x', s=3, color='blue') # 'bo-' means blue circles connected by lines
                title='Tagging'
                plt.title(title)
                plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
                plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
                plt.style.use('dark_background')
                ax = plt.gca()
                ax.set_facecolor(background_color)
                ax.spines['top'].set_color(layer_color)
                ax.spines['bottom'].set_color(layer_color)
                ax.spines['right'].set_color(layer_color)
                ax.spines['left'].set_color(layer_color)
                ax.tick_params(axis='x', colors=layer_color)
                ax.tick_params(axis='y', colors=layer_color)
                ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
                plt.savefig('Layers.png')
                # add slide
                slide = prs.slides.add_slide(slide_layout)
                img_path = 'Layers.png'
                left = Inches(Left_centering)  # Adjust position as needed in icnhes
                top = Inches(Top_centering)  # Adjust position as needed in inches
                slide.shapes.add_picture(img_path, left, top)
                if close_image==1:
                 plt.close('all')

                # prs.save('Layers.pptx')



    if len(x_datum)>0:
                x_layer=x_datum
                y_layer=y_datum
                Layer_Name_plot='DATUM LINE'

                # Plotting the polygon
                plt.figure(figsize=(image_size, image_size))
                plt.scatter(x_layer, y_layer,marker='.', s=5, color='blue') # 'bo-' means blue circles connected by lines
                title='Tagging'
                plt.title(title)
                plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
                plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
                plt.style.use('dark_background')
                ax = plt.gca()
                ax.set_facecolor(background_color)
                ax.spines['top'].set_color(layer_color)
                ax.spines['bottom'].set_color(layer_color)
                ax.spines['right'].set_color(layer_color)
                ax.spines['left'].set_color(layer_color)
                ax.tick_params(axis='x', colors=layer_color)
                ax.tick_params(axis='y', colors=layer_color)
                ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
                plt.savefig('Layers.png')
                # add slide
                slide = prs.slides.add_slide(slide_layout)
                img_path = 'Layers.png'
                left = Inches(Left_centering)  # Adjust position as needed in icnhes
                top = Inches(Top_centering)  # Adjust position as needed in inches
                slide.shapes.add_picture(img_path, left, top)
                if close_image==1:
                 plt.close('all')

                # prs.save('Layers.pptx')


    # Save the PowerPoint presentation
    prs.save('Layers.pptx')


    print('total layers = ',layer_qty)
def read_and_plot_layerv11(file_path,image_size,Left_centering,Top_centering,Axis_Limit,scale_mode,Reduce_factor,background_color,layer_color,close_image,pattern_mode,raspberry):
    ### error correction using absolute error to define layers, doesn work yet
    import ezdxf
    import pandas as pd
    import matplotlib.pyplot as plt
    from pptx import Presentation
    from pptx.util import Inches


    ##............READING THE DXF
    doc = ezdxf.readfile(file_path)

    # POWERPOINT SLIDE
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank slide layout

    # Extract entities (splines)
    msp = doc.modelspace()
    for entity in msp:
        # Get the type of the entity
        print(entity.dxftype())
    splines = msp.query('SPLINE')
    polylines= msp.query('POLYLINE')
    lw_polylines= msp.query('LWPOLYLINE')
    i=0
    x=[]
    y=[]
    x_tag=[]
    y_tag=[]
    x_datum=[]
    y_datum=[]
    ####.................... Iterate over splines/polylines/lw-polyinies
    for spline in splines:
            # Check if the spline is a BSpline
            if spline.dxftype() == 'SPLINE':  # here we look for entity type SPLINE
                if spline.dxf.color==1 or spline.dxf.color==3:   # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline._control_points

                    # Print control points
                    for point in control_points:
                        i=i+1
                        print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline


    for spline in polylines:
        # Check if the spline is a BSpline
            if spline.dxftype() == 'POLYLINE':  # here we look for entity type SPLINE
                if spline.dxf.color==1 or spline.dxf.color==3:   # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline.points()
                    # Print control points
                    for point in control_points:
                        i=i+1
                        print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline

    qty_spline=0
    Layer_Names=[]
    for spline in lw_polylines:
        # Check if the spline is a BSpline
            if spline.dxftype() == 'LWPOLYLINE':  # here we look for entity type SPLINE
                 print(spline.dxf.color)
                 layer_name = spline.dxf.layer  # Get the layer name
                 layer = doc.layers.get(layer_name)  # Get the layer object
                 color = layer.dxf.color  # Get the color from the layer
                 print (color)
                 ######## GET THE GREEN A RED LAYERS
                 if color==1 or color==3:    # check if its not the Frame
                    # Get control points of the spline
                    control_points = spline.get_points('xy')

                    # Print control points
                    for point in control_points:
                        i=i+1
                        print("Vertice:", point,'vertice number', i)
                        x.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                        y.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline
                    #####................................................
                    if qty_spline == 0:
                        Layer_Names.append(layer_name)
                    elif Layer_Names[-1] != layer_name:
                        Layer_Names.append(layer_name)
                    qty_spline=qty_spline+1

                 ################### GET THE TAGGING LINES
                 if color == 7 or color == 255:
                     control_points = spline.get_points('xy')
                     for point in control_points:
                            i=i+1
                            print("Vertice:", point,'vertice number', i)
                            x_tag.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                            y_tag.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline

                 ################### GET THE DATUM LINES
                 if color == 141:
                     control_points = spline.get_points('xy')
                     for point in control_points:
                            i=i+1
                            print("Vertice:", point,'vertice number', i)
                            x_datum.append(round(point[0],8)) # Here we extract the x-coordenates of the vertices of the spline
                            y_datum.append(round(point[1],8)) # Here we extract the y-coordenates of the vertices of the spline




            #####...................1st scale reduction  OPTIONAL...........................
    # Reduce_factor=1
    x = [i / Reduce_factor for i in x]
    y = [i / Reduce_factor for i in y]

    excel_table_calib = r'C:\Users\Juan Pablo Lopez\PycharmProjects\ProjectJP\Table_Calib.xlsx'
    # .....LOAD THE TABLE CALIBRATION DATA.......................................................................................................
    if raspberry==1:
        X_table = pd.read_excel(excel_table_calib, sheet_name='X_axis_RP', header=0)
        Y_table = pd.read_excel(excel_table_calib, sheet_name='Y_axis_RP', header=0)
    else:
        X_table = pd.read_excel(excel_table_calib, sheet_name='X_axis', header=0)
        Y_table = pd.read_excel(excel_table_calib, sheet_name='Y_axis', header=0)

    if scale_mode==1:
        for i in range(0,len(x)):
            target_value = x[i]
            df=X_table
            # Find indices of bracketing points
            idx = df['Real_Measure'].searchsorted(target_value)

            # Check if target is within data range
            if idx == 0 or idx == len(x):
                print("Target value outside data range for interpolation")
            else:
                # Extract bracketed values
                a_prev = df.loc[idx - 1, 'Real_Measure']
                a_next = df.loc[idx, 'Real_Measure']
                b_prev = df.loc[idx - 1, 'Autocad_Measure']
                b_next = df.loc[idx, 'Autocad_Measure']

                # Perform linear interpolation
                interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                # print(f"Interpolated value for {target_value}: {interpolated_value}")
                x[i]=interpolated_value

        for i in range(0,len(y)):
            target_value = y[i]
            df=Y_table
            # Find indices of bracketing points
            idx = df['Real_Measure'].searchsorted(target_value)

            # Check if target is within data range
            if idx == 0 or idx == len(y):
                print("Target value outside data range for interpolation")
            else:
                # Extract bracketed values
                a_prev = df.loc[idx - 1, 'Real_Measure']
                a_next = df.loc[idx, 'Real_Measure']
                b_prev = df.loc[idx - 1, 'Autocad_Measure']
                b_next = df.loc[idx, 'Autocad_Measure']

                # Perform linear interpolation
                interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                # print(f"Interpolated value for {target_value}: {interpolated_value}")
                y[i]=interpolated_value

        if len(x_tag)>0:
            for i in range(0,len(x_tag)):
                target_value = x_tag[i]
                df=X_table
                # Find indices of bracketing points
                idx = df['Real_Measure'].searchsorted(target_value)

                # Check if target is within data range
                if idx == 0 or idx == len(x_tag):
                    print("Target value outside data range for interpolation")
                else:
                    # Extract bracketed values
                    a_prev = df.loc[idx - 1, 'Real_Measure']
                    a_next = df.loc[idx, 'Real_Measure']
                    b_prev = df.loc[idx - 1, 'Autocad_Measure']
                    b_next = df.loc[idx, 'Autocad_Measure']

                    # Perform linear interpolation
                    interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                    # print(f"Interpolated value for {target_value}: {interpolated_value}")
                    x_tag[i]=interpolated_value

            for i in range(0,len(y_tag)):
                target_value = y_tag[i]
                df=Y_table
                # Find indices of bracketing points
                idx = df['Real_Measure'].searchsorted(target_value)

                # Check if target is within data range
                if idx == 0 or idx == len(y_tag):
                    print("Target value outside data range for interpolation")
                else:
                    # Extract bracketed values
                    a_prev = df.loc[idx - 1, 'Real_Measure']
                    a_next = df.loc[idx, 'Real_Measure']
                    b_prev = df.loc[idx - 1, 'Autocad_Measure']
                    b_next = df.loc[idx, 'Autocad_Measure']

                    # Perform linear interpolation
                    interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                    # print(f"Interpolated value for {target_value}: {interpolated_value}")
                    y_tag[i]=interpolated_value

        if len(x_datum) > 0:
            for i in range(0, len(x_datum)):
                target_value = x_datum[i]
                df = X_table
                # Find indices of bracketing points
                idx = df['Real_Measure'].searchsorted(target_value)

                # Check if target is within data range
                if idx == 0 or idx == len(x_tag):
                    print("Target value outside data range for interpolation")
                else:
                    # Extract bracketed values
                    a_prev = df.loc[idx - 1, 'Real_Measure']
                    a_next = df.loc[idx, 'Real_Measure']
                    b_prev = df.loc[idx - 1, 'Autocad_Measure']
                    b_next = df.loc[idx, 'Autocad_Measure']

                    # Perform linear interpolation
                    interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                    # print(f"Interpolated value for {target_value}: {interpolated_value}")
                    x_datum[i] = interpolated_value

            for i in range(0, len(y_datum)):
                target_value = y_datum[i]
                df = Y_table
                # Find indices of bracketing points
                idx = df['Real_Measure'].searchsorted(target_value)

                # Check if target is within data range
                if idx == 0 or idx == len(y_tag):
                    print("Target value outside data range for interpolation")
                else:
                    # Extract bracketed values
                    a_prev = df.loc[idx - 1, 'Real_Measure']
                    a_next = df.loc[idx, 'Real_Measure']
                    b_prev = df.loc[idx - 1, 'Autocad_Measure']
                    b_next = df.loc[idx, 'Autocad_Measure']

                    # Perform linear interpolation
                    interpolated_value = b_prev + ((target_value - a_prev) / (a_next - a_prev)) * (b_next - b_prev)
                    # print(f"Interpolated value for {target_value}: {interpolated_value}")
                    y_datum[i] = interpolated_value

    ############# BEGGIN THE PLOTTING

    Axis_Limit_x = Axis_Limit
    Axis_Limit_y = Axis_Limit

    if pattern_mode==1:
        x=x_tag
        y=y_tag
        x_tag=[]
        y_tag=[]

    x_init=x[0]
    y_init=y[0]
    z=0
    layer_qty=0
    check=0

    for i in range(0,len(x)):    # loop through all the vertices of the splines found
        #print(i)
        error_x=abs(x[i]-x_init)
        error_y = abs(y[i] - y_init)
        if error_x<80 and error_y<80:   # layer detection condition is when there is a closed loop, initial and final points are equal
            check=check+1 # the initial point of the loop
            if check==2: # the final point of the loop
                if pattern_mode==1:
                    Layer_Name_plot = 'PATTERN'
                else:
                    Layer_Name_plot=Layer_Names[layer_qty]
                x_layer = x[z:i + 1]
                y_layer = y[z:i + 1]
                z=i+1
                layer_qty=layer_qty+1
                check=0
                if i<(len(x)-1):
                 x_init = x[i+1]
                 y_init = y[i+1]
                print('Layer ', layer_qty)
                print('X_coordenates ', x_layer)
                print('y_coordenates ', y_layer)

                # Plotting the polygon
                plt.figure(figsize=(image_size, image_size))
                plt.plot(x_layer, y_layer, color=layer_color, linestyle='-', linewidth=0.1, marker='o', markersize=0.1) # 'bo-' means blue circles connected by lines
                plt.fill(x_layer, y_layer,facecolor = layer_color ,alpha=1)  # Fill the polygon
                title='Layer '+str(layer_qty)
                plt.title(title)
                #plt.grid(True, linestyle='--', linewidth=0.5)  # Set linewidth to 1.5
                plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
                plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
                plt.style.use('dark_background')
                ax = plt.gca()
                ax.set_facecolor(background_color)
                ax.spines['top'].set_color(layer_color)
                ax.spines['bottom'].set_color(layer_color)
                ax.spines['right'].set_color(layer_color)
                ax.spines['left'].set_color(layer_color)
                ax.tick_params(axis='x', colors=layer_color)
                ax.tick_params(axis='y', colors=layer_color)
                ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
                plt.savefig('Layers.png')
                # add slide
                slide = prs.slides.add_slide(slide_layout)
                img_path = 'Layers.png'
                left = Inches(Left_centering)  # Adjust position as needed in icnhes
                top = Inches(Top_centering)  # Adjust position as needed in inches
                slide.shapes.add_picture(img_path, left, top)
                if close_image==1:
                 plt.close('all')

    if len(x_tag)>0:
                x_layer=x_tag
                y_layer=y_tag
                Layer_Name_plot='TAGGING'

                # Plotting the polygon
                plt.figure(figsize=(image_size, image_size))
                plt.scatter(x_layer, y_layer,marker='x', s=3, color='blue') # 'bo-' means blue circles connected by lines
                title='Tagging'
                plt.title(title)
                plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
                plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
                plt.style.use('dark_background')
                ax = plt.gca()
                ax.set_facecolor(background_color)
                ax.spines['top'].set_color(layer_color)
                ax.spines['bottom'].set_color(layer_color)
                ax.spines['right'].set_color(layer_color)
                ax.spines['left'].set_color(layer_color)
                ax.tick_params(axis='x', colors=layer_color)
                ax.tick_params(axis='y', colors=layer_color)
                ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
                plt.savefig('Layers.png')
                # add slide
                slide = prs.slides.add_slide(slide_layout)
                img_path = 'Layers.png'
                left = Inches(Left_centering)  # Adjust position as needed in icnhes
                top = Inches(Top_centering)  # Adjust position as needed in inches
                slide.shapes.add_picture(img_path, left, top)
                if close_image==1:
                 plt.close('all')

                # prs.save('Layers.pptx')



    if len(x_datum)>0:
                x_layer=x_datum
                y_layer=y_datum
                Layer_Name_plot='DATUM LINE'

                # Plotting the polygon
                plt.figure(figsize=(image_size, image_size))
                plt.scatter(x_layer, y_layer,marker='.', s=5, color='blue') # 'bo-' means blue circles connected by lines
                title='Tagging'
                plt.title(title)
                plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
                plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
                plt.style.use('dark_background')
                ax = plt.gca()
                ax.set_facecolor(background_color)
                ax.spines['top'].set_color(layer_color)
                ax.spines['bottom'].set_color(layer_color)
                ax.spines['right'].set_color(layer_color)
                ax.spines['left'].set_color(layer_color)
                ax.tick_params(axis='x', colors=layer_color)
                ax.tick_params(axis='y', colors=layer_color)
                ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
                plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
                plt.savefig('Layers.png')
                # add slide
                slide = prs.slides.add_slide(slide_layout)
                img_path = 'Layers.png'
                left = Inches(Left_centering)  # Adjust position as needed in icnhes
                top = Inches(Top_centering)  # Adjust position as needed in inches
                slide.shapes.add_picture(img_path, left, top)
                if close_image==1:
                 plt.close('all')

                # prs.save('Layers.pptx')


    # Save the PowerPoint presentation
    prs.save('Layers.pptx')


    print('total layers = ',layer_qty)



file_path='c97_CUT11_Tagging.dxf'     ####
image_size = 12.7  # in inches    12.7
Left_centering = -1.55  # in inches   -1.55
Top_centering = -4.9  # in inches   -4.9
Axis_Limit=8000  #in MM
Reduce_factor=1  #  default = 1, if not it is used for scaling down the original image by a factor, eg 10,100,1000
background_color='white'
layer_color='blue'
close_image=1   # 1- close all images, 2- open all images
scale_mode=1   # 1-apply scale, 2- no scale
pattern_mode=0  # 1- loading a pattern dxf file, 2-  normal dxf
raspberry=1      # 1- using the raspberry pi, 2- using the laptop


#### V10-LATEST VERSION, DATUM LINES ADDED in the form of scatter
read_and_plot_layerv10(file_path,image_size,Left_centering,Top_centering,Axis_Limit,scale_mode,Reduce_factor,background_color,layer_color,close_image,pattern_mode,raspberry)