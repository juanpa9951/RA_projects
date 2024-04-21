import ezdxf
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches


##............READING THE DXF AND EXTRACTING THE POLYGON LAYERS
# Load DXF file
#doc = ezdxf.readfile("36_LRU12.dxf")
doc = ezdxf.readfile("spline1.dxf")


# Extract entities (spline)
msp = doc.modelspace()
splines = msp.query('SPLINE')
i=0
x=[]
y=[]
# Iterate over splines
for spline in splines:
    # Check if the spline is a BSpline
    if spline.dxftype() == 'SPLINE':
        # Get control points
        control_points = spline._control_points

        # Print control points
        for point in control_points:
            i=i+1
            print("Vertex:", point,'verte number', i)
            x.append(round(point[0],2))
            y.append(round(point[1]))


data = {'X': x,
        'Y': y}

# Create DataFrame
df = pd.DataFrame(data)

# Display DataFrame
print(df)

df_no_duplicates = df.drop_duplicates()

# Display DataFrame without duplicates
print(df_no_duplicates)


# PLOTTING THE POLYGON

# Coordinates of the vertices of the polygon
# x = [0, 2, 4, 2]  # x-coordinates
# y = [1, 3, 2, 1]  # y-coordinates

# Plotting the polygon
plt.figure(figsize=(7, 7))
#plt.plot(x, y, 'bo-')  # 'bo-' means blue circles connected by lines
plt.plot(x, y, color='red', linestyle='-', linewidth=2, marker='o', markersize=1)
plt.fill(x, y, alpha=1)  # Fill the polygon
#plt.xlabel('X-axis')
#plt.ylabel('Y-axis')
plt.title('Layer 1')
plt.grid(True, linestyle='--', linewidth=0.5)  # Set linewidth to 1.5

# plt.axis('equal')  # Equal aspect ratio
# Set fixed size for x and y axes
plt.xlim(0, 4000)  # Set x-axis limit from 0 to 4
plt.ylim(0, 4000)  # Set y-axis limit from 0 to 4
plt.savefig('plot.png')
plt.show()

# POWERPOINT SLIDE
# Create a PowerPoint presentation
prs = Presentation()
# select slide Layout
slide_layout = prs.slide_layouts[6]  # Blank slide layout
#add slide
slide = prs.slides.add_slide(slide_layout)
# Add the plot image to the slide
img_path = 'plot.png'
left = Inches(0)  # Adjust position as needed in icnhes
top = Inches(0)   # Adjust position as needed in inches
#slide.shapes.add_picture(img_path, left, top, width=Inches(5))  # Adjust width as needed
slide.shapes.add_picture(img_path, left, top)

#ADD ANOTHER SLIDE....................................
# Plotting the polygon
plt.figure(figsize=(7, 7))
#plt.plot(x, y, 'bo-')  # 'bo-' means blue circles connected by lines
plt.plot(x, y, color='blue', linestyle='-', linewidth=2, marker='o', markersize=1)
plt.fill(x, y, alpha=1)  # Fill the polygon
#plt.xlabel('X-axis')
#plt.ylabel('Y-axis')
plt.title('Layer 1')
plt.grid(True, linestyle='--', linewidth=0.5)  # Set linewidth to 1.5
# plt.axis('equal')  # Equal aspect ratio
# Set fixed size for x and y axes
plt.xlim(0, 7000)  # Set x-axis limit from 0 to 4
plt.ylim(0, 7000)  # Set y-axis limit from 0 to 4
plt.savefig('plot.png')
plt.show()
slide = prs.slides.add_slide(slide_layout)
# Add the plot image to the slide
img_path = 'plot.png'
left = Inches(0)  # Adjust position as needed in inches
top = Inches(0)   # Adjust position as needed in inches
#slide.shapes.add_picture(img_path, left, top, width=Inches(5))  # Adjust width as needed
slide.shapes.add_picture(img_path, left, top)
#....................................................................

# Save the PowerPoint presentation
prs.save('plot.pptx')

