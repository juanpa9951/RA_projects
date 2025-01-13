def read_and_plot_PSv0(file_path,name,image_size,Left_centering,Top_centering,Axis_Limit,scale_mode,background_color,layer_color,close_image,Destination_path,Axis_Limit_Y,filling,reorder_points):

    import ezdxf
    import pandas as pd
    import matplotlib.pyplot as plt
    from pptx import Presentation
    from pptx.util import Inches
    from dxf_surface_map import euclidean_distance,find_closest_tupleV2
    from dxf_ProcessKit_polygon import polygon_sorter


    ##............READING THE DXF
    doc = ezdxf.readfile(file_path)

    # POWERPOINT SLIDE
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank slide layout

    # Extract entities (splines)
    msp = doc.modelspace()
    for entity in msp:
        # Get the type of the entity
        print(entity.dxftype())
    splines = msp.query('SPLINE')
    polylines= msp.query('POLYLINE')
    lw_polylines= msp.query('LWPOLYLINE')
    lines = msp.query('LINE')
    i=0
    x=[]
    y=[]
    x_omega=[]
    y_omega=[]
    color=[]

    ####.................... Iterate over splines/polylines/lw-polyinies/lines


    qty_spline=0
    Layer_Names=[]
    Layer_points=[]

    for spline in lw_polylines:
        # Check if the spline is a BSpline
            if spline.dxftype() == 'LWPOLYLINE':  # here we look for entity type SPLINE
                layer_name = spline.dxf.layer  # Get the layer name
                layer = doc.layers.get(layer_name)  # Get the layer object
                #color = layer.dxf.color  # Get the color from the layer
                control_points = spline.get_points('xy')

                # Print control points
                for point in control_points:
                    i = i + 1
                    # print("Vertice:", point,'vertice number', i)
                    x_omega.append(round(point[0], 8))  # Here we extract the x-coordenates of the vertices of the spline
                    y_omega.append(round(point[1], 8))  # Here we extract the y-coordenates of the vertices of the spline


    for spline in lines:
        # Check if the spline is a BSpline
            if spline.dxftype() == 'LINE':  # here we look for entity type LINE
                # Get control points of the spline

                layer_name = spline.dxf.layer  # Get the layer name
                # layer = doc.layers.get(layer_name)  # Get the layer object
                # color = layer.dxf.color  # Get the color from the layer

                control_points = [spline.dxf.start, spline.dxf.end]
                for point in control_points:
                    x.append(round(point[0],6))  # Here we extract the x-coordenates of the vertices of the spline rounded to 8 decimals
                    y.append(round(point[1],6))  # Here we extract the y-coordenates of the vertices of the spline  8 rounded to 8 decimals
                    color.append(spline.dxf.color)
                    Layer_points.append(layer_name)

                #####................................................
                if qty_spline == 0:
                    Layer_Names.append(layer_name)
                elif Layer_Names[-1] != layer_name:
                    Layer_Names.append(layer_name)
                qty_spline=qty_spline+1

    #####....................................SHIFT POSITION .............................
    #### BRING TO THE ORIGIN
    cx1=min(x)
    cy1=min(y)
    x = [i -cx1 for i in x]  # BRING X VALUES TO THE ORIGIN
    y = [i -cy1 for i in y]  # BRING Y VALUES TO THE ORIGIN
    x_omega = [i -cx1 for i in x_omega]  # BRING X VALUES TO THE ORIGIN
    y_omega = [i -cy1 for i in y_omega]  # BRING Y VALUES TO THE ORIGIN


    #### APPLY THE OFFSET (DESFASE)
    x_lim=7000
    y_lim=5000
    cx=x_lim-max(x)
    cy=y_lim-max(y)
    x = [i + cx for i in x]  # apply offset desfase
    y = [i + cy for i in y]  # apply offset desfase
    x_omega = [i + cx for i in x_omega]  # apply offset desfase
    y_omega = [i + cy for i in y_omega]  # apply offset desfase
    ########..............................................................................................................

    # .....LOAD THE TABLE CALIBRATION DATA.......................................................................................................
    excel_table_calib = r'C:\Users\Juan Pablo Lopez\OneDrive - Rewair A S\Documents\Proyector OPTOMA\table_calib_optoma.xlsx'
    Surface_map = pd.read_excel(excel_table_calib, sheet_name='Surface', header=0)
    tuples_list_real=[]
    tuples_list_autocad=[]
    for i in range(0,len(Surface_map)):
        x_tup_autocad = Surface_map['Autocad_X'][i]
        y_tup_autocad = Surface_map['Autocad_Y'][i]
        tuples_list_autocad.append((x_tup_autocad,y_tup_autocad))
        x_tup_real = Surface_map['Real_X'][i]
        y_tup_real = Surface_map['Real_Y'][i]
        tuples_list_real.append((x_tup_real,y_tup_real))

    #####......................... APPLY THE CALIBRATION......................................
    if scale_mode==1:
        for i in range(0,len(x)):
            input_tuple = (x[i],y[i])
            print(input_tuple)
            xy_tup,idx_low,idx_high,tup_low,tup_high = find_closest_tupleV2(tuples_list_real,tuples_list_autocad,input_tuple)
            x[i] = xy_tup[0]
            y[i] = xy_tup[1]
            #print("Original ",input_tuple," Converted ",xy_tup," tup_low=",tup_low," idx low ",idx_low," tup_high=",tup_high, " idx high ",idx_high)

        if len(x_omega)>0:
            for i in range(0,len(x_omega)):
                input_tuple = (x_omega[i], y_omega[i])
                # print(i,input_tuple)
                xy_tup, idx_low, idx_high, tup_low, tup_high = find_closest_tupleV2(tuples_list_real, tuples_list_autocad, input_tuple)
                x_omega[i] = xy_tup[0]
                y_omega[i] = xy_tup[1]
                # print("Original ", input_tuple, " Converted ", xy_tup, " tup_low=", tup_low, " idx low ", idx_low," tup_high=", tup_high, " idx high ", idx_high)

    ############# BEGGIN THE PLOTTING...............................

    Axis_Limit_x = Axis_Limit
    Axis_Limit_y = Axis_Limit_Y



    layer_qty=0
    x_layer=[]
    y_layer=[]


    seen = set()
    Layer_Names = [x for x in Layer_Names if not (x in seen or seen.add(x))]


    for Layer_i in Layer_Names:    # loop through all the vertices of the splines found
        for i in range(0, len(x)):
            if Layer_points[i]==Layer_i:
                x_layer.append(x[i])
                y_layer.append(y[i])

        x_layer,y_layer=polygon_sorter(x_layer,y_layer)

        Layer_Name_plot = Layer_i
        layer_qty = layer_qty + 1
        # Plotting the polygon
        plt.figure(figsize=(image_size, image_size))
        plt.scatter(x_layer, y_layer, marker='o', s=4, color='blue')  #  s=0.1  'bo-' means blue circles connected by lines
        #plt.fill(x_layer, y_layer, facecolor=layer_color, alpha=1)  # Fill the polygon
        title = 'Layer '+str(layer_qty)
        plt.title(title)
        plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
        plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
        plt.style.use('dark_background')
        ax = plt.gca()
        ax.set_facecolor(background_color)
        ax.spines['top'].set_color(layer_color)
        ax.spines['bottom'].set_color(layer_color)
        ax.spines['right'].set_color(layer_color)
        ax.spines['left'].set_color(layer_color)
        ax.tick_params(axis='x', colors=layer_color)
        ax.tick_params(axis='y', colors=layer_color)
        ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
        plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
        plt.savefig('Layers.png')
        # add slide
        slide = prs.slides.add_slide(slide_layout)
        img_path = 'Layers.png'
        left = Inches(Left_centering)  # Adjust position as needed in icnhes
        top = Inches(Top_centering)  # Adjust position as needed in inches
        slide.shapes.add_picture(img_path, left, top)
        if close_image == 1:
            plt.close('all')
        #print(title)
        x_layer = []   ### reset x_layer
        y_layer = []   ### reset y_layer




    # if len(x_omega)>0:
    #             x_layer=x_omega
    #             y_layer=y_omega
    #             Layer_Name_plot='OMEGA'
    #
    #             # Plotting the polygon
    #             plt.figure(figsize=(image_size, image_size))
    #             plt.scatter(x_layer, y_layer,marker='x', s=3, color='blue') # 'bo-' means blue circles connected by lines
    #             title='Omega'
    #             plt.title(title)
    #             plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
    #             plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
    #             plt.style.use('dark_background')
    #             ax = plt.gca()
    #             ax.set_facecolor(background_color)
    #             ax.spines['top'].set_color(layer_color)
    #             ax.spines['bottom'].set_color(layer_color)
    #             ax.spines['right'].set_color(layer_color)
    #             ax.spines['left'].set_color(layer_color)
    #             ax.tick_params(axis='x', colors=layer_color)
    #             ax.tick_params(axis='y', colors=layer_color)
    #             ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
    #             plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
    #             plt.savefig('Layers.png')
    #             # add slide
    #             slide = prs.slides.add_slide(slide_layout)
    #             img_path = 'Layers.png'
    #             left = Inches(Left_centering)  # Adjust position as needed in icnhes
    #             top = Inches(Top_centering)  # Adjust position as needed in inches
    #             slide.shapes.add_picture(img_path, left, top)
    #             if close_image==1:
    #              plt.close('all')

                # prs.save('Layers.pptx')

    # Save the PowerPoint presentation
    Name0=name[:-4]
    Name2 = Name0 + '.pptx'
    out_Name = f"{Destination_path}\{Name2}"
    prs.save(out_Name)


    print('total layers = ',layer_qty)
def read_and_plot_PSv1(file_path,name,image_size,Left_centering,Top_centering,Axis_Limit,scale_mode,background_color,layer_color,close_image,Destination_path,Axis_Limit_Y,filling,reorder_points):

    import ezdxf
    import pandas as pd
    import matplotlib.pyplot as plt
    from pptx import Presentation
    from pptx.util import Inches
    from dxf_surface_map import euclidean_distance,find_closest_tupleV2
    from dxf_ProcessKit_polygon import polygon_sorter


    ##............READING THE DXF
    doc = ezdxf.readfile(file_path)

    # POWERPOINT SLIDE
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank slide layout

    # Extract entities (splines)
    msp = doc.modelspace()
    for entity in msp:
        # Get the type of the entity
        print(entity.dxftype())
    splines = msp.query('SPLINE')
    polylines= msp.query('POLYLINE')
    lw_polylines= msp.query('LWPOLYLINE')
    lines = msp.query('LINE')
    i=0
    x=[]
    y=[]
    x_omega=[]
    y_omega=[]
    color=[]

    ####.................... Iterate over splines/polylines/lw-polyinies/lines


    qty_spline=0
    Layer_Names=[]
    Layer_points=[]

    for spline in lines:
        # Check if the spline is a BSpline
            if spline.dxftype() == 'LINE':  # here we look for entity type LINE
                # Get control points of the spline

                layer_name = spline.dxf.layer  # Get the layer name
                # layer = doc.layers.get(layer_name)  # Get the layer object
                # color = layer.dxf.color  # Get the color from the layer

                control_points = [spline.dxf.start, spline.dxf.end]
                for point in control_points:
                    x.append(round(point[0],6))  # Here we extract the x-coordenates of the vertices of the spline rounded to 8 decimals
                    y.append(round(point[1],6))  # Here we extract the y-coordenates of the vertices of the spline  8 rounded to 8 decimals
                    color.append(spline.dxf.color)
                    Layer_points.append(layer_name)
                    Layer_Names.append(layer_name)



    for spline in splines:
        # Check if the spline is a BSpline

        if spline.dxftype() == 'SPLINE':  # here we look for entity type SPLINE
                control_points = spline._control_points
                layer_name = spline.dxf.layer
                # EXTRACT control points, i.e  VERTICES OF THE SPLINES
                for point in control_points:
                    x.append(round(point[0], 6))  # Here we extract the x-coordenates of the vertices of the spline rounded to 8 decimals
                    y.append(round(point[1], 6))  # Here we extract the y-coordenates of the vertices of the spline  8 rounded to 8 decimals
                    color.append(spline.dxf.color)
                    Layer_points.append(layer_name)
                    Layer_Names.append(layer_name)


    for spline in lw_polylines:
        # Check if the spline is a BSpline
            if spline.dxftype() == 'LWPOLYLINE':  # here we look for entity type SPLINE
                layer_name = spline.dxf.layer  # Get the layer name
                control_points = spline.get_points('xy')

                # Print control points
                for point in control_points:
                    i = i + 1
                    # print("Vertice:", point,'vertice number', i)
                    x.append(round(point[0], 6))  # Here we extract the x-coordenates of the vertices of the spline
                    y.append(round(point[1], 6))  # Here we extract the y-coordenates of the vertices of the spline
                    Layer_points.append(layer_name)
                    Layer_Names.append(layer_name)


    #### REMOVE DUPLICATES IN THE LAYERS NAMES
    seen = set()
    Layer_Names = [x for x in Layer_Names if not (x in seen or seen.add(x))]

    #####....................................SHIFT POSITION .............................
    #### BRING TO THE ORIGIN
    cx1=min(x)
    cy1=min(y)
    x = [i -cx1 for i in x]  # BRING X VALUES TO THE ORIGIN
    y = [i -cy1 for i in y]  # BRING Y VALUES TO THE ORIGIN
    x_omega = [i -cx1 for i in x_omega]  # BRING X VALUES TO THE ORIGIN
    y_omega = [i -cy1 for i in y_omega]  # BRING Y VALUES TO THE ORIGIN


    #### APPLY THE OFFSET (DESFASE)
    x_lim=7000
    y_lim=5000
    cx=x_lim-max(x)
    cy=y_lim-max(y)
    x = [i + cx for i in x]  # apply offset desfase
    y = [i + cy for i in y]  # apply offset desfase
    x_omega = [i + cx for i in x_omega]  # apply offset desfase
    y_omega = [i + cy for i in y_omega]  # apply offset desfase
    ########..............................................................................................................

    # .....LOAD THE TABLE CALIBRATION DATA.......................................................................................................
    excel_table_calib = r'C:\Users\Juan Pablo Lopez\OneDrive - Rewair A S\Documents\Proyector OPTOMA\table_calib_optoma.xlsx'
    Surface_map = pd.read_excel(excel_table_calib, sheet_name='Surface', header=0)
    tuples_list_real=[]
    tuples_list_autocad=[]
    for i in range(0,len(Surface_map)):
        x_tup_autocad = Surface_map['Autocad_X'][i]
        y_tup_autocad = Surface_map['Autocad_Y'][i]
        tuples_list_autocad.append((x_tup_autocad,y_tup_autocad))
        x_tup_real = Surface_map['Real_X'][i]
        y_tup_real = Surface_map['Real_Y'][i]
        tuples_list_real.append((x_tup_real,y_tup_real))

    #####......................... APPLY THE CALIBRATION......................................
    if scale_mode==1:
        for i in range(0,len(x)):
            input_tuple = (x[i],y[i])
            print(input_tuple)
            xy_tup,idx_low,idx_high,tup_low,tup_high = find_closest_tupleV2(tuples_list_real,tuples_list_autocad,input_tuple)
            x[i] = xy_tup[0]
            y[i] = xy_tup[1]
            #print("Original ",input_tuple," Converted ",xy_tup," tup_low=",tup_low," idx low ",idx_low," tup_high=",tup_high, " idx high ",idx_high)


    ############# BEGGIN THE PLOTTING...............................

    Axis_Limit_x = Axis_Limit
    Axis_Limit_y = Axis_Limit_Y



    layer_qty=0
    x_layer=[]
    y_layer=[]

    for Layer_i in Layer_Names:    # loop through all the vertices of the splines found
        for i in range(0, len(x)):
            if Layer_points[i]==Layer_i:
                x_layer.append(x[i])
                y_layer.append(y[i])

        x_layer,y_layer=polygon_sorter(x_layer,y_layer)

        Layer_Name_plot = Layer_i
        layer_qty = layer_qty + 1
        # Plotting the polygon
        plt.figure(figsize=(image_size, image_size))
        plt.scatter(x_layer, y_layer, marker='d', s=7, color='blue')  #  s=0.1  'bo-' means blue circles connected by lines
        #plt.fill(x_layer, y_layer, facecolor=layer_color, alpha=1)  # Fill the polygon
        title = 'Layer '+str(layer_qty)
        plt.title(title)
        plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
        plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
        plt.style.use('dark_background')
        ax = plt.gca()
        ax.set_facecolor(background_color)
        ax.spines['top'].set_color(layer_color)
        ax.spines['bottom'].set_color(layer_color)
        ax.spines['right'].set_color(layer_color)
        ax.spines['left'].set_color(layer_color)
        ax.tick_params(axis='x', colors=layer_color)
        ax.tick_params(axis='y', colors=layer_color)
        ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
        plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
        plt.savefig('Layers.png')
        # add slide
        slide = prs.slides.add_slide(slide_layout)
        img_path = 'Layers.png'
        left = Inches(Left_centering)  # Adjust position as needed in icnhes
        top = Inches(Top_centering)  # Adjust position as needed in inches
        slide.shapes.add_picture(img_path, left, top)
        if close_image == 1:
            plt.close('all')
        #print(title)
        x_layer = []   ### reset x_layer
        y_layer = []   ### reset y_layer




    # if len(x_omega)>0:
    #             x_layer=x_omega
    #             y_layer=y_omega
    #             Layer_Name_plot='OMEGA'
    #
    #             # Plotting the polygon
    #             plt.figure(figsize=(image_size, image_size))
    #             plt.scatter(x_layer, y_layer,marker='x', s=3, color='blue') # 'bo-' means blue circles connected by lines
    #             title='Omega'
    #             plt.title(title)
    #             plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
    #             plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
    #             plt.style.use('dark_background')
    #             ax = plt.gca()
    #             ax.set_facecolor(background_color)
    #             ax.spines['top'].set_color(layer_color)
    #             ax.spines['bottom'].set_color(layer_color)
    #             ax.spines['right'].set_color(layer_color)
    #             ax.spines['left'].set_color(layer_color)
    #             ax.tick_params(axis='x', colors=layer_color)
    #             ax.tick_params(axis='y', colors=layer_color)
    #             ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
    #             plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
    #             plt.savefig('Layers.png')
    #             # add slide
    #             slide = prs.slides.add_slide(slide_layout)
    #             img_path = 'Layers.png'
    #             left = Inches(Left_centering)  # Adjust position as needed in icnhes
    #             top = Inches(Top_centering)  # Adjust position as needed in inches
    #             slide.shapes.add_picture(img_path, left, top)
    #             if close_image==1:
    #              plt.close('all')

                # prs.save('Layers.pptx')

    # Save the PowerPoint presentation
    Name0=name[:-4]
    Name2 = Name0 + '.pptx'
    out_Name = f"{Destination_path}\{Name2}"
    prs.save(out_Name)


    print('total layers = ',layer_qty)
def read_and_plot_PSv2(file_path,name,image_size,Left_centering,Top_centering,Axis_Limit,scale_mode,background_color,layer_color,close_image,Destination_path,Axis_Limit_Y,filling,reorder_points):
    ###
    import ezdxf
    import pandas as pd
    import matplotlib.pyplot as plt
    from pptx import Presentation
    from pptx.util import Inches
    from dxf_surface_map import euclidean_distance,find_closest_tupleV2
    from dxf_ProcessKit_polygon import polygon_sorter


    ##............READING THE DXF
    doc = ezdxf.readfile(file_path)

    # POWERPOINT SLIDE
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank slide layout

    # Extract entities (splines)
    msp = doc.modelspace()
    for entity in msp:
        # Get the type of the entity
        print(entity.dxftype())
    splines = msp.query('SPLINE')
    polylines= msp.query('POLYLINE')
    lw_polylines= msp.query('LWPOLYLINE')
    lines = msp.query('LINE')
    i=0
    x=[]
    y=[]
    x_omega=[]
    y_omega=[]
    color=[]

    ####.................... Iterate over splines/polylines/lw-polyinies/lines


    qty_spline=0
    Layer_Names=[]
    Layer_points=[]

    for spline in lines:
        # Check if the spline is a BSpline
            if spline.dxftype() == 'LINE':  # here we look for entity type LINE
                # Get control points of the spline

                layer_name = spline.dxf.layer  # Get the layer name
                # layer = doc.layers.get(layer_name)  # Get the layer object
                # color = layer.dxf.color  # Get the color from the layer

                control_points = [spline.dxf.start, spline.dxf.end]
                for point in control_points:
                    x.append(round(point[0],6))  # Here we extract the x-coordenates of the vertices of the spline rounded to 8 decimals
                    y.append(round(point[1],6))  # Here we extract the y-coordenates of the vertices of the spline  8 rounded to 8 decimals
                    color.append(spline.dxf.color)
                    Layer_points.append(layer_name)
                    Layer_Names.append(layer_name)



    for spline in splines:
        # Check if the spline is a BSpline

        if spline.dxftype() == 'SPLINE':  # here we look for entity type SPLINE
                control_points = spline._control_points
                layer_name = spline.dxf.layer
                # EXTRACT control points, i.e  VERTICES OF THE SPLINES
                for point in control_points:
                    x.append(round(point[0], 6))  # Here we extract the x-coordenates of the vertices of the spline rounded to 8 decimals
                    y.append(round(point[1], 6))  # Here we extract the y-coordenates of the vertices of the spline  8 rounded to 8 decimals
                    color.append(spline.dxf.color)
                    Layer_points.append(layer_name)
                    Layer_Names.append(layer_name)


    for spline in lw_polylines:
        # Check if the spline is a BSpline
            if spline.dxftype() == 'LWPOLYLINE':  # here we look for entity type SPLINE
                layer_name = spline.dxf.layer  # Get the layer name
                control_points = spline.get_points('xy')

                # Print control points
                for point in control_points:
                    i = i + 1
                    # print("Vertice:", point,'vertice number', i)
                    x.append(round(point[0], 6))  # Here we extract the x-coordenates of the vertices of the spline
                    y.append(round(point[1], 6))  # Here we extract the y-coordenates of the vertices of the spline
                    Layer_points.append(layer_name)
                    Layer_Names.append(layer_name)


    #### REMOVE DUPLICATES IN THE LAYERS NAMES
    seen = set()
    Layer_Names = [x for x in Layer_Names if not (x in seen or seen.add(x))]

    #####....................................SHIFT POSITION .............................
    #### BRING TO THE ORIGIN
    cx1=min(x)
    cy1=min(y)
    x = [i -cx1 for i in x]  # BRING X VALUES TO THE ORIGIN
    y = [i -cy1 for i in y]  # BRING Y VALUES TO THE ORIGIN
    x_omega = [i -cx1 for i in x_omega]  # BRING X VALUES TO THE ORIGIN
    y_omega = [i -cy1 for i in y_omega]  # BRING Y VALUES TO THE ORIGIN


    #### APPLY THE OFFSET (DESFASE)
    x_lim=9000
    y_lim=4900
    cx=x_lim-max(x)
    cy=y_lim-max(y)
    x = [i + cx for i in x]  # apply offset desfase
    y = [i + cy for i in y]  # apply offset desfase
    x_omega = [i + cx for i in x_omega]  # apply offset desfase
    y_omega = [i + cy for i in y_omega]  # apply offset desfase
    ########..............................................................................................................

    # .....LOAD THE TABLE CALIBRATION DATA.......................................................................................................
    excel_table_calib = r'C:\Users\Juan Pablo Lopez\OneDrive - Rewair A S\Documents\Proyector OPTOMA\table_calib_optoma.xlsx'
    Surface_map = pd.read_excel(excel_table_calib, sheet_name='Surface', header=0)
    tuples_list_real=[]
    tuples_list_autocad=[]
    for i in range(0,len(Surface_map)):
        x_tup_autocad = Surface_map['Autocad_X'][i]
        y_tup_autocad = Surface_map['Autocad_Y'][i]
        tuples_list_autocad.append((x_tup_autocad,y_tup_autocad))
        x_tup_real = Surface_map['Real_X'][i]
        y_tup_real = Surface_map['Real_Y'][i]
        tuples_list_real.append((x_tup_real,y_tup_real))

    #####......................... APPLY THE CALIBRATION......................................
    if scale_mode==1:
        for i in range(0,len(x)):
            input_tuple = (x[i],y[i])
            print(input_tuple)
            xy_tup,idx_low,idx_high,tup_low,tup_high = find_closest_tupleV2(tuples_list_real,tuples_list_autocad,input_tuple)
            x[i] = xy_tup[0]
            y[i] = xy_tup[1]
            #print("Original ",input_tuple," Converted ",xy_tup," tup_low=",tup_low," idx low ",idx_low," tup_high=",tup_high, " idx high ",idx_high)


    ############# BEGGIN THE PLOTTING...............................

    Axis_Limit_x = Axis_Limit
    Axis_Limit_y = Axis_Limit_Y


    layer_qty=0
    x_layer=[]
    y_layer=[]

    for Layer_i in Layer_Names:    # loop through all the vertices of the splines found
        for i in range(0, len(x)):
            if Layer_points[i]==Layer_i:
                x_layer.append(x[i])
                y_layer.append(y[i])
        if (reorder_points==1) and (Layer_i[-1]!="c"):
          x_layer,y_layer=polygon_sorter(x_layer,y_layer)

        Layer_Name_plot = Layer_i
        layer_qty = layer_qty + 1
        # Plotting the polygon
        plt.figure(figsize=(image_size, image_size))
        plt.scatter(x_layer, y_layer, marker='o', s=1, color='blue')  #  s=0.1  'bo-' means blue circles connected by lines
        if ("Omega" not in Layer_i) and (filling ==1) :
         plt.fill(x_layer, y_layer, facecolor=layer_color, alpha=1)  # Fill the polygon
        title = 'Layer '+str(layer_qty)
        plt.title(title)
        plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
        plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
        plt.style.use('dark_background')
        ax = plt.gca()
        ax.set_facecolor(background_color)
        ax.spines['top'].set_color(layer_color)
        ax.spines['bottom'].set_color(layer_color)
        ax.spines['right'].set_color(layer_color)
        ax.spines['left'].set_color(layer_color)
        ax.tick_params(axis='x', colors=layer_color)
        ax.tick_params(axis='y', colors=layer_color)
        ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
        plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
        plt.savefig('Layers.png')
        # add slide
        slide = prs.slides.add_slide(slide_layout)
        img_path = 'Layers.png'
        left = Inches(Left_centering)  # Adjust position as needed in icnhes
        top = Inches(Top_centering)  # Adjust position as needed in inches
        slide.shapes.add_picture(img_path, left, top)
        if close_image == 1:
            plt.close('all')
        #print(title)
        x_layer = []   ### reset x_layer
        y_layer = []   ### reset y_layer




    # if len(x_omega)>0:
    #             x_layer=x_omega
    #             y_layer=y_omega
    #             Layer_Name_plot='OMEGA'
    #
    #             # Plotting the polygon
    #             plt.figure(figsize=(image_size, image_size))
    #             plt.scatter(x_layer, y_layer,marker='x', s=3, color='blue') # 'bo-' means blue circles connected by lines
    #             title='Omega'
    #             plt.title(title)
    #             plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
    #             plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
    #             plt.style.use('dark_background')
    #             ax = plt.gca()
    #             ax.set_facecolor(background_color)
    #             ax.spines['top'].set_color(layer_color)
    #             ax.spines['bottom'].set_color(layer_color)
    #             ax.spines['right'].set_color(layer_color)
    #             ax.spines['left'].set_color(layer_color)
    #             ax.tick_params(axis='x', colors=layer_color)
    #             ax.tick_params(axis='y', colors=layer_color)
    #             ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
    #             plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
    #             plt.savefig('Layers.png')
    #             # add slide
    #             slide = prs.slides.add_slide(slide_layout)
    #             img_path = 'Layers.png'
    #             left = Inches(Left_centering)  # Adjust position as needed in icnhes
    #             top = Inches(Top_centering)  # Adjust position as needed in inches
    #             slide.shapes.add_picture(img_path, left, top)
    #             if close_image==1:
    #              plt.close('all')

                # prs.save('Layers.pptx')

    # Save the PowerPoint presentation
    Name0=name[:-4]
    Name2 = Name0 + '.pptx'
    out_Name = f"{Destination_path}\{Name2}"
    prs.save(out_Name)


    print('total layers = ',layer_qty)
def read_and_plot_PSv3_doubles(file_path,name,image_size,Left_centering,Top_centering,Axis_Limit,scale_mode,background_color,layer_color,close_image,Destination_path,Axis_Limit_Y,filling,reorder_points,d1,d2):
    ### this is for combining shapes into the same plot WITH MIXED REORDERED AND NOT-REORDERED LAYERS
    ### this works looking for the Layers inside the dxf with names that end in "d1" and "d2", and it joins them into 1 single slide
    ###  if it gives trouble try re-ordering the points
    import ezdxf
    import pandas as pd
    import matplotlib.pyplot as plt
    from pptx import Presentation
    from pptx.util import Inches
    from dxf_surface_map import euclidean_distance,find_closest_tupleV2
    from dxf_ProcessKit_polygon import polygon_sorter
    from matplotlib.patches import Polygon


    ##............READING THE DXF
    doc = ezdxf.readfile(file_path)

    # POWERPOINT SLIDE
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank slide layout

    # Extract entities (splines)
    msp = doc.modelspace()
    for entity in msp:
        # Get the type of the entity
        print(entity.dxftype())
    splines = msp.query('SPLINE')
    polylines= msp.query('POLYLINE')
    lw_polylines= msp.query('LWPOLYLINE')
    lines = msp.query('LINE')
    i=0
    x=[]
    y=[]
    x_omega=[]
    y_omega=[]
    color=[]

    ####.................... Iterate over splines/polylines/lw-polyinies/lines


    qty_spline=0
    Layer_Names=[]
    Layer_points=[]



    for spline in msp:

            if spline.dxftype() == 'LINE':  # here we look for entity type LINE
                # Get control points of the spline

                layer_name = spline.dxf.layer  # Get the layer name
                # layer = doc.layers.get(layer_name)  # Get the layer object
                # color = layer.dxf.color  # Get the color from the layer

                control_points = [spline.dxf.start, spline.dxf.end]
                for point in control_points:
                    x.append(round(point[0],6))  # Here we extract the x-coordenates of the vertices of the spline rounded to 8 decimals
                    y.append(round(point[1],6))  # Here we extract the y-coordenates of the vertices of the spline  8 rounded to 8 decimals
                    color.append(spline.dxf.color)
                    Layer_points.append(layer_name)
                    Layer_Names.append(layer_name)


            elif spline.dxftype() == 'SPLINE':  # here we look for entity type SPLINE
                control_points = spline._control_points
                layer_name = spline.dxf.layer
                # EXTRACT control points, i.e  VERTICES OF THE SPLINES
                for point in control_points:
                    x.append(round(point[0], 6))  # Here we extract the x-coordenates of the vertices of the spline rounded to 8 decimals
                    y.append(round(point[1], 6))  # Here we extract the y-coordenates of the vertices of the spline  8 rounded to 8 decimals
                    color.append(spline.dxf.color)
                    Layer_points.append(layer_name)
                    Layer_Names.append(layer_name)



            elif spline.dxftype() == 'LWPOLYLINE':  # here we look for entity type SPLINE
                layer_name = spline.dxf.layer  # Get the layer name
                control_points = spline.get_points('xy')

                # Print control points
                for point in control_points:
                    i = i + 1
                    # print("Vertice:", point,'vertice number', i)
                    x.append(round(point[0], 6))  # Here we extract the x-coordenates of the vertices of the spline
                    y.append(round(point[1], 6))  # Here we extract the y-coordenates of the vertices of the spline
                    Layer_points.append(layer_name)
                    Layer_Names.append(layer_name)


    #### REMOVE DUPLICATES IN THE LAYERS NAMES
    seen = set()
    Layer_Names = [x for x in Layer_Names if not (x in seen or seen.add(x))]

    #####....................................SHIFT POSITION .............................
    ###### be very careful if using this, all layers must have same origin of coordenates
    #### BRING TO THE ORIGIN
    # cx1=min(x)
    # cy1=min(y)
    # x = [i -cx1 for i in x]  # BRING X VALUES TO THE ORIGIN
    # y = [i -cy1 for i in y]  # BRING Y VALUES TO THE ORIGIN



    # #### APPLY THE OFFSET (DESFASE)
    # x_lim=9000
    # y_lim=4900
    # cx=x_lim-max(x)
    # cy=y_lim-max(y)
    # x = [i + cx for i in x]  # apply offset desfase
    # y = [i + cy for i in y]  # apply offset desfase

    ########..............................................................................................................

    # .....LOAD THE TABLE CALIBRATION DATA.......................................................................................................
    excel_table_calib = r'C:\Users\Juan Pablo Lopez\OneDrive - Rewair A S\Documents\Proyector OPTOMA\table_calib_optoma.xlsx'
    Surface_map = pd.read_excel(excel_table_calib, sheet_name='Surface', header=0)
    tuples_list_real=[]
    tuples_list_autocad=[]
    for i in range(0,len(Surface_map)):
        x_tup_autocad = Surface_map['Autocad_X'][i]
        y_tup_autocad = Surface_map['Autocad_Y'][i]
        tuples_list_autocad.append((x_tup_autocad,y_tup_autocad))
        x_tup_real = Surface_map['Real_X'][i]
        y_tup_real = Surface_map['Real_Y'][i]
        tuples_list_real.append((x_tup_real,y_tup_real))

    #####......................... APPLY THE CALIBRATION......................................
    if scale_mode==1:
        for i in range(0,len(x)):
            input_tuple = (x[i],y[i])
            print(input_tuple)
            xy_tup,idx_low,idx_high,tup_low,tup_high = find_closest_tupleV2(tuples_list_real,tuples_list_autocad,input_tuple)
            x[i] = xy_tup[0]
            y[i] = xy_tup[1]
            #print("Original ",input_tuple," Converted ",xy_tup," tup_low=",tup_low," idx low ",idx_low," tup_high=",tup_high, " idx high ",idx_high)


    ############# BEGGIN THE PLOTTING...............................

    Axis_Limit_x = Axis_Limit
    Axis_Limit_y = Axis_Limit_Y


    layer_qty=0
    x_layer_1=[]
    y_layer_1=[]
    x_layer_2=[]
    y_layer_2=[]
    poly1=[]
    poly2=[]


    for i in range(0, len(x)):
        if Layer_points[i][-2:]=="d1":
            x_layer_1.append(x[i])
            y_layer_1.append(y[i])
            poly1.append((x[i],y[i]))
            layer_name=Layer_points[i]
        elif Layer_points[i][-2:]=="d2":
            x_layer_2.append(x[i])
            y_layer_2.append(y[i])
            poly2.append((x[i], y[i]))

    if (reorder_points == 1):   #### this if you need to re-order points
        if d1==1:
            x_layer_1, y_layer_1 = polygon_sorter(x_layer_1, y_layer_1)
            poly1=[]
            for i in range(0, len(x_layer_1)):
                poly1.append((x_layer_1[i], y_layer_1[i]))
        if d2==1:
            x_layer_2, y_layer_2 = polygon_sorter(x_layer_2, y_layer_2)
            poly2=[]
            for i in range(0, len(x_layer_2)):
                poly2.append((x_layer_2[i], y_layer_2[i]))

    Layer_Name_plot = layer_name
    layer_qty = layer_qty + 1
    # Plotting the polygon
    plt.figure(figsize=(image_size, image_size))
    poly1 = Polygon(poly1, closed=True, edgecolor='b', facecolor='blue', linewidth=2, alpha=1)
    poly2 = Polygon(poly2, closed=True, edgecolor='b', facecolor='blue', linewidth=2, alpha=1)
    #plt.fill(x_layer, y_layer, facecolor=layer_color, alpha=1)  # Fill the polygon
    title = 'Layer '+str(layer_qty)
    plt.title(title)
    plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
    plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
    plt.style.use('dark_background')
    ax = plt.gca()
    ax.set_facecolor(background_color)
    ax.spines['top'].set_color(layer_color)
    ax.spines['bottom'].set_color(layer_color)
    ax.spines['right'].set_color(layer_color)
    ax.spines['left'].set_color(layer_color)
    ax.tick_params(axis='x', colors=layer_color)
    ax.tick_params(axis='y', colors=layer_color)
    ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
    ax.add_patch(poly1)
    ax.add_patch(poly2)
    plt.gca().set_aspect('equal', adjustable='box')
    plt.text(2000, 50, Layer_Name_plot, fontsize=15, color='red')
    plt.savefig('Layers.png')
    # add slide
    slide = prs.slides.add_slide(slide_layout)
    img_path = 'Layers.png'
    left = Inches(Left_centering)  # Adjust position as needed in icnhes
    top = Inches(Top_centering)  # Adjust position as needed in inches
    slide.shapes.add_picture(img_path, left, top)
    if close_image == 1:
        plt.close('all')
    #print(title)



    # Save the PowerPoint presentation
    Name0=name[:-4]
    Name2 = Name0 + '.pptx'
    out_Name = f"{Destination_path}\{Name2}"
    prs.save(out_Name)


    print('total layers = ',layer_qty)
def read_and_plot_PSv3(file_path,name,image_size,Left_centering,Top_centering,Axis_Limit,scale_mode,background_color,layer_color,close_image,Destination_path,Axis_Limit_Y,filling,reorder_points,d1,d2):
    ### Last official, this is PSV2 with entities filling X-Y as they appear in the drawing
    import ezdxf
    import pandas as pd
    import matplotlib.pyplot as plt
    from pptx import Presentation
    from pptx.util import Inches
    from dxf_surface_map import euclidean_distance,find_closest_tupleV2
    from dxf_ProcessKit_polygon import polygon_sorter


    ##............READING THE DXF
    doc = ezdxf.readfile(file_path)

    # POWERPOINT SLIDE
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank slide layout

    # Extract entities (splines)
    msp = doc.modelspace()
    for entity in msp:
        # Get the type of the entity
        print(entity.dxftype())
    splines = msp.query('SPLINE')
    polylines= msp.query('POLYLINE')
    lw_polylines= msp.query('LWPOLYLINE')
    lines = msp.query('LINE')
    i=0
    x=[]
    y=[]
    x_omega=[]
    y_omega=[]
    color=[]

    ####.................... Iterate over splines/polylines/lw-polyinies/lines


    qty_spline=0
    Layer_Names=[]
    Layer_points=[]



    for spline in msp:

            if spline.dxftype() == 'LINE':  # here we look for entity type LINE
                # Get control points of the spline

                layer_name = spline.dxf.layer  # Get the layer name
                # layer = doc.layers.get(layer_name)  # Get the layer object
                # color = layer.dxf.color  # Get the color from the layer

                control_points = [spline.dxf.start, spline.dxf.end]
                for point in control_points:
                    x.append(round(point[0],6))  # Here we extract the x-coordenates of the vertices of the spline rounded to 8 decimals
                    y.append(round(point[1],6))  # Here we extract the y-coordenates of the vertices of the spline  8 rounded to 8 decimals
                    color.append(spline.dxf.color)
                    Layer_points.append(layer_name)
                    Layer_Names.append(layer_name)


            elif spline.dxftype() == 'SPLINE':  # here we look for entity type SPLINE
                control_points = spline._control_points
                layer_name = spline.dxf.layer
                # EXTRACT control points, i.e  VERTICES OF THE SPLINES
                for point in control_points:
                    x.append(round(point[0], 6))  # Here we extract the x-coordenates of the vertices of the spline rounded to 8 decimals
                    y.append(round(point[1], 6))  # Here we extract the y-coordenates of the vertices of the spline  8 rounded to 8 decimals
                    color.append(spline.dxf.color)
                    Layer_points.append(layer_name)
                    Layer_Names.append(layer_name)



            elif spline.dxftype() == 'LWPOLYLINE':  # here we look for entity type SPLINE
                layer_name = spline.dxf.layer  # Get the layer name
                control_points = spline.get_points('xy')

                # Print control points
                for point in control_points:
                    i = i + 1
                    # print("Vertice:", point,'vertice number', i)
                    x.append(round(point[0], 6))  # Here we extract the x-coordenates of the vertices of the spline
                    y.append(round(point[1], 6))  # Here we extract the y-coordenates of the vertices of the spline
                    Layer_points.append(layer_name)
                    Layer_Names.append(layer_name)


    #### REMOVE DUPLICATES IN THE LAYERS NAMES
    seen = set()
    Layer_Names = [x for x in Layer_Names if not (x in seen or seen.add(x))]

    #####....................................SHIFT POSITION .............................
    ###### be very careful if using this, all layers must have same origin of coordenates
    # #### BRING TO THE ORIGIN
    # cx1=min(x)
    # cy1=min(y)
    # x = [i -cx1 for i in x]  # BRING X VALUES TO THE ORIGIN
    # y = [i -cy1 for i in y]  # BRING Y VALUES TO THE ORIGIN
    #
    #
    #
    # #### APPLY THE OFFSET (DESFASE)
    # x_lim=10000
    # y_lim=4400
    # cx=x_lim-max(x)
    # cy=y_lim-max(y)
    # x = [i + cx for i in x]  # apply offset desfase
    # y = [i + cy for i in y]  # apply offset desfase

    ########..............................................................................................................

    # .....LOAD THE TABLE CALIBRATION DATA.......................................................................................................
    excel_table_calib = r'C:\Users\Juan Pablo Lopez\OneDrive - Rewair A S\Documents\Proyector OPTOMA\table_calib_optoma.xlsx'
    Surface_map = pd.read_excel(excel_table_calib, sheet_name='Surface', header=0)
    tuples_list_real=[]
    tuples_list_autocad=[]
    for i in range(0,len(Surface_map)):
        x_tup_autocad = Surface_map['Autocad_X'][i]
        y_tup_autocad = Surface_map['Autocad_Y'][i]
        tuples_list_autocad.append((x_tup_autocad,y_tup_autocad))
        x_tup_real = Surface_map['Real_X'][i]
        y_tup_real = Surface_map['Real_Y'][i]
        tuples_list_real.append((x_tup_real,y_tup_real))

    #####......................... APPLY THE CALIBRATION......................................
    if scale_mode==1:
        for i in range(0,len(x)):
            input_tuple = (x[i],y[i])
            print(input_tuple)
            xy_tup,idx_low,idx_high,tup_low,tup_high = find_closest_tupleV2(tuples_list_real,tuples_list_autocad,input_tuple)
            x[i] = xy_tup[0]
            y[i] = xy_tup[1]
            #print("Original ",input_tuple," Converted ",xy_tup," tup_low=",tup_low," idx low ",idx_low," tup_high=",tup_high, " idx high ",idx_high)


    ############# BEGGIN THE PLOTTING...............................

    Axis_Limit_x = Axis_Limit
    Axis_Limit_y = Axis_Limit_Y


    layer_qty=0
    x_layer=[]
    y_layer=[]

    for Layer_i in Layer_Names:    # loop through all the vertices of the splines found
        for i in range(0, len(x)):
            if Layer_points[i]==Layer_i:
                x_layer.append(x[i])
                y_layer.append(y[i])
        if (reorder_points==1) and (Layer_i[-1]!="c"):
          x_layer,y_layer=polygon_sorter(x_layer,y_layer)

        Layer_Name_plot = Layer_i
        layer_qty = layer_qty + 1
        # Plotting the polygon
        plt.figure(figsize=(image_size, image_size))
        plt.scatter(x_layer, y_layer, marker='o', s=0.1, color='blue')  # s=1  'bo-' means blue circles connected by lines
        if filling==1:
           plt.fill(x_layer, y_layer, facecolor=layer_color, alpha=1)  # Fill the polygon
        title = 'Layer '+str(layer_qty)
        plt.title(title)
        plt.xlim(0, Axis_Limit_x)  # Set x-axis limit from 0 to 4
        plt.ylim(0, Axis_Limit_y)  # Set y-axis limit from 0 to 4
        plt.style.use('dark_background')
        ax = plt.gca()
        ax.set_facecolor(background_color)
        ax.spines['top'].set_color(layer_color)
        ax.spines['bottom'].set_color(layer_color)
        ax.spines['right'].set_color(layer_color)
        ax.spines['left'].set_color(layer_color)
        ax.tick_params(axis='x', colors=layer_color)
        ax.tick_params(axis='y', colors=layer_color)
        ax.set_title(ax.get_title(), color=layer_color)  # Update title with green color
        plt.text(500, 1800, Layer_Name_plot, fontsize=35, color='red')   # fontsize=15
        plt.savefig('Layers.png')
        # add slide
        slide = prs.slides.add_slide(slide_layout)
        img_path = 'Layers.png'
        left = Inches(Left_centering)  # Adjust position as needed in icnhes
        top = Inches(Top_centering)  # Adjust position as needed in inches
        slide.shapes.add_picture(img_path, left, top)
        if close_image == 1:
            plt.close('all')
        #print(title)
        x_layer = []   ### reset x_layer
        y_layer = []   ### reset y_layer


    # Save the PowerPoint presentation
    Name0=name[:-4]
    Name2 = Name0 + '.pptx'
    out_Name = f"{Destination_path}\{Name2}"
    prs.save(out_Name)


    print('total layers = ',layer_qty)



Origin_path=r'C:\Users\Juan Pablo Lopez\OneDrive - Rewair A S\Desktop\PFILES\Python_versions\LG projecting files\origin'
Destination_path=r'C:\Users\Juan Pablo Lopez\OneDrive - Rewair A S\Desktop\PFILES\Python_versions\LG projecting files\destiny'
image_size = 17  # in inches    12.9
Left_centering = -3.7  # in inches   -1.6
Top_centering = -8.85  # in inches     -4.0
Axis_Limit = 9800  # in MM    10000, last 10015
Axis_Limit_Y=9800  # ONLY FOR V19 and following,  10000, last 10020
background_color = 'white'
layer_color = 'blue'
close_image = 1  # 1- close all images, 0- open all images
scale_mode = 1  # 1-apply scale, 0- no scale
filling=1     # 1- fill the polygons      , 0 - just scatter plot
reorder_points=1   # 1- reorder points to guess polygon,   0- no reorder
d1=1    ### reorder layer d1 =1  ONLY FOR DOUBLES CODE
d2=1    ### reorder layer d2 =1  ONLY FOR DOUBLES CODE


import os
DXF_Names = os.listdir(Origin_path)  # this extracts the names of the files inside source folder
for k1 in DXF_Names:  # this is for eliminating the files that are not dxf
    if k1.endswith("dxf") == False:  # dxf
        Path2Remove = f"{Origin_path}\{k1}"
        os.remove(Path2Remove)
DXF_Names = os.listdir(Origin_path)  # this is to update the names after deleting non dxf files
print(" TOTAL OF ", len(DXF_Names), " DXF UPLOADED")

for name in DXF_Names:
    dxf_file = f"{Origin_path}\{name}"
    #print(name)
    #print(dxf_file)     ### LAST OFFICIAL  PSV3
    read_and_plot_PSv3(dxf_file, name, image_size, Left_centering, Top_centering, Axis_Limit, scale_mode, background_color, layer_color, close_image,Destination_path,Axis_Limit_Y,filling,reorder_points,d1,d2)
    print(name+' fully converted')
print('All files converted')